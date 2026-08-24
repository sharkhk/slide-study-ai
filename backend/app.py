import io
import re
import uuid
import json
import os
import time
import secrets
import threading
import traceback as _tb
import logging
import requests as http
from concurrent.futures import ThreadPoolExecutor, as_completed
from collections import OrderedDict

# ── File logger ────────────────────────────────────────────────────────────────
_LOG_FILE = os.path.join(os.path.dirname(__file__), "debug.log")
logging.basicConfig(
    filename=_LOG_FILE, level=logging.DEBUG,
    format="%(asctime)s %(levelname)s %(message)s"
)
_log = logging.getLogger("app")
from flask import Flask, request, jsonify, send_file, send_from_directory, Response, stream_with_context
from flask_cors import CORS
from pptx import Presentation
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.platypus import (
    BaseDocTemplate, Frame, PageTemplate,
    Paragraph, Spacer, Table, TableStyle,
    HRFlowable, KeepTogether
)
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.pdfbase import pdfmetrics as _pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont as _TTFont

DIST         = os.path.join(os.path.dirname(__file__), "dist")
GROQ_API_KEY = os.environ.get("GROQ_API_KEY", "")
GROQ_MODEL   = os.environ.get("GROQ_MODEL", "llama-3.1-8b-instant")
OLLAMA_URL   = os.environ.get("OLLAMA_URL", "http://localhost:11434")
OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "mistral")

# ── Supabase / Auth ────────────────────────────────────────────────────────────
SUPABASE_URL              = os.environ.get("SUPABASE_URL", "")
SUPABASE_ANON_KEY         = os.environ.get("SUPABASE_ANON_KEY", "")
SUPABASE_SERVICE_ROLE_KEY = os.environ.get("SUPABASE_SERVICE_ROLE_KEY", "")
SUPABASE_JWT_SECRET       = os.environ.get("SUPABASE_JWT_SECRET", "")

# ── Stripe ─────────────────────────────────────────────────────────────────────
STRIPE_SECRET_KEY    = os.environ.get("STRIPE_SECRET_KEY", "")
STRIPE_PUBLISHABLE_KEY = os.environ.get("STRIPE_PUBLISHABLE_KEY", "")
STRIPE_WEBHOOK_SECRET = os.environ.get("STRIPE_WEBHOOK_SECRET", "")
STRIPE_PRICE_ID      = os.environ.get("STRIPE_PRICE_ID", "")
APP_URL              = os.environ.get("APP_URL", "https://slide-study-ai.onrender.com")

_AUTH_ENABLED = bool(SUPABASE_URL and SUPABASE_JWT_SECRET)  # False in local dev without Supabase

DETAIL = {
    "brief":    {"slide_chars": 400,  "max_slides": 30,  "keywords": "8-10",  "bullets": "2-4",  "n_flash": 6,  "n_mcq": 5,  "num_predict": 2048},
    "standard": {"slide_chars": 700,  "max_slides": 60,  "keywords": "18-25", "bullets": "3-8",  "n_flash": 14, "n_mcq": 10, "num_predict": 4096},
    "detailed": {"slide_chars": 1200, "max_slides": 120, "keywords": "25-35", "bullets": "5-12", "n_flash": 20, "n_mcq": 15, "num_predict": 6000},
}

# ── Arabic PDF support ─────────────────────────────────────────────────────────
_ARABIC_FONT      = "NotoNaskhArabic"
_ARABIC_FONT_PATH = "/tmp/NotoNaskhArabic.ttf"
_arabic_font_ok   = False
_arabic_font_lock = threading.Lock()

def _ensure_arabic_font():
    global _arabic_font_ok
    with _arabic_font_lock:
        if _arabic_font_ok:
            return True
        try:
            if not os.path.exists(_ARABIC_FONT_PATH):
                r = http.get(
                    "https://github.com/googlefonts/noto-fonts/raw/main/hinted/ttf/NotoNaskhArabic/NotoNaskhArabic-Regular.ttf",
                    timeout=30, allow_redirects=True
                )
                r.raise_for_status()
                with open(_ARABIC_FONT_PATH, "wb") as fh:
                    fh.write(r.content)
            _pdfmetrics.registerFont(_TTFont(_ARABIC_FONT, _ARABIC_FONT_PATH))
            _arabic_font_ok = True
            return True
        except Exception as exc:
            print(f"Arabic font error: {exc}", flush=True)
            return False

def _ar(text):
    """Reshape + bidi-flip Arabic for correct visual display in LTR PDF renderer."""
    if not text:
        return text
    try:
        import arabic_reshaper
        from bidi.algorithm import get_display
        return get_display(arabic_reshaper.reshape(str(text)))
    except ImportError:
        return str(text)

# ── Supabase client (lazy, uses service-role key → bypasses RLS) ───────────────
_sb_client = None
_sb_lock   = threading.Lock()

def _get_sb():
    global _sb_client
    if _sb_client:
        return _sb_client
    with _sb_lock:
        if _sb_client:
            return _sb_client
        if SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY:
            try:
                from supabase import create_client
                _sb_client = create_client(SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY)
            except Exception as e:
                print(f"Supabase init error: {e}", flush=True)
    return _sb_client

# ── Stripe setup ───────────────────────────────────────────────────────────────
if STRIPE_SECRET_KEY:
    try:
        import stripe as _stripe
        _stripe.api_key = STRIPE_SECRET_KEY
    except ImportError:
        _stripe = None
else:
    _stripe = None

# ── JWT helpers ────────────────────────────────────────────────────────────────
def _get_bearer(req):
    auth = req.headers.get("Authorization", "")
    return auth[7:] if auth.startswith("Bearer ") else None

def _verify_jwt(token):
    """Verify a Supabase-issued JWT. Returns user_id (str) or None."""
    if not token or not SUPABASE_JWT_SECRET:
        return None
    try:
        import jwt as _pyjwt
        payload = _pyjwt.decode(
            token, SUPABASE_JWT_SECRET,
            algorithms=["HS256"], audience="authenticated"
        )
        return payload.get("sub")
    except Exception:
        return None

# ── Token helpers ──────────────────────────────────────────────────────────────
def _consume_token(user_id):
    """
    Atomically consume one token via Supabase RPC.
    Returns (ok: bool, tokens_remaining: int, reason: str)
    """
    sb = _get_sb()
    if not sb:
        return True, 999, ""   # Supabase not configured → dev mode, always allow
    try:
        r = sb.rpc("consume_token", {"p_user_id": user_id}).execute()
        d = r.data if isinstance(r.data, dict) else {}
        ok = d.get("success", False)
        return ok, d.get("tokens_remaining", 0), d.get("reason", "")
    except Exception as exc:
        print(f"consume_token error: {exc}", flush=True)
        return False, 0, "db_error"

def _get_user(user_id):
    """Fetch full user row from Supabase."""
    sb = _get_sb()
    if not sb:
        return None
    try:
        r = sb.table("users").select("*").eq("id", user_id).single().execute()
        return r.data
    except Exception:
        return None

def _get_or_create_referral_code(user_id):
    """Return this user's referral code, generating one if not yet set."""
    import hashlib
    sb = _get_sb()
    if not sb:
        return None
    try:
        r = sb.table("users").select("referral_code").eq("id", user_id).single().execute()
        code = (r.data or {}).get("referral_code")
        if code:
            return code
        code = hashlib.sha256(user_id.encode()).hexdigest()[:8].upper()
        sb.table("users").update({"referral_code": code}).eq("id", user_id).execute()
        return code
    except Exception:
        return None

def _award_referral(new_subscriber_id, sb):
    """Award 10 tokens to the referrer when their referee first subscribes."""
    try:
        r = sb.table("users").select("referred_by, referral_paid").eq("id", new_subscriber_id).single().execute()
        if not r.data:
            return
        referrer_id  = r.data.get("referred_by")
        already_paid = r.data.get("referral_paid", False)
        if not referrer_id or already_paid:
            return
        ref_row = sb.table("users").select("tokens_remaining").eq("id", referrer_id).single().execute()
        if ref_row.data:
            new_bal = (ref_row.data.get("tokens_remaining") or 0) + 10
            sb.table("users").update({"tokens_remaining": new_bal}).eq("id", referrer_id).execute()
        sb.table("users").update({"referral_paid": True}).eq("id", new_subscriber_id).execute()
        _log.info(f"Referral reward: 10 tokens → {referrer_id} (subscriber={new_subscriber_id})")
    except Exception as exc:
        _log.error(f"_award_referral error: {exc}")

def _auth_check(req):
    """
    Extract + verify JWT from request.
    Returns (user_id, error_response_tuple | None).
    When _AUTH_ENABLED is False (local dev), always returns ('dev', None).
    """
    if not _AUTH_ENABLED:
        return "dev", None
    tok = _get_bearer(req)
    uid = _verify_jwt(tok)
    if not uid:
        return None, (jsonify({"error": "Sign in required", "code": "auth_required"}), 401)
    return uid, None

app = Flask(__name__, static_folder=DIST, static_url_path="")
app.config['SECRET_KEY'] = secrets.token_hex(32)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50 MB upload limit
_cors_origins = [o for o in {
    os.environ.get("APP_URL", "").rstrip("/"),
    "http://localhost:3000",   # vite dev server
    "http://localhost:5000",   # local flask
} if o]
CORS(app, origins=_cors_origins)

# ── Security headers on every response ────────────────────────────────────────────────────────
@app.after_request
def _security_headers(resp):
    resp.headers.setdefault("X-Content-Type-Options", "nosniff")
    resp.headers.setdefault("X-Frame-Options", "DENY")
    resp.headers.setdefault("Referrer-Policy", "strict-origin-when-cross-origin")
    resp.headers.setdefault("Permissions-Policy", "camera=(), microphone=(), geolocation=()")
    if request.is_secure or request.headers.get("X-Forwarded-Proto", "").lower() == "https":
        resp.headers.setdefault("Strict-Transport-Security", "max-age=31536000; includeSubDomains")
    if request.path.startswith("/api/"):
        # study material / summaries of sensitive data must never be cached
        resp.headers.setdefault("Cache-Control", "no-store")
    return resp

# Pre-load Arabic font in background so the first Arabic PDF request isn't blocked
threading.Thread(target=_ensure_arabic_font, daemon=True).start()

# ── Per-IP rate limiting for the summarize endpoint ───────────────────────────
_rate_limit_lock = threading.Lock()
_rate_limit      = {}   # ip -> [timestamp, ...]
_RATE_WINDOW     = 60   # seconds
_RATE_MAX        = 5    # requests per window per IP

def _check_rate_limit(ip, scope="main", limit=_RATE_MAX):
    """Return True if request is allowed, False if rate-limited."""
    if _is_private(ip):
        return True
    key = f"{scope}:{ip}"
    now = time.time()
    with _rate_limit_lock:
        times = [t for t in _rate_limit.get(key, []) if now - t < _RATE_WINDOW]
        if len(times) >= limit:
            _rate_limit[key] = times
            return False
        times.append(now)
        _rate_limit[key] = times
    return True

# ── Visitor tracking ──────────────────────────────────────────────────────────
ADMIN_TOKEN  = os.environ.get("ADMIN_TOKEN", "")
if not ADMIN_TOKEN:
    ADMIN_TOKEN = secrets.token_hex(16)
    _log.warning("ADMIN_TOKEN not set — generated a random ephemeral token. "
                 "Set the ADMIN_TOKEN env var to access admin endpoints.")
_visitors    = []
_vis_lock    = threading.Lock()
_blocked_ips = set()   # IPs that are blocked from using the app

_PRIVATE_RANGES = (
    "127.", "::1", "10.", "192.168.", "172.16.", "172.17.", "172.18.",
    "172.19.", "172.20.", "172.21.", "172.22.", "172.23.", "172.24.",
    "172.25.", "172.26.", "172.27.", "172.28.", "172.29.", "172.30.",
    "172.31.", "0.0.0.0",
)

def _is_private(ip):
    return any(ip.startswith(p) for p in _PRIVATE_RANGES)

def _admin_ok():
    """Admin auth: X-Admin-Token header (preferred) or ?token= query,
    compared in constant time."""
    supplied = request.headers.get("X-Admin-Token") or request.args.get("token") or ""
    return secrets.compare_digest(supplied, ADMIN_TOKEN)

def _safe_err(e):
    """User-facing error text: pass through intentional ValueErrors, hide internals."""
    if isinstance(e, ValueError):
        return str(e)
    _log.error("internal error: %s\n%s", e, _tb.format_exc())
    return "Processing failed — please try again."

def _client_ip():
    """Real client IP. Forwarded headers (CF-Connecting-IP / X-Forwarded-For)
    are only trusted when the direct peer is our own proxy (private/loopback);
    a client connecting directly cannot spoof another identity."""
    ra = request.remote_addr or "unknown"
    if not _is_private(ra):
        return ra
    fwd = (request.headers.get("CF-Connecting-IP") or
           request.headers.get("X-Forwarded-For", "").split(",")[0].strip())
    return fwd or ra

def _enrich_geo(entry, ip):
    """Background thread: full geolocation via ip-api.com."""
    if _is_private(ip):
        with _vis_lock:
            entry["country"] = "Local / LAN"
            entry["region"]  = ""
            entry["city"]    = "localhost"
            entry["isp"]     = "private network"
            entry["lat"]     = ""
            entry["lon"]     = ""
        return
    try:
        r = http.get(
            f"http://ip-api.com/json/{ip}"
            f"?fields=status,country,countryCode,regionName,city,isp,org,lat,lon",
            timeout=5
        )
        d = r.json()
        if d.get("status") != "success":
            return
        with _vis_lock:
            entry["country"] = d.get("country", "")
            entry["region"]  = d.get("regionName", "")
            entry["city"]    = d.get("city", "")
            entry["isp"]     = d.get("org") or d.get("isp", "")
            entry["lat"]     = d.get("lat", "")
            entry["lon"]     = d.get("lon", "")
    except Exception:
        pass

@app.before_request
def track_visitor():
    skip = ("/assets/", "/favicon", "/admin")
    if any(request.path.startswith(s) for s in skip):
        # still enforce block on non-admin paths
        pass
    ip = _client_ip()

    # Block check — return 403 immediately for blocked IPs (except admin itself)
    if ip in _blocked_ips and not request.path.startswith("/admin"):
        from flask import abort
        abort(403)

    if any(request.path.startswith(s) for s in ("/assets/", "/favicon")):
        return

    entry = {
        "time":    time.strftime("%Y-%m-%d %H:%M:%S UTC", time.gmtime()),
        "ip":      ip,
        "country": request.headers.get("CF-IPCountry", ""),
        "region":  "",
        "city":    "",
        "isp":     "",
        "lat":     "",
        "lon":     "",
        "path":    request.path,
        "method":  request.method,
        "ua":      (request.headers.get("User-Agent") or "")[:160],
    }
    with _vis_lock:
        _visitors.insert(0, entry)
        if len(_visitors) > 1000:
            _visitors.pop()
    # Always enrich — CF headers don't give city/ISP/coords
    threading.Thread(target=_enrich_geo, args=(entry, ip), daemon=True).start()


@app.route("/admin/block", methods=["POST"])
def admin_block():
    if not _admin_ok():
        return jsonify({"error": "Unauthorized"}), 401
    ip = request.json.get("ip", "").strip()
    if not ip:
        return jsonify({"error": "No IP"}), 400
    with _vis_lock:
        _blocked_ips.add(ip)
    return jsonify({"ok": True, "blocked": ip})


@app.route("/admin/unblock", methods=["POST"])
def admin_unblock():
    if not _admin_ok():
        return jsonify({"error": "Unauthorized"}), 401
    ip = request.json.get("ip", "").strip()
    with _vis_lock:
        _blocked_ips.discard(ip)
    return jsonify({"ok": True, "unblocked": ip})


@app.route("/admin/clear", methods=["POST"])
def admin_clear_log():
    if not _admin_ok():
        return jsonify({"error": "Unauthorized"}), 401
    with _vis_lock:
        _visitors.clear()
    return jsonify({"ok": True})


@app.route("/admin")
def admin_page():
    if not _admin_ok():
        return ("<h2 style='font-family:sans-serif;margin:2rem'>🔒 Unauthorized — "
                "add <code>?token=YOUR_TOKEN</code> to the URL</h2>"), 401
    token = ADMIN_TOKEN
    with _vis_lock:
        vis_copy     = list(_visitors)
        blocked_copy = set(_blocked_ips)

    def flag(cc):
        # Convert 2-letter country code to emoji flag
        if not cc or len(cc) != 2:
            return ""
        return "".join(chr(0x1F1E6 + ord(c) - ord('A')) for c in cc.upper())

    rows = ""
    for v in vis_copy:
        blocked = v["ip"] in blocked_copy
        cc = ""
        # Try to extract country code from country name via CF header stored separately
        loc_parts = [p for p in [v.get("city",""), v.get("region",""), v.get("country","")] if p]
        location  = ", ".join(loc_parts) if loc_parts else "—"
        map_link  = ""
        if v.get("lat") and v.get("lon"):
            map_link = f'<a href="https://maps.google.com/?q={v["lat"]},{v["lon"]}" target="_blank" style="color:#4f8ef7;font-size:11px">📍 map</a>'
        block_btn = (
            f'<button onclick="unblock(\'{v["ip"]}\')" '
            f'style="background:#16a34a;color:#fff;border:none;padding:3px 10px;border-radius:6px;cursor:pointer;font-size:12px">✓ Unblock</button>'
            if blocked else
            f'<button onclick="blockIp(\'{v["ip"]}\')" '
            f'style="background:#dc2626;color:#fff;border:none;padding:3px 10px;border-radius:6px;cursor:pointer;font-size:12px">⛔ Block</button>'
        )
        row_style = "background:#1a0a0a" if blocked else ""
        rows += f"""
          <tr style="{row_style}">
            <td style="color:#8aa0c8;white-space:nowrap">{v['time']}</td>
            <td><b style="{'color:#f87171' if blocked else ''}">{v['ip']}</b>
                {'<span style="background:#7f1d1d;color:#fca5a5;padding:1px 6px;border-radius:4px;font-size:10px;margin-left:4px">BLOCKED</span>' if blocked else ''}
            </td>
            <td>{v.get('country','—')}</td>
            <td>{location} {map_link}</td>
            <td style="max-width:180px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;color:#8aa0c8">{v.get('isp','—')}</td>
            <td style="color:#6b7fa8">{v['path']}</td>
            <td style="max-width:200px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;font-size:11px;color:#4a5f80">{v['ua']}</td>
            <td>{block_btn}</td>
          </tr>"""

    blocked_section = ""
    if blocked_copy:
        blocked_rows = "".join(
            f'<tr><td style="color:#f87171;padding:6px 12px">{ip}</td>'
            f'<td><button onclick="unblock(\'{ip}\')" style="background:#16a34a;color:#fff;border:none;padding:2px 10px;border-radius:6px;cursor:pointer;font-size:12px">Unblock</button></td></tr>'
            for ip in sorted(blocked_copy)
        )
        blocked_section = f"""
        <div style="margin:1.5rem 2rem;background:#1a0a0a;border:1px solid #7f1d1d;border-radius:10px;padding:1rem">
          <h3 style="margin:0 0 0.75rem;color:#f87171;font-size:0.95rem">⛔ Blocked IPs ({len(blocked_copy)})</h3>
          <table style="border-collapse:collapse;font-size:13px"><tbody>{blocked_rows}</tbody></table>
        </div>"""

    return f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8">
<title>Admin — Alimne</title>
<meta http-equiv="refresh" content="20">
<style>
  *{{box-sizing:border-box;margin:0;padding:0}}
  body{{font-family:'Segoe UI',system-ui,sans-serif;background:#050d1a;color:#e8f0ff;min-height:100vh}}
  .topbar{{display:flex;align-items:center;justify-content:space-between;padding:1rem 2rem;
           background:#0a1628;border-bottom:1px solid #1a3a6e;position:sticky;top:0;z-index:10}}
  .topbar h1{{font-size:1rem;font-weight:700;display:flex;align-items:center;gap:.6rem}}
  .badge{{background:#4f8ef7;color:#fff;padding:2px 10px;border-radius:20px;font-size:12px}}
  .badge.red{{background:#dc2626}}
  .actions{{display:flex;gap:.5rem}}
  .btn{{padding:6px 14px;border-radius:8px;border:none;cursor:pointer;font-size:13px;font-weight:600;transition:opacity .2s}}
  .btn:hover{{opacity:.8}}
  .btn-red{{background:#dc2626;color:#fff}}
  .btn-green{{background:#16a34a;color:#fff}}
  .btn-gray{{background:#1e3a5f;color:#8aa0c8}}
  .wrap{{overflow-x:auto}}
  table{{width:100%;border-collapse:collapse;font-size:13px}}
  th{{background:#0f2040;padding:10px 12px;text-align:left;font-weight:600;color:#8aa0c8;
      border-bottom:1px solid #1a3a6e;position:sticky;top:57px;white-space:nowrap}}
  td{{padding:9px 12px;border-bottom:1px solid #0d1e35;vertical-align:middle}}
  tr:hover td{{background:#0a1e38}}
  .empty{{padding:3rem;text-align:center;color:#4a5f80;font-size:0.9rem}}
  #toast{{position:fixed;bottom:1.5rem;right:1.5rem;background:#16a34a;color:#fff;
          padding:.6rem 1.2rem;border-radius:8px;font-size:13px;display:none;z-index:100}}
</style></head>
<body>
<div class="topbar">
  <h1>📊 Alimne — Admin
    <span class="badge">{len(vis_copy)} visits</span>
    {f'<span class="badge red">⛔ {len(blocked_copy)} blocked</span>' if blocked_copy else ''}
  </h1>
  <div class="actions">
    <span style="font-size:12px;color:#4a5f80;align-self:center">Auto-refresh: 20s</span>
    <button class="btn btn-gray" onclick="location.reload()">↻ Refresh</button>
    <button class="btn btn-red" onclick="clearLog()">🗑 Clear Log</button>
  </div>
</div>
{blocked_section}
<div class="wrap">
<table>
  <thead><tr>
    <th>Time (UTC)</th><th>IP Address</th><th>Country</th>
    <th>Location</th><th>ISP / Org</th><th>Path</th><th>User Agent</th><th>Action</th>
  </tr></thead>
  <tbody>
  {rows if rows else '<tr><td colspan="8" class="empty">No visitors recorded yet.</td></tr>'}
  </tbody>
</table>
</div>
<div id="toast"></div>
<script>
const TOKEN = "{token}";
function toast(msg, color="#16a34a"){{
  const t = document.getElementById("toast");
  t.textContent = msg; t.style.background = color; t.style.display = "block";
  setTimeout(()=>t.style.display="none", 2500);
}}
async function blockIp(ip){{
  if(!confirm("Block " + ip + "?\\nThis will 403 all their requests immediately.")) return;
  const r = await fetch("/admin/block?token=" + TOKEN, {{
    method:"POST", headers:{{"Content-Type":"application/json"}},
    body: JSON.stringify({{ip}})
  }});
  if(r.ok){{ toast("⛔ Blocked: " + ip, "#dc2626"); setTimeout(()=>location.reload(),1200); }}
}}
async function unblock(ip){{
  const r = await fetch("/admin/unblock?token=" + TOKEN, {{
    method:"POST", headers:{{"Content-Type":"application/json"}},
    body: JSON.stringify({{ip}})
  }});
  if(r.ok){{ toast("✓ Unblocked: " + ip); setTimeout(()=>location.reload(),1200); }}
}}
async function clearLog(){{
  if(!confirm("Clear all visitor log entries?")) return;
  const r = await fetch("/admin/clear?token=" + TOKEN, {{method:"POST"}});
  if(r.ok){{ toast("🗑 Log cleared"); setTimeout(()=>location.reload(),1200); }}
}}
</script>
</body></html>"""

_SAFE_NAME = re.compile(r'[^\w\-. ]')
_JOB_ID_RE = re.compile(r'^[0-9a-f]{32}$')

def _safe_name(s, maxlen=80):
    return _SAFE_NAME.sub('_', str(s))[:maxlen]

def _valid_job(job_id):
    return bool(_JOB_ID_RE.match(str(job_id)))

def _he(s):
    return str(s).replace('&','&amp;').replace('<','&lt;').replace('>','&gt;').replace('"','&quot;')

# ── Job store (memory ONLY — uploads and guides are never written to disk)
_JOB_TTL   = 5400  # 90 minutes

# Purge any job files left on disk by older versions that persisted to /tmp
import shutil as _shutil
_shutil.rmtree(os.path.join(os.sep, "tmp", "slide-study-jobs"), ignore_errors=True)

_jobs      = OrderedDict()
_jobs_lock = threading.RLock()  # reentrant — get_job may acquire while route holds it

# Ollama can only run one inference at a time locally. Serialise all calls so
# the ThreadPoolExecutor doesn't flood it, which causes truncated JSON output.
_ollama_sem      = threading.Semaphore(1)
_DEBUG_RAW = os.environ.get("DEBUG_RAW") == "1"  # never dump user content unless explicitly enabled
_ollama_raw_lock = threading.Lock()  # protect concurrent debug-file writes

def store_job(job_id, pdf_bytes, md_text, guide, slides, filename):
    ts = time.time()
    with _jobs_lock:
        _jobs[job_id] = {
            "pdf": pdf_bytes, "md": md_text,
            "guide": guide,   "slides": slides,
            "filename": filename, "ts": ts
        }
        stale = [k for k, v in list(_jobs.items()) if time.time() - v["ts"] > _JOB_TTL]
        for k in stale:
            del _jobs[k]

def get_job(job_id):
    with _jobs_lock:
        job = _jobs.get(job_id)
        if job and time.time() - job["ts"] > _JOB_TTL:
            del _jobs[job_id]
            return None
    return job

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY        = colors.HexColor('#0a1628')
NAVY_MID    = colors.HexColor('#1a3a6e')
NAVY_LIGHT  = colors.HexColor('#2e5ca8')
ACCENT      = colors.HexColor('#4f8ef7')
KW_BG       = colors.HexColor('#eef3ff')
ROW_ALT     = colors.HexColor('#f4f7ff')
BORDER      = colors.HexColor('#c5d8ff')
TEXT        = colors.HexColor('#0d1b2e')
TEXT_LIGHT  = colors.HexColor('#4a5f80')
WHITE       = colors.white
CARD_Q      = colors.HexColor('#1a3a6e')
CARD_A      = colors.HexColor('#f8faff')
GREEN       = colors.HexColor('#16a34a')
SECTION_BG  = colors.HexColor('#f0f5ff')


# ── Ollama helpers ─────────────────────────────────────────────────────────────

def ollama_running():
    if GROQ_API_KEY:
        return True
    try:
        return http.get(f"{OLLAMA_URL}/api/tags", timeout=3).status_code == 200
    except Exception:
        return False

def ollama_models():
    if GROQ_API_KEY:
        return [GROQ_MODEL]
    try:
        return [m["name"] for m in http.get(f"{OLLAMA_URL}/api/tags", timeout=3).json().get("models", [])]
    except Exception:
        return []

def _extract_json(text):
    text = re.sub(r'^```(?:json)?\s*', '', text.strip(), flags=re.IGNORECASE)
    text = re.sub(r'\s*```$', '', text).strip()

    # ── Pass 1: try the whole text as-is ────────────────────────────────────
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass

    # ── Find the start of the first JSON object or array ───────────────────
    obj_pos = text.find('{')
    arr_pos = text.find('[')
    if obj_pos == -1 and arr_pos == -1:
        raise ValueError("No JSON found in model output")
    if arr_pos != -1 and (obj_pos == -1 or arr_pos < obj_pos):
        start, open_c, close_c = arr_pos, '[', ']'
    else:
        start, open_c, close_c = obj_pos, '{', '}'

    # ── Pass 2: string-aware bracket scan ──────────────────────────────────
    # Tracks whether we're inside a JSON string so { } inside strings are ignored
    depth, in_str, escaped, end = 0, False, False, -1
    stack = []          # track every opening bracket for repair below
    for i, ch in enumerate(text[start:], start):
        if escaped:
            escaped = False
            continue
        if ch == '\\' and in_str:
            escaped = True
            continue
        if ch == '"':
            in_str = not in_str
            continue
        if in_str:
            continue
        if ch in ('{', '['):
            depth += 1
            stack.append(ch)
        elif ch in ('}', ']'):
            depth -= 1
            if stack:
                stack.pop()
            if depth == 0:
                end = i
                break

    if end != -1:
        try:
            return json.loads(text[start:end + 1])
        except json.JSONDecodeError:
            pass

    # ── Pass 3: JSON is truncated — repair by closing open brackets ─────────
    # Re-scan to find the deepest valid position and close everything open
    pairs = {'{': '}', '[': ']'}
    depth2, in_str2, escaped2 = 0, False, False
    open_stack = []
    last_good_pos = start  # last position where depth was 0 between top-level items
    for i, ch in enumerate(text[start:], start):
        if escaped2:
            escaped2 = False
            continue
        if ch == '\\' and in_str2:
            escaped2 = True
            continue
        if ch == '"':
            in_str2 = not in_str2
            continue
        if in_str2:
            continue
        if ch in ('{', '['):
            depth2 += 1
            open_stack.append(ch)
        elif ch in ('}', ']'):
            depth2 -= 1
            if open_stack:
                open_stack.pop()

    # Trim trailing incomplete entry: remove everything after the last comma
    # at depth==1 so the partial last entry is dropped
    snippet = text[start:].rstrip()
    # Remove trailing comma + anything after it (the cut-off entry)
    snippet = re.sub(r',\s*[^,\]\}]*$', '', snippet)
    # Close all still-open brackets
    closing = ''.join(pairs[c] for c in reversed(open_stack))
    # Make sure we haven't over-trimmed: if snippet is just the opening char, add empty body
    repaired = snippet + closing
    try:
        return json.loads(repaired)
    except json.JSONDecodeError:
        pass

    # ── Pass 4: nuclear fallback — try every possible closing combination ──
    for suffix in ['}', ']}', '}}', '"]}}', '"]}', '"}']:
        try:
            return json.loads(snippet + suffix)
        except json.JSONDecodeError:
            continue

    raise ValueError(f"Cannot parse model output (length={len(text)}). "
                     "Model may have returned incomplete JSON.")

def _call_ollama(prompt, retries=3, num_predict=4096):
    if GROQ_API_KEY:
        return _call_groq(prompt, retries=retries, max_tokens=num_predict)
    payload = {
        "model": OLLAMA_MODEL,
        "format": "json",
        "messages": [
            {"role": "system", "content": "You output only valid JSON. No markdown, no explanation."},
            {"role": "user",   "content": prompt},
        ],
        "stream": False,
        "options": {"temperature": 0, "num_predict": num_predict},
    }
    last_err = None
    for attempt in range(retries):
        try:
            with _ollama_sem:
                r = http.post(f"{OLLAMA_URL}/api/chat", json=payload, timeout=300)
                r.raise_for_status()
                raw_content = r.json()["message"]["content"]

            if _DEBUG_RAW:
                with _ollama_raw_lock:
                    with open(os.path.join(os.path.dirname(__file__), "ollama_raw.txt"), "w", encoding="utf-8") as _f:
                        _f.write(raw_content)

            return _extract_json(raw_content)
        except Exception as e:
            last_err = e
            _log.warning("OLLAMA attempt %d/%d failed: %s", attempt + 1, retries, e)
            if attempt < retries - 1:
                time.sleep(1.5 * (attempt + 1))
    raise last_err


def _call_groq(prompt, retries=5, max_tokens=2048):
    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": GROQ_MODEL,
        "messages": [
            {"role": "system", "content": "You output only valid JSON. No markdown, no explanation."},
            {"role": "user",   "content": prompt},
        ],
        "temperature": 0,
        "max_tokens": max_tokens,
    }
    last_err = None
    for attempt in range(retries):
        try:
            r = http.post(
                "https://api.groq.com/openai/v1/chat/completions",
                json=payload, headers=headers, timeout=120
            )
            if r.status_code == 429:
                wait = int(r.headers.get("retry-after", 10))
                _log.warning("GROQ rate limited — waiting %ds", wait)
                time.sleep(wait)
                continue
            r.raise_for_status()
            raw_content = r.json()["choices"][0]["message"]["content"]
            return _extract_json(raw_content)
        except Exception as e:
            last_err = e
            _log.warning("GROQ attempt %d/%d failed: %s", attempt + 1, retries, e)
            if attempt < retries - 1:
                time.sleep(2 * (attempt + 1))
    raise last_err


# ── Three-pass AI processing ──────────────────────────────────────────────────

def _as_dict(result, list_key=None):
    """If the model returned a list instead of a dict, coerce it."""
    if isinstance(result, dict):
        return result
    if isinstance(result, list):
        # If list contains dicts, pick the first one
        first_dict = next((x for x in result if isinstance(x, dict)), None)
        if first_dict:
            return first_dict
        # Otherwise wrap under the given key
        if list_key:
            return {list_key: result}
    return {}


def _detect_language(content):
    """Detect 'ar' or 'en' from slides list or plain text. Based on Arabic char ratio."""
    if isinstance(content, list):
        sample = " ".join(
            f"{s.get('title', '')} {s.get('content', '')}" for s in content[:20]
        )
    else:
        sample = str(content)[:6000]
    arabic = sum(1 for c in sample if '؀' <= c <= 'ۿ')
    alpha  = sum(1 for c in sample if c.isalpha())
    return "ar" if alpha > 0 and arabic / alpha > 0.25 else "en"


def pass1_overview(slides, language, dcfg=None):
    """Get title, objectives, section groupings, and keywords."""
    dcfg = dcfg or DETAIL["standard"]
    lang = "in Arabic" if language == "ar" else "in English"
    # Include ALL slides for grouping; use title-only for slides beyond max_slides
    # to keep the LLM input within token limits for large presentations.
    outline = "".join(
        f"Slide {s['slide_num']}: {s['title']}\n{s['content'][:dcfg['slide_chars']]}\n\n"
        if i < dcfg["max_slides"] else
        f"Slide {s['slide_num']}: {s['title']}\n"
        for i, s in enumerate(slides)
    )
    result = _call_ollama(f"""Create a study guide overview {lang} from this PowerPoint outline.

{outline}

Return JSON:
{{
  "title": "TOPIC IN CAPS",
  "subtitle": "course/session description",
  "objectives": ["objective 1", "objective 2", "objective 3"],
  "sections": [
    {{"title": "SECTION NAME", "slide_nums": [1, 2, 3]}}
  ],
  "keywords": [
    {{"term": "Term", "definition": "Specific 1-sentence definition or role description"}}
  ]
}}

Rules:
- Group related slides into 3-7 logical sections
- objectives: extract from learning objectives slide or infer from content
- keywords: {dcfg['keywords']} terms — every key concept, acronym, role, process mentioned
- Output JSON only""", num_predict=dcfg["num_predict"])
    return _as_dict(result)


def pass2_section(title, section_slides, language, dcfg=None):
    """Get detailed bullets + optional comparison table for one section."""
    dcfg = dcfg or DETAIL["standard"]
    lang = "in Arabic" if language == "ar" else "in English"
    content = "".join(
        f"Slide {s['slide_num']}: {s['title']}\n{s['content']}\n\n"
        for s in section_slides
    )
    result = _call_ollama(f"""Extract detailed exam study notes {lang} for the section "{title}".

{content}

Return JSON:
{{
  "bullets": [
    "Full specific fact, definition, or step from the content",
    "Another detailed point — include names, numbers, roles, processes"
  ],
  "table": {{
    "headers": ["Column 1", "Column 2"],
    "rows": [["value", "value"]]
  }}
}}

Rules:
- bullets: {dcfg['bullets']} specific, exam-worthy facts from the slides
- table: include ONLY if content has roles/comparisons/structured lists; otherwise omit the table field entirely
- Output JSON only""", num_predict=dcfg["num_predict"])
    result = _as_dict(result, list_key="bullets")
    # If model returned a flat list of strings under "bullets", normalise each element
    bullets = result.get("bullets", [])
    result["bullets"] = [str(b) for b in bullets if b]
    return result


def pass3_flashcards(guide, language, dcfg=None):
    """Generate Q&A flash cards from the guide content."""
    dcfg = dcfg or DETAIL["standard"]
    lang = "in Arabic" if language == "ar" else "in English"
    kw_text  = ", ".join(k["term"] for k in guide.get("keywords", [])[:20] if isinstance(k, dict))
    sec_text = " | ".join(s["title"] for s in guide.get("sections", []) if isinstance(s, dict))
    bullets_ctx = ""
    for sec in guide.get("sections", []):
        if isinstance(sec, dict) and sec.get("bullets"):
            bullets_ctx += f"\n{sec.get('title','')}:\n" + "\n".join(f"- {b}" for b in sec["bullets"])
    result = _call_ollama(f"""Create exam flash cards {lang} for a study guide about: {guide.get('title', '')}.

Topics: {sec_text}
Key terms: {kw_text}
Study content:{bullets_ctx}

Return JSON:
{{
  "flashcards": [
    {{"q": "Question that tests a key concept?", "a": "Clear, concise answer"}},
    {{"q": "Define [term]?", "a": "Definition"}},
    {{"q": "What are the responsibilities of [role]?", "a": "Specific responsibilities"}}
  ]
}}

Rules:
- Create exactly {dcfg['n_flash']} flash cards
- Base EVERY question directly on the study content provided above — no generic questions
- Mix definition questions, "what is" questions, "name the" questions, and role/responsibility questions
- Answers must be specific, referencing actual details from the content
- Output JSON only""", num_predict=dcfg["num_predict"])
    # Model may return the array directly instead of wrapping it
    if isinstance(result, list):
        return {"flashcards": [x for x in result if isinstance(x, dict)]}
    return _as_dict(result, list_key="flashcards")


def pass4_mcq(guide, language, dcfg=None):
    """Generate multiple-choice quiz questions."""
    dcfg   = dcfg or DETAIL["standard"]
    lang   = "in Arabic" if language == "ar" else "in English"
    n      = dcfg["n_mcq"]
    kw_text  = ", ".join(k["term"] for k in guide.get("keywords", [])[:20] if isinstance(k, dict))
    sec_text = " | ".join(s["title"] for s in guide.get("sections", []) if isinstance(s, dict))
    bullets_ctx = ""
    for sec in guide.get("sections", []):
        if isinstance(sec, dict) and sec.get("bullets"):
            bullets_ctx += f"\n{sec.get('title','')}:\n" + "\n".join(f"- {b}" for b in sec["bullets"])
    result = _call_ollama(f"""Create {n} multiple-choice exam questions {lang} for: {guide.get('title','')}.

Topics: {sec_text}
Key terms: {kw_text}
Study content:{bullets_ctx}

Return JSON:
{{
  "mcqs": [
    {{
      "q": "Question text?",
      "options": ["A. option one", "B. option two", "C. option three", "D. option four"],
      "answer": "A",
      "explanation": "Why A is correct"
    }}
  ]
}}

Rules:
- Exactly {n} questions, 4 options each (A B C D)
- answer: just the letter
- Mix easy and hard questions
- JSON only""", num_predict=dcfg["num_predict"])
    if isinstance(result, list):
        return {"mcqs": [x for x in result if isinstance(x, dict)]}
    return _as_dict(result, list_key="mcqs")


def _sections_parallel(sections, content_slides, language, dcfg):
    """Process sections sequentially (Groq rate limits prevent safe concurrency).
    Yields plain dicts. Flashcards+MCQ are parallelized separately."""
    n = len(sections)
    for i, sec in enumerate(sections):
        yield {"step": "section", "msg": f"Section {i+1}/{n}: {sec.get('title', '')}…"}
        nums = set(sec.get("slide_nums", []))
        sl = [s for s in content_slides if s["slide_num"] in nums]
        if not sl:
            chunk = max(1, len(content_slides) // n)
            start = i * chunk
            end   = start + chunk if i < n - 1 else len(content_slides)
            sl    = content_slides[start:end] or content_slides
        try:
            det = pass2_section(sec.get("title", ""), sl, language, dcfg)
            sec["bullets"] = det.get("bullets", [])
            if isinstance(det.get("table"), dict):
                sec["table"] = det["table"]
        except Exception:
            _log.error("pass2 [%s] error:\n%s", sec.get("title", "?"), _tb.format_exc())
            sec["bullets"] = []
        yield {"step": "section", "msg": f"Sections: {i+1}/{n} done…"}


def _flashcards_mcq_parallel(overview, language, dcfg):
    """Run pass3 then pass4 sequentially (Groq free tier rate limits concurrent calls).
    Yields plain dicts."""
    yield {"step": "flashcards", "msg": "Generating flash cards…"}
    try:
        overview["flashcards"] = pass3_flashcards(overview, language, dcfg).get("flashcards", [])
    except Exception:
        _log.error("pass3 error:\n%s", _tb.format_exc())
        overview["flashcards"] = []
    yield {"step": "flashcards", "msg": "Flash cards ready…"}

    yield {"step": "mcq", "msg": "Generating quiz…"}
    try:
        overview["mcqs"] = pass4_mcq(overview, language, dcfg).get("mcqs", [])
    except Exception:
        _log.error("pass4 error:\n%s", _tb.format_exc())
        overview["mcqs"] = []
    yield {"step": "mcq", "msg": "Quiz ready…"}


def build_markdown(guide):
    """Convert guide dict to a Markdown string."""
    lines = [f"# {guide.get('title', 'Study Guide')}", ""]
    if guide.get("subtitle"):
        lines += [f"*{guide['subtitle']}*", ""]
    objs = [o for o in guide.get("objectives", []) if isinstance(o, str)]
    if objs:
        lines += ["## Learning Objectives", ""]
        for o in objs: lines.append(f"- {o}")
        lines.append("")
    for sec in guide.get("sections", []):
        if not isinstance(sec, dict): continue
        lines += [f"## {sec.get('title', '')}", ""]
        for b in sec.get("bullets", []): lines.append(f"- {b}")
        tbl = sec.get("table")
        if isinstance(tbl, dict) and tbl.get("headers") and tbl.get("rows"):
            lines.append("")
            lines.append("| " + " | ".join(tbl["headers"]) + " |")
            lines.append("|" + "|".join(["---"] * len(tbl["headers"])) + "|")
            for row in tbl["rows"]:
                lines.append("| " + " | ".join(str(c) for c in row) + " |")
        lines.append("")
    kws = [k for k in guide.get("keywords", []) if isinstance(k, dict)]
    if kws:
        lines += ["## Keywords Cheatsheet", ""]
        for k in kws: lines.append(f"**{k.get('term','')}** — {k.get('definition','')}")
        lines.append("")
    fcs = [f for f in guide.get("flashcards", []) if isinstance(f, dict)]
    if fcs:
        lines += ["## Flash Cards", ""]
        for i, fc in enumerate(fcs, 1):
            lines += [f"**Q{i}:** {fc.get('q','')}", f"**A:** {fc.get('a','')}", ""]
    mcqs = [m for m in guide.get("mcqs", []) if isinstance(m, dict)]
    if mcqs:
        lines += ["## Multiple Choice Questions", ""]
        for i, m in enumerate(mcqs, 1):
            lines.append(f"**{i}. {m.get('q','')}**")
            for opt in m.get("options", []): lines.append(f"   {opt}")
            lines += [f"   ✓ **{m.get('answer','')}** — {m.get('explanation','')}", ""]
    return "\n".join(lines)


def ask_ollama(slides, language, progress_cb=None):
    """Orchestrate three-pass processing with optional progress callback."""
    content_slides = [s for s in slides if s["content"].strip() or s["title"].strip()]

    if progress_cb: progress_cb("overview", f"Analysing structure ({len(content_slides)} slides)…")
    overview = pass1_overview(content_slides, language)

    # Normalize sections/keywords (mistral may return plain strings)
    raw_sec = overview.get("sections", [])
    overview["sections"] = [
        s if isinstance(s, dict) else {"title": str(s), "slide_nums": []}
        for s in (raw_sec if isinstance(raw_sec, list) else [])
    ]
    raw_kw = overview.get("keywords", [])
    overview["keywords"] = [
        k if isinstance(k, dict) else {"term": str(k), "definition": ""}
        for k in (raw_kw if isinstance(raw_kw, list) else [])
    ]
    if not overview["sections"]:
        overview["sections"] = [{"title": overview.get("title", "Overview"), "slide_nums": [s["slide_num"] for s in content_slides]}]

    sections = overview["sections"]
    n = len(sections)
    for i, sec in enumerate(sections):
        if progress_cb: progress_cb("section", f"Building section {i+1} of {n}: {sec.get('title','')}…")
        nums = set(sec.get("slide_nums", []))
        sl = [s for s in content_slides if s["slide_num"] in nums]
        if not sl:
            chunk = max(1, len(content_slides) // n)
            start = i * chunk
            end   = start + chunk if i < n - 1 else len(content_slides)
            sl    = content_slides[start:end] or content_slides
        detail = pass2_section(sec.get("title", ""), sl, language)
        sec["bullets"] = detail.get("bullets", [])
        if isinstance(detail.get("table"), dict):
            sec["table"] = detail["table"]

    if progress_cb: progress_cb("flashcards", "Generating flash cards…")
    try:
        fc = pass3_flashcards(overview, language)
        overview["flashcards"] = fc.get("flashcards", [])
    except Exception:
        overview["flashcards"] = []

    return overview


# ── File extraction — PPTX, PPT binary, PDF ───────────────────────────────────

def _extract_pptx_raw(raw):
    """Open as pptx; on relationship errors fall back to raw ZIP XML walk."""
    try:
        prs = Presentation(io.BytesIO(raw))
    except Exception as e:
        if "officeDocument" not in str(e) and "relationship" not in str(e):
            raise ValueError(f"Cannot open PowerPoint file: {e}")
        # Fallback: read slide XML directly from the ZIP
        import zipfile, xml.etree.ElementTree as ET
        slides = []
        with zipfile.ZipFile(io.BytesIO(raw)) as z:
            slide_files = sorted(
                [n for n in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml', n)],
                key=lambda x: int(re.findall(r'\d+', x)[-1])
            )
            for idx, sf in enumerate(slide_files, 1):
                parts = [
                    t.text.strip()
                    for t in ET.fromstring(z.read(sf)).iter(
                        '{http://schemas.openxmlformats.org/drawingml/2006/main}t'
                    )
                    if t.text and t.text.strip()
                ]
                if parts:
                    slides.append({"slide_num": idx, "title": parts[0][:80], "content": "\n".join(parts)})
        if not slides:
            raise ValueError("No readable content found in this file")
        return slides

    slides = []
    for i, slide in enumerate(prs.slides, 1):
        ts    = slide.shapes.title
        title = ts.text.strip() if ts and ts.text.strip() else f"Slide {i}"
        body  = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape != ts and shape.text.strip()
        ]
        slides.append({"slide_num": i, "title": title, "content": "\n".join(body)})
    return slides


def _extract_ppt_binary(raw):
    """Extract text from old binary OLE2 .ppt via record-level parsing."""
    import struct
    try:
        import olefile
    except ImportError:
        raise ValueError("olefile not installed — run: pip install olefile")

    ole = olefile.OleFileIO(io.BytesIO(raw))
    if not ole.exists('PowerPoint Document'):
        raise ValueError("Not a valid binary PowerPoint file")

    stream = ole.openstream('PowerPoint Document').read()
    texts  = []

    def parse(buf, start, end):
        i = start
        while i + 8 <= end:
            rec_ver  = struct.unpack_from('<H', buf, i)[0] & 0x0F
            rec_type = struct.unpack_from('<H', buf, i + 2)[0]
            rec_len  = struct.unpack_from('<I', buf, i + 4)[0]
            data_end = i + 8 + rec_len
            if data_end > end:
                break
            if rec_type == 0x0FA0:       # TextCharsAtom — UTF-16LE
                t = buf[i+8:data_end].decode('utf-16-le', errors='ignore').strip()
                if t: texts.append(t)
            elif rec_type == 0x0FA8:     # TextBytesAtom — Latin-1
                t = buf[i+8:data_end].decode('latin-1', errors='ignore').strip()
                if t: texts.append(t)
            elif rec_ver == 0xF:         # Container — recurse into children
                parse(buf, i + 8, data_end)
            i = data_end

    parse(stream, 0, len(stream))

    if not texts:
        raise ValueError("No readable text found in the binary .ppt file")

    # Group into pseudo-slides of ~5 text blocks
    slides = []
    for n, chunk in enumerate([texts[j:j+5] for j in range(0, len(texts), 5)], 1):
        slides.append({"slide_num": n, "title": chunk[0][:80], "content": "\n".join(chunk)})
    return slides


def _extract_pdf(raw):
    """Extract text from PDF, one entry per page."""
    try:
        import pypdf
    except ImportError:
        raise ValueError("pypdf not installed — run: pip install pypdf")

    reader = pypdf.PdfReader(io.BytesIO(raw))
    slides = []
    for i, page in enumerate(reader.pages, 1):
        text  = (page.extract_text() or "").strip()
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        if lines:
            slides.append({"slide_num": i, "title": lines[0][:80], "content": "\n".join(lines)})

    if not slides:
        raise ValueError(
            "No readable text found in the PDF — it may be a scanned/image-based file."
        )
    return slides


def _extract_docx(raw):
    """Extract text from .docx (Open XML Word) files, grouped by headings."""
    try:
        from docx import Document
    except ImportError:
        raise ValueError("python-docx not installed — run: pip install python-docx")

    doc = Document(io.BytesIO(raw))
    slides, current_title, current_lines, slide_num = [], "Document", [], 1
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style.name.startswith('Heading'):
            if current_lines:
                slides.append({"slide_num": slide_num, "title": current_title, "content": "\n".join(current_lines)})
                slide_num += 1
                current_lines = []
            current_title = text[:80]
        else:
            current_lines.append(text)
    if current_lines:
        slides.append({"slide_num": slide_num, "title": current_title, "content": "\n".join(current_lines)})
    if not slides:
        raise ValueError("No readable text found in the Word document.")
    return slides


def _extract_doc_ole(raw):
    """Extract text from old binary .doc via OLE2 stream decoding."""
    import olefile
    ole = olefile.OleFileIO(io.BytesIO(raw))
    if not ole.exists('WordDocument'):
        raise ValueError("Not a valid binary Word (.doc) file.")
    stream = ole.openstream('WordDocument').read()
    # Word stores main text as UTF-16-LE; extract printable runs
    try:
        text = stream.decode('utf-16-le', errors='ignore')
    except Exception:
        text = stream.decode('latin-1', errors='ignore')
    # Keep only printable chars + Arabic/newlines, strip control chars
    text = re.sub(r'[^\x20-\x7E؀-ۿÀ-ɏ\n\r\t]+', ' ', text)
    text = re.sub(r'[ \t]{3,}', '  ', text).strip()
    lines = [l.strip() for l in text.splitlines() if len(l.strip()) > 3]
    if not lines:
        raise ValueError(
            "Could not extract readable text from this .doc file. "
            "Try saving it as .docx and uploading again."
        )
    # Group lines into pseudo-slides of 15 lines each
    slides = []
    for n, chunk in enumerate([lines[i:i+15] for i in range(0, len(lines), 15)], 1):
        slides.append({"slide_num": n, "title": chunk[0][:80], "content": "\n".join(chunk)})
    return slides


def _extract_txt(raw):
    """Extract text from a plain-text file (.txt)."""
    for enc in ('utf-8', 'utf-16', 'latin-1', 'cp1256'):
        try:
            text = raw.decode(enc, errors='strict')
            break
        except (UnicodeDecodeError, LookupError):
            continue
    else:
        text = raw.decode('utf-8', errors='replace')
    text = text.strip()
    if not text:
        raise ValueError("The text file appears to be empty.")
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    # Group into pseudo-slides of 20 lines
    slides = []
    for n, chunk in enumerate([lines[i:i+20] for i in range(0, len(lines), 20)], 1):
        slides.append({"slide_num": n, "title": chunk[0][:80], "content": "\n".join(chunk)})
    return slides


def extract_slides(file_stream, filename=""):
    """Detect format from magic bytes (and filename for .txt) and dispatch."""
    raw = file_stream.read()
    fname = (filename or "").lower()

    if raw[:4] == b'PK\x03\x04':                         # ZIP-based (pptx or docx)
        import zipfile
        try:
            with zipfile.ZipFile(io.BytesIO(raw)) as z:
                names = z.namelist()
            if any(n.startswith('word/') for n in names):
                return _extract_docx(raw)
        except Exception:
            pass
        return _extract_pptx_raw(raw)

    if raw[:8] == b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1':  # OLE2 (ppt or doc)
        try:
            import olefile
            with olefile.OleFileIO(io.BytesIO(raw)) as ole:
                if ole.exists('WordDocument'):
                    return _extract_doc_ole(raw)
        except Exception:
            pass
        return _extract_ppt_binary(raw)

    if raw[:4] == b'%PDF':                                # PDF
        return _extract_pdf(raw)

    if fname.endswith('.txt') or fname.endswith('.md'):   # Plain text
        return _extract_txt(raw)

    # Last resort: try decoding as plain text
    try:
        decoded = raw.decode('utf-8', errors='strict')
        if len(decoded.strip()) > 50:
            return _extract_txt(raw)
    except UnicodeDecodeError:
        pass

    raise ValueError(
        "Unrecognised file format. Supported: .pptx, .ppt, .pdf, .docx, .doc, .txt"
    )


# ── PDF builder ────────────────────────────────────────────────────────────────

def _st(name, parent, **kw):
    s = ParagraphStyle(name, parent=parent)
    for k, v in kw.items(): setattr(s, k, v)
    return s

def _page_num(canvas, doc):
    canvas.saveState()
    canvas.setFont('Helvetica', 7.5)
    canvas.setFillColor(TEXT_LIGHT)
    n = canvas.getPageNumber()
    canvas.drawRightString(doc.width + doc.leftMargin, 0.55*cm, f"{n}")
    canvas.drawString(doc.leftMargin, 0.55*cm, doc.title_str if hasattr(doc, 'title_str') else '')
    canvas.restoreState()


def build_pdf(guide, language, out_filename="study_guide"):
    if not isinstance(guide, dict):
        guide = {}
    guide["sections"]   = [s for s in guide.get("sections",   []) if isinstance(s, dict)]
    guide["keywords"]   = [k for k in guide.get("keywords",   []) if isinstance(k, dict)]
    guide["flashcards"] = [f for f in guide.get("flashcards", []) if isinstance(f, dict)]
    guide["objectives"] = [o for o in guide.get("objectives", []) if isinstance(o, str)]

    is_ar = (language == "ar")
    ar_ok = is_ar and _ensure_arabic_font()
    AF    = _ARABIC_FONT if ar_ok else "Helvetica"
    AFB   = _ARABIC_FONT if ar_ok else "Helvetica-Bold"
    ALIGN = TA_RIGHT if is_ar else TA_LEFT

    def T(text):
        return _ar(str(text)) if ar_ok else str(text)

    L = {
        "objectives": T("الأهداف التعليمية") if is_ar else "◆  LEARNING OBJECTIVES",
        "obj_bullet": "• " if is_ar else "◆  ",
        "contents":   T("المحتويات") if is_ar else "CONTENTS",
        "kw_head":    T("قاموس المصطلحات") if is_ar else "■  KEYWORDS CHEATSHEET",
        "kw_append":  [(T("قاموس المصطلحات"), ""), (T("بطاقات المراجعة"), "")] if is_ar
                      else [("KEYWORDS CHEATSHEET", ""), ("FLASH CARDS", "")],
        "fc_head":    T("بطاقات المراجعة") if is_ar else "■  FLASH CARDS",
        "sec_bullet": "" if is_ar else "■  ",
        "bul_bullet": "• " if is_ar else "▸  ",
        "q_pre":      T("سؤال: ") if is_ar else "Q:  ",
        "guide":      T("دليل الدراسة بالذكاء الاصطناعي") if is_ar else "AI Exam Study Guide",
        "luck":       T("حظ سعيد!") if is_ar else "Good luck!",
    }

    buf = io.BytesIO()
    W   = 17.4*cm

    class Doc(BaseDocTemplate):
        pass

    doc = Doc(buf, pagesize=A4,
              leftMargin=1.8*cm, rightMargin=1.8*cm,
              topMargin=1.8*cm, bottomMargin=1.5*cm)
    doc.title_str = guide.get("title", "")

    frame = Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height - 0.3*cm,
                  id='main', leftPadding=0, rightPadding=0, topPadding=0, bottomPadding=0)
    doc.addPageTemplates([PageTemplate(id='main', frames=[frame], onPage=_page_num)])

    base = getSampleStyleSheet()["Normal"]
    ST = {
        "h_title":  _st("HT",  base, fontSize=20, fontName=AFB, textColor=WHITE,        alignment=TA_CENTER, leading=26),
        "h_sub":    _st("HS",  base, fontSize=9,  fontName=AF,  textColor=colors.HexColor("#a8c4e8"), alignment=TA_CENTER),
        "obj_head": _st("OH",  base, fontSize=11, fontName=AFB, textColor=NAVY,         spaceBefore=4, spaceAfter=4, alignment=ALIGN),
        "obj_item": _st("OI",  base, fontSize=9.5,fontName=AF,  textColor=TEXT,         leftIndent=0 if is_ar else 14, rightIndent=14 if is_ar else 0, spaceAfter=3, leading=16, alignment=ALIGN),
        "toc_title":_st("TOT", base, fontSize=11, fontName=AFB, textColor=NAVY,         spaceAfter=6, alignment=ALIGN),
        "toc_item": _st("TOI", base, fontSize=9.5,fontName=AF,  textColor=TEXT,         leftIndent=0 if is_ar else 10, rightIndent=10 if is_ar else 0, spaceAfter=2, alignment=ALIGN),
        "sec_title":_st("SCT", base, fontSize=11, fontName=AFB, textColor=WHITE,        alignment=ALIGN),
        "bullet":   _st("BL",  base, fontSize=9.5,fontName=AF,  textColor=TEXT,         leftIndent=0 if is_ar else 12, rightIndent=12 if is_ar else 0, spaceAfter=3, leading=16, alignment=ALIGN),
        "tbl_hdr":  _st("TH",  base, fontSize=9,  fontName=AFB, textColor=WHITE,        alignment=ALIGN),
        "tbl_cell": _st("TC",  base, fontSize=9,  fontName=AF,  textColor=TEXT,         leading=14, alignment=ALIGN),
        "kw_term":  _st("KT",  base, fontSize=9,  fontName=AFB, textColor=NAVY_MID,     alignment=ALIGN),
        "kw_def":   _st("KD",  base, fontSize=9,  fontName=AF,  textColor=TEXT,         leading=14, alignment=ALIGN),
        "kw_head":  _st("KH",  base, fontSize=11, fontName=AFB, textColor=WHITE,        alignment=TA_CENTER),
        "fc_q":     _st("FCQ", base, fontSize=9,  fontName=AFB, textColor=WHITE,        leading=14, alignment=ALIGN),
        "fc_a":     _st("FCA", base, fontSize=9,  fontName=AF,  textColor=TEXT,         leading=14, alignment=ALIGN),
        "fc_head":  _st("FCH", base, fontSize=11, fontName=AFB, textColor=WHITE,        alignment=TA_CENTER),
        "footer":   _st("FT",  base, fontSize=8,  fontName=AF,  textColor=TEXT_LIGHT,   alignment=TA_CENTER),
    }

    elems = []

    # ── Header ────────────────────────────────────────────────────────────────
    raw_title = guide.get("title", "Study Guide")
    title    = T(raw_title.upper() if not is_ar else raw_title)
    subtitle = T(guide.get("subtitle", "Exam Study Guide"))
    hdr = Table([
        [Paragraph(title, ST["h_title"])],
        [Paragraph(f"{subtitle}  ·  {OLLAMA_MODEL}", ST["h_sub"])],
    ], colWidths=[W])
    hdr.setStyle(TableStyle([
        ("BACKGROUND",    (0,0), (-1,-1), NAVY),
        ("ALIGN",         (0,0), (-1,-1), "CENTER"),
        ("TOPPADDING",    (0,0), (0,0),   16),
        ("BOTTOMPADDING", (0,1), (0,1),   14),
        ("TOPPADDING",    (0,1), (0,1),   2),
        ("LEFTPADDING",   (0,0), (-1,-1), 12),
        ("RIGHTPADDING",  (0,0), (-1,-1), 12),
    ]))
    elems.append(hdr)
    elems.append(Spacer(1, 0.3*cm))

    # ── Learning Objectives ───────────────────────────────────────────────────
    objectives = guide.get("objectives", [])
    if objectives:
        rows = [[Paragraph(L["objectives"], ST["obj_head"])]]
        for o in objectives:
            rows.append([Paragraph(f"{L['obj_bullet']}{T(o)}", ST["obj_item"])])
        t = Table(rows, colWidths=[W])
        t.setStyle(TableStyle([
            ("BACKGROUND",    (0,0), (-1,0),  SECTION_BG),
            ("BOX",           (0,0), (-1,-1), 0.8, BORDER),
            ("LINEBELOW",     (0,0), (-1,0),  0.8, BORDER),
            ("TOPPADDING",    (0,0), (-1,-1), 5),
            ("BOTTOMPADDING", (0,0), (-1,-1), 5),
            ("LEFTPADDING",   (0,0), (-1,-1), 10),
            ("RIGHTPADDING",  (0,0), (-1,-1), 10),
        ]))
        elems.append(t)
        elems.append(Spacer(1, 0.3*cm))

    # ── Table of Contents ─────────────────────────────────────────────────────
    sections = guide.get("sections", [])
    if sections:
        toc_rows = [[Paragraph(L["contents"], ST["toc_title"])]]
        all_items = [(T(s["title"]), "") for s in sections] + L["kw_append"]
        for i, (name, _) in enumerate(all_items, 1):
            dot_row = Table(
                [[Paragraph(f"{i}.  {name}", ST["toc_item"]), Paragraph("", ST["toc_item"])]],
                colWidths=[W*0.85, W*0.15]
            )
            dot_row.setStyle(TableStyle([
                ("LINEBELOW",     (0,0), (-1,-1), 0.3, BORDER),
                ("TOPPADDING",    (0,0), (-1,-1), 3),
                ("BOTTOMPADDING", (0,0), (-1,-1), 3),
                ("LEFTPADDING",   (0,0), (-1,-1), 6),
                ("RIGHTPADDING",  (0,0), (-1,-1), 6),
            ]))
            toc_rows.append([dot_row])
        toc = Table(toc_rows, colWidths=[W])
        toc.setStyle(TableStyle([
            ("BACKGROUND",    (0,0), (-1,0),  SECTION_BG),
            ("BOX",           (0,0), (-1,-1), 0.8, BORDER),
            ("TOPPADDING",    (0,0), (-1,-1), 4),
            ("BOTTOMPADDING", (0,0), (-1,-1), 4),
            ("LEFTPADDING",   (0,0), (-1,-1), 0),
        ]))
        elems.append(toc)
        elems.append(Spacer(1, 0.35*cm))

    # ── Sections ──────────────────────────────────────────────────────────────
    for idx, sec in enumerate(sections, 1):
        block = []

        sec_title_text = T(sec.get("title", ""))
        sec_hdr = Table(
            [[Paragraph(
                f"{L['sec_bullet']}{idx} · {sec_title_text if is_ar else sec_title_text.upper()}",
                ST["sec_title"]
            )]],
            colWidths=[W]
        )
        sec_hdr.setStyle(TableStyle([
            ("BACKGROUND",    (0,0), (-1,-1), NAVY_MID),
            ("TOPPADDING",    (0,0), (-1,-1), 7),
            ("BOTTOMPADDING", (0,0), (-1,-1), 7),
            ("LEFTPADDING",   (0,0), (-1,-1), 10),
            ("RIGHTPADDING",  (0,0), (-1,-1), 10),
        ]))
        block.append(sec_hdr)

        bullets = sec.get("bullets", [])
        if bullets:
            bdata = [[Paragraph(f"{L['bul_bullet']}{T(b)}", ST["bullet"])] for b in bullets]
            bt = Table(bdata, colWidths=[W])
            bt.setStyle(TableStyle([
                ("TOPPADDING",    (0,0), (-1,-1), 5),
                ("BOTTOMPADDING", (0,0), (-1,-1), 4),
                ("LEFTPADDING",   (0,0), (-1,-1), 10),
                ("RIGHTPADDING",  (0,0), (-1,-1), 10),
                ("LINEBELOW",     (0,0), (-1,-2), 0.3, colors.HexColor("#dde8ff")),
                ("BOX",           (0,0), (-1,-1), 0.5, BORDER),
            ]))
            block.append(bt)

        tbl = sec.get("table")
        if isinstance(tbl, dict) and tbl.get("headers") and tbl.get("rows"):
            headers = tbl["headers"]
            n_cols  = len(headers)
            col_w   = W / n_cols
            tbl_rows = [[Paragraph(T(h), ST["tbl_hdr"]) for h in headers]]
            for ri, row in enumerate(tbl["rows"]):
                padded = (list(row) + [""] * n_cols)[:n_cols]
                tbl_rows.append([Paragraph(T(str(c)), ST["tbl_cell"]) for c in padded])
            inner = Table(tbl_rows, colWidths=[col_w]*n_cols)
            ts = [
                ("BACKGROUND",    (0,0), (-1,0),  NAVY_LIGHT),
                ("TOPPADDING",    (0,0), (-1,-1), 5),
                ("BOTTOMPADDING", (0,0), (-1,-1), 5),
                ("LEFTPADDING",   (0,0), (-1,-1), 7),
                ("GRID",          (0,0), (-1,-1), 0.4, BORDER),
            ]
            for ri in range(1, len(tbl_rows)):
                if ri % 2 == 0:
                    ts.append(("BACKGROUND", (0,ri), (-1,ri), ROW_ALT))
            inner.setStyle(TableStyle(ts))
            block.append(Spacer(1, 0.15*cm))
            block.append(inner)

        block.append(Spacer(1, 0.3*cm))
        elems.append(KeepTogether(block))

    # ── Keywords Cheatsheet ───────────────────────────────────────────────────
    keywords = guide.get("keywords", [])
    if keywords:
        kw_hdr = Table(
            [[Paragraph(L["kw_head"], ST["kw_head"])]],
            colWidths=[W]
        )
        kw_hdr.setStyle(TableStyle([
            ("BACKGROUND",    (0,0), (-1,-1), NAVY),
            ("TOPPADDING",    (0,0), (-1,-1), 8),
            ("BOTTOMPADDING", (0,0), (-1,-1), 8),
            ("LEFTPADDING",   (0,0), (-1,-1), 10),
        ]))
        elems.append(kw_hdr)

        if is_ar:
            # Arabic: definition left, term right (visual RTL order)
            lc, rc = W*0.60, W*0.40
            kw_rows = [[
                Paragraph(T(k.get("definition", "")), ST["kw_def"]),
                Paragraph(T(k.get("term", "")),       ST["kw_term"]),
            ] for k in keywords]
        else:
            lc, rc = W*0.27, W*0.73
            kw_rows = [[
                Paragraph(k.get("term", ""),       ST["kw_term"]),
                Paragraph(k.get("definition", ""), ST["kw_def"]),
            ] for k in keywords]
        kw_t = Table(kw_rows, colWidths=[lc, rc])
        kts = [
            ("VALIGN",        (0,0), (-1,-1), "TOP"),
            ("TOPPADDING",    (0,0), (-1,-1), 5),
            ("BOTTOMPADDING", (0,0), (-1,-1), 5),
            ("LEFTPADDING",   (0,0), (-1,-1), 8),
            ("RIGHTPADDING",  (0,0), (-1,-1), 8),
            ("LINEBELOW",     (0,0), (-1,-2), 0.3, BORDER),
            ("BOX",           (0,0), (-1,-1), 0.5, BORDER),
        ]
        for ri in range(len(kw_rows)):
            if ri % 2 == 0:
                kts.append(("BACKGROUND", (0,ri), (-1,ri), KW_BG))
        kw_t.setStyle(TableStyle(kts))
        elems.append(kw_t)
        elems.append(Spacer(1, 0.35*cm))

    # ── Flash Cards ───────────────────────────────────────────────────────────
    flashcards = guide.get("flashcards", [])
    if flashcards:
        fc_hdr = Table(
            [[Paragraph(L["fc_head"], ST["fc_head"])]],
            colWidths=[W]
        )
        fc_hdr.setStyle(TableStyle([
            ("BACKGROUND",    (0,0), (-1,-1), NAVY),
            ("TOPPADDING",    (0,0), (-1,-1), 8),
            ("BOTTOMPADDING", (0,0), (-1,-1), 8),
            ("LEFTPADDING",   (0,0), (-1,-1), 10),
        ]))
        elems.append(fc_hdr)
        elems.append(Spacer(1, 0.2*cm))

        cw = (W - 0.3*cm) / 2
        pairs = [flashcards[i:i+2] for i in range(0, len(flashcards), 2)]
        for pair in pairs:
            row_cells = []
            for fc in pair:
                card = Table([
                    [Paragraph(f"{L['q_pre']}{T(fc.get('q',''))}", ST["fc_q"])],
                    [Paragraph(T(fc.get('a','')), ST["fc_a"])],
                ], colWidths=[cw])
                card.setStyle(TableStyle([
                    ("BACKGROUND",    (0,0), (-1,0),  CARD_Q),
                    ("BACKGROUND",    (0,1), (-1,1),  CARD_A),
                    ("TOPPADDING",    (0,0), (-1,-1), 7),
                    ("BOTTOMPADDING", (0,0), (-1,-1), 7),
                    ("LEFTPADDING",   (0,0), (-1,-1), 9),
                    ("RIGHTPADDING",  (0,0), (-1,-1), 9),
                    ("BOX",           (0,0), (-1,-1), 0.8, BORDER),
                    ("LINEBELOW",     (0,0), (-1,0),  0.5, BORDER),
                ]))
                row_cells.append(card)
            if len(row_cells) == 1:
                row_cells.append(Spacer(cw, 1))
            grid_row = Table([row_cells], colWidths=[cw, cw], hAlign='LEFT')
            grid_row.setStyle(TableStyle([("LEFTPADDING",(0,0),(-1,-1),0),("RIGHTPADDING",(0,0),(-1,-1),0)]))
            elems.append(grid_row)
            elems.append(Spacer(1, 0.2*cm))

    # ── Footer ────────────────────────────────────────────────────────────────
    elems.append(HRFlowable(width="100%", thickness=0.5, color=BORDER))
    elems.append(Spacer(1, 0.1*cm))
    elems.append(Paragraph(
        f"{T(guide.get('title',''))}  ·  {L['guide']}  ·  {OLLAMA_MODEL}  ·  {L['luck']}",
        ST["footer"]
    ))

    doc.multiBuild(elems)
    buf.seek(0)
    return buf


# ── Public config endpoint ────────────────────────────────────────────────────
@app.route("/api/config")
def api_config():
    """Return public keys the frontend needs to initialise Supabase and Stripe."""
    return jsonify({
        "supabase_url":          SUPABASE_URL,
        "supabase_anon_key":     SUPABASE_ANON_KEY,
        "stripe_publishable_key": STRIPE_PUBLISHABLE_KEY,
        "auth_enabled":          _AUTH_ENABLED,
    })


# ── Auth — current user ────────────────────────────────────────────────────────
@app.route("/api/auth/me")
def auth_me():
    uid, err = _auth_check(request)
    if err:
        return err
    if uid == "dev":
        return jsonify({"email": "dev@local", "name": "Dev", "tokens_remaining": 999,
                        "subscription_status": "active", "referral_code": "DEVLOCAL"})
    user = _get_user(uid)
    if not user:
        return jsonify({"error": "User not found"}), 404
    ref_code = user.get("referral_code") or _get_or_create_referral_code(uid)
    return jsonify({
        "id":                      user["id"],
        "email":                   user["email"],
        "name":                    user.get("name") or "",
        "avatar_url":              user.get("avatar_url") or "",
        "tokens_remaining":        user.get("tokens_remaining", 0),
        "subscription_status":     user.get("subscription_status", "free"),
        "subscription_period_end": str(user.get("subscription_period_end") or ""),
        "referral_code":           ref_code or "",
    })


# ── Referral — apply code ─────────────────────────────────────────────────────
@app.route("/api/referral/apply", methods=["POST"])
def referral_apply():
    uid, err = _auth_check(request)
    if err:
        return err
    if uid == "dev":
        return jsonify({"success": False, "reason": "dev_mode"})
    code = (request.get_json(silent=True) or {}).get("code", "").strip().upper()
    if not code:
        return jsonify({"success": False, "reason": "no_code"}), 400
    sb = _get_sb()
    if not sb:
        return jsonify({"success": False, "reason": "unavailable"})
    try:
        # Find referrer by code (cannot self-refer)
        ref = sb.table("users").select("id").eq("referral_code", code).neq("id", uid).execute()
        if not ref.data:
            return jsonify({"success": False, "reason": "invalid_code"})
        referrer_id = ref.data[0]["id"]
        # Apply only if not already referred
        sb.table("users").update({"referred_by": referrer_id}).eq("id", uid).is_("referred_by", "null").execute()
        # Remove stored code from client regardless (avoid re-tries)
        return jsonify({"success": True})
    except Exception as exc:
        _log.error(f"referral_apply error: {exc}")
        return jsonify({"success": False, "reason": "db_error"}), 500


# ── Referral — stats ───────────────────────────────────────────────────────────
@app.route("/api/referral/stats")
def referral_stats():
    uid, err = _auth_check(request)
    if err:
        return err
    if uid == "dev":
        return jsonify({"total": 0, "paid": 0, "tokens_earned": 0})
    sb = _get_sb()
    if not sb:
        return jsonify({"total": 0, "paid": 0, "tokens_earned": 0})
    try:
        rows = sb.table("users").select("referral_paid").eq("referred_by", uid).execute()
        total = len(rows.data) if rows.data else 0
        paid  = sum(1 for r in (rows.data or []) if r.get("referral_paid"))
        return jsonify({"total": total, "paid": paid, "tokens_earned": paid * 10})
    except Exception:
        return jsonify({"total": 0, "paid": 0, "tokens_earned": 0})


# ── Stripe — create checkout session ──────────────────────────────────────────
@app.route("/api/stripe/checkout", methods=["POST"])
def stripe_checkout():
    uid, err = _auth_check(request)
    if err:
        return err
    if not _stripe or not STRIPE_PRICE_ID:
        return jsonify({"error": "Payments not configured"}), 503

    user = _get_user(uid)
    if not user:
        return jsonify({"error": "User not found"}), 404

    try:
        # Reuse existing Stripe customer or create a new one
        customer_id = user.get("stripe_customer_id")
        if not customer_id:
            cust = _stripe.Customer.create(
                email=user["email"],
                name=user.get("name", ""),
                metadata={"supabase_id": uid},
            )
            customer_id = cust["id"]
            _get_sb().table("users").update(
                {"stripe_customer_id": customer_id}
            ).eq("id", uid).execute()

        session = _stripe.checkout.Session.create(
            customer=customer_id,
            payment_method_types=["card"],
            line_items=[{"price": STRIPE_PRICE_ID, "quantity": 1}],
            mode="subscription",
            success_url=f"{APP_URL}/?sub=success",
            cancel_url=f"{APP_URL}/?sub=canceled",
            metadata={"user_id": uid},
        )
        return jsonify({"url": session.url})
    except Exception as exc:
        return jsonify({"error": str(exc)}), 500


# ── Stripe — customer billing portal ──────────────────────────────────────────
@app.route("/api/stripe/portal", methods=["POST"])
def stripe_portal():
    uid, err = _auth_check(request)
    if err:
        return err
    if not _stripe:
        return jsonify({"error": "Payments not configured"}), 503

    user = _get_user(uid)
    if not user or not user.get("stripe_customer_id"):
        return jsonify({"error": "No billing account found"}), 404

    try:
        portal = _stripe.billing_portal.Session.create(
            customer=user["stripe_customer_id"],
            return_url=APP_URL,
        )
        return jsonify({"url": portal.url})
    except Exception as exc:
        return jsonify({"error": str(exc)}), 500


# ── Stripe — webhook ───────────────────────────────────────────────────────────
@app.route("/api/stripe/webhook", methods=["POST"])
def stripe_webhook():
    if not _stripe:
        return jsonify({"error": "Payments not configured"}), 503

    payload   = request.get_data()
    sig       = request.headers.get("Stripe-Signature", "")
    try:
        event = _stripe.Webhook.construct_event(payload, sig, STRIPE_WEBHOOK_SECRET)
    except ValueError:
        return jsonify({"error": "Invalid payload"}), 400
    except _stripe.error.SignatureVerificationError:
        return jsonify({"error": "Invalid signature"}), 400

    sb  = _get_sb()
    typ = event["type"]

    if typ == "checkout.session.completed":
        sess    = event["data"]["object"]
        user_id = (sess.get("metadata") or {}).get("user_id")
        if user_id and sb:
            sb.table("users").update({
                "subscription_status": "active",
                "tokens_remaining":    20,
                "tokens_month":        time.strftime("%Y-%m"),
            }).eq("id", user_id).execute()
            _award_referral(user_id, sb)   # reward referrer if applicable

    elif typ in ("customer.subscription.updated", "customer.subscription.deleted"):
        sub         = event["data"]["object"]
        cust_id     = sub.get("customer")
        status      = sub.get("status", "")
        period_end  = sub.get("current_period_end")
        if sb and cust_id:
            update = {"subscription_id": sub.get("id", "")}
            if status == "active":
                update["subscription_status"] = "active"
                if period_end:
                    from datetime import datetime, timezone
                    update["subscription_period_end"] = datetime.fromtimestamp(
                        period_end, tz=timezone.utc
                    ).isoformat()
            elif status in ("canceled", "unpaid", "past_due"):
                update["subscription_status"] = status
            sb.table("users").update(update).eq("stripe_customer_id", cust_id).execute()

    elif typ == "invoice.payment_succeeded":
        inv     = event["data"]["object"]
        cust_id = inv.get("customer")
        if inv.get("billing_reason") == "subscription_cycle" and sb and cust_id:
            sb.table("users").update({
                "tokens_remaining": 20,
                "tokens_month":     time.strftime("%Y-%m"),
            }).eq("stripe_customer_id", cust_id).execute()

    return jsonify({"ok": True})


# ── SSE streaming endpoint ─────────────────────────────────────────────────────

def _sse(data):
    return f"data: {json.dumps(data)}\n\n"

@app.route("/api/summarize-stream", methods=["POST"])
def summarize_stream():
    # ── Per-IP rate limit (checked first, before any other validation) ─────────
    if not _check_rate_limit(_client_ip()):
        return jsonify({"error": "Too many requests. Please wait a minute before trying again."}), 429

    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    # ── Auth + token gate ──────────────────────────────────────────────────────
    uid, err = _auth_check(request)
    if err:
        return err
    ok, tok_left, reason = _consume_token(uid)
    if not ok:
        return jsonify({
            "error": "You have no tokens left. Upgrade to continue.",
            "code": "no_tokens", "tokens_remaining": 0
        }), 402

    f            = request.files["file"]
    lang_param   = request.form.get("language", "auto")
    out_name     = _safe_name(request.form.get("filename", f.filename.rsplit(".", 1)[0]))
    detail_level = request.form.get("detail", "standard")
    dcfg         = DETAIL.get(detail_level, DETAIL["standard"])

    _ALLOWED_EXT = (".pptx", ".ppt", ".pdf", ".docx", ".doc", ".txt")
    if not f.filename.lower().endswith(_ALLOWED_EXT):
        return jsonify({"error": "Unsupported file type. Supported: .pptx, .ppt, .pdf, .docx, .doc, .txt"}), 400

    if not ollama_running():
        return jsonify({"error": "AI service is not configured. Set GROQ_API_KEY."}), 503

    file_bytes = f.read()

    def generate():
        try:
            yield _sse({"step": "extract", "msg": "Extracting content…"})
            slides = extract_slides(io.BytesIO(file_bytes), filename=f.filename)
            if not any(s["content"] or s["title"] for s in slides):
                yield _sse({"error": "No readable content in this file"}); return

            total = len([s for s in slides if s["content"].strip()])

            # Auto-detect language from slide content
            language = lang_param if lang_param in ("ar", "en") else _detect_language(slides)
            lang_label = "Arabic" if language == "ar" else "English"
            yield _sse({"step": "extract", "msg": f"Found {total} content slides ({lang_label}) — analysing…", "language": language})

            # Pass 1
            yield _sse({"step": "overview", "msg": "Analysing structure and keywords…"})
            overview = pass1_overview(
                [s for s in slides if s["content"].strip() or s["title"].strip()],
                language, dcfg
            )

            # Normalize: mistral sometimes returns sections/keywords as plain strings
            raw_sections = overview.get("sections", [])
            overview["sections"] = [
                s if isinstance(s, dict) else {"title": str(s), "slide_nums": []}
                for s in (raw_sections if isinstance(raw_sections, list) else [])
            ]
            raw_keywords = overview.get("keywords", [])
            overview["keywords"] = [
                k if isinstance(k, dict) else {"term": str(k), "definition": ""}
                for k in (raw_keywords if isinstance(raw_keywords, list) else [])
            ]

            content_slides = [s for s in slides if s["content"].strip()]

            # Fallback: if no sections produced, wrap all content into one
            if not overview["sections"]:
                overview["sections"] = [{
                    "title": overview.get("title", "Content Overview"),
                    "slide_nums": [s["slide_num"] for s in content_slides]
                }]

            sections = overview["sections"]

            for evt in _sections_parallel(sections, content_slides, language, dcfg):
                yield _sse(evt)

            for evt in _flashcards_mcq_parallel(overview, language, dcfg):
                yield _sse(evt)

            # Build PDF + Markdown
            yield _sse({"step": "pdf", "msg": "Building PDF & Markdown…"})
            pdf_buf   = build_pdf(overview, language, out_name)
            pdf_bytes = pdf_buf.read()
            md_text   = build_markdown(overview)

            job_id = uuid.uuid4().hex
            store_job(job_id, pdf_bytes, md_text, overview, slides, f"{out_name}_study_guide.pdf")

            yield _sse({"step": "done", "job_id": job_id,
                        "sections":   len(sections),
                        "keywords":   len(overview.get("keywords",   [])),
                        "flashcards": len(overview.get("flashcards", [])),
                        "mcqs":       len(overview.get("mcqs",       []))})

        except Exception as e:
            _log.error("GENERATE_ERROR: %s\n%s", e, _tb.format_exc())
            yield _sse({"error": _safe_err(e)})

    return Response(
        stream_with_context(generate()),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"}
    )


@app.route("/api/download/<job_id>")
def download_job(job_id):
    if not _valid_job(job_id):
        return jsonify({"error": "Invalid job ID"}), 400
    fmt      = request.args.get("format", "pdf")
    filename = _safe_name(request.args.get("filename", "study_guide"))
    with _jobs_lock:
        job = get_job(job_id)   # keep in store — TTL handles cleanup
    if not job:
        return jsonify({"error": "File not found or expired"}), 404
    if fmt == "md":
        content = (job.get("md") or "").encode("utf-8")
        return send_file(io.BytesIO(content), mimetype="text/markdown",
                         as_attachment=True, download_name=f"{filename}.md")
    return send_file(io.BytesIO(job["pdf"]), mimetype="application/pdf",
                     as_attachment=True, download_name=job.get("filename", f"{filename}.pdf"))


@app.route("/api/guide/<job_id>")
def get_guide(job_id):
    if not _valid_job(job_id):
        return jsonify({"error": "Invalid job ID"}), 400
    with _jobs_lock:
        job = get_job(job_id)
    if not job:
        return jsonify({"error": "Session expired — re-upload the file"}), 404
    guide = job.get("guide", {})
    return jsonify({
        "title":      guide.get("title", ""),
        "flashcards": guide.get("flashcards", []),
        "mcqs":       guide.get("mcqs", []),
        "keywords":   guide.get("keywords",   []),
        "objectives": guide.get("objectives", []),
    })


@app.route("/api/chat/<job_id>", methods=["POST"])
def chat_with_slides(job_id):
    if not _check_rate_limit(_client_ip(), scope="chat", limit=20):
        return jsonify({"error": "Too many requests. Please wait a minute."}), 429
    uid, err = _auth_check(request)
    if err:
        return err
    if not _valid_job(job_id):
        return jsonify({"error": "Invalid job ID"}), 400
    data     = request.json or {}
    question = data.get("question", "")[:500].strip()
    language = "ar" if data.get("language") == "ar" else "en"
    if not question:
        return jsonify({"error": "No question provided"}), 400
    with _jobs_lock:
        job = get_job(job_id)
    if not job:
        return jsonify({"error": "Session expired — re-upload the file"}), 404
    slides  = job.get("slides", [])
    context = "\n\n".join(
        f"[{s['title']}]\n{s['content'][:400]}"
        for s in slides[:30] if s.get("content")
    )
    lang = "in Arabic" if language == "ar" else "in English"
    try:
        result = _call_ollama(
            f"""Answer this question {lang} using ONLY the study material below.

Material:
{context}

Question: {question}

Return JSON: {{"answer": "your detailed answer"}}

Rules:
- Answer from the material only; if not covered say so clearly
- Be specific and concise
- JSON only""", num_predict=1024)
        return jsonify({"answer": result.get("answer", "No answer found in the slides.")})
    except Exception as e:
        return jsonify({"error": _safe_err(e)}), 500


@app.route("/api/download-zip", methods=["POST"])
def download_zip():
    import zipfile as zf
    job_ids = [jid for jid in (request.json or {}).get("job_ids", []) if _valid_job(jid)]
    buf = io.BytesIO()
    with zf.ZipFile(buf, "w", zf.ZIP_DEFLATED) as z:
        with _jobs_lock:
            for jid in job_ids:
                job = _jobs.get(jid)
                if job and job.get("pdf"):
                    z.writestr(job["filename"], job["pdf"])
    buf.seek(0)
    return send_file(buf, mimetype="application/zip", as_attachment=True,
                     download_name="study_guides.zip")


@app.route("/api/view/md/<job_id>")
def view_md(job_id):
    if not _valid_job(job_id):
        return "<h2 style='font-family:sans-serif;padding:2rem'>Invalid job ID</h2>", 400
    with _jobs_lock:
        job = get_job(job_id)
    if not job:
        return "<h2 style='font-family:sans-serif;padding:2rem'>Guide not found or expired (10 min TTL)</h2>", 404
    title = _he(job["guide"].get("title", "Study Guide"))
    md = job["md"]

    def _md_to_html(text):
        lines, out = text.split('\n'), []
        i = 0
        while i < len(lines):
            l = lines[i]
            if l.startswith('# '):
                out.append(f'<h1>{_he(l[2:])}</h1>')
            elif l.startswith('## '):
                out.append(f'<h2>{_he(l[3:])}</h2>')
            elif l.startswith('### '):
                out.append(f'<h3>{_he(l[4:])}</h3>')
            elif l.startswith('- ') or l.startswith('* '):
                out.append(f'<li>{_inline(_he(l[2:]))}</li>')
            elif l.startswith('| ') and '|' in l[2:]:
                rows, align = [], l
                while i < len(lines) and lines[i].startswith('|'):
                    if not re.match(r'^\|[-| :]+\|$', lines[i]):
                        cells = [c.strip() for c in lines[i].strip('|').split('|')]
                        tag = 'th' if rows == [] else 'td'
                        out.append('<tr>' + ''.join(f'<{tag}>{_inline(_he(c))}</{tag}>' for c in cells) + '</tr>')
                    i += 1
                i -= 1
            elif l.strip() == '':
                out.append('<br>')
            else:
                out.append(f'<p>{_inline(_he(l))}</p>')
            i += 1
        return '\n'.join(out)

    def _inline(t):
        t = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', t)
        t = re.sub(r'\*(.+?)\*',     r'<em>\1</em>',         t)
        t = re.sub(r'`(.+?)`',       r'<code>\1</code>',     t)
        return t

    body = _md_to_html(md)
    return f"""<!DOCTYPE html><html><head><meta charset="UTF-8">
<title>{title}</title>
<style>
  body{{font-family:'Segoe UI',system-ui,sans-serif;max-width:820px;margin:0 auto;padding:2rem;
       background:#f8faff;color:#0a1628;line-height:1.7}}
  h1{{font-size:1.7rem;color:#1a3a6e;border-bottom:2px solid #c5d8ff;padding-bottom:.5rem}}
  h2{{font-size:1.2rem;color:#2e5ca8;margin-top:1.8rem;border-left:4px solid #4f8ef7;padding-left:.75rem}}
  h3{{font-size:1rem;color:#1a3a6e}}
  li{{margin-bottom:.35rem}}
  table{{border-collapse:collapse;width:100%;margin:1rem 0}}
  th{{background:#1a3a6e;color:#fff;padding:8px 12px;text-align:left}}
  td{{padding:7px 12px;border-bottom:1px solid #dde8ff}}
  tr:nth-child(even) td{{background:#f0f5ff}}
  code{{background:#e8f0ff;padding:1px 5px;border-radius:4px;font-size:.9em}}
  strong{{color:#1a3a6e}}
  @media print{{body{{background:#fff}}}}
</style></head><body>
{body}
<hr style="margin-top:2rem;border-color:#c5d8ff">
<p style="font-size:.8rem;color:#8aa0c8;text-align:center">Generated by Alimne (علّمني) · {OLLAMA_MODEL}</p>
</body></html>"""


def _page_shell(title, body_css, body_html, script=""):
    return f"""<!DOCTYPE html><html><head><meta charset="UTF-8">
<title>{title}</title>
<style>
  *{{box-sizing:border-box;margin:0;padding:0}}
  body{{font-family:'Segoe UI',system-ui,sans-serif;background:#050d1a;color:#e8f0ff;min-height:100vh}}
  .top{{background:#0a1628;border-bottom:1px solid #1a3a6e;padding:1rem 2rem;
        display:flex;align-items:center;justify-content:space-between}}
  .top h1{{font-size:1rem;font-weight:700;color:#e8f0ff}}
  .badge{{background:#4f8ef7;color:#fff;padding:2px 10px;border-radius:20px;font-size:12px}}
  .wrap{{max-width:760px;margin:0 auto;padding:2rem 1.25rem}}
  {body_css}
</style></head><body>
<div class="top"><h1>📖 {title}</h1></div>
<div class="wrap">{body_html}</div>
<script>{script}</script>
</body></html>"""


@app.route("/api/view/cards/<job_id>")
def view_cards(job_id):
    if not _valid_job(job_id):
        return "<h2 style='font-family:sans-serif;padding:2rem'>Invalid job ID</h2>", 400
    with _jobs_lock:
        job = get_job(job_id)
    if not job:
        return "<h2 style='font-family:sans-serif;padding:2rem'>Guide not found or expired</h2>", 404
    guide = job["guide"]
    title = _he(guide.get("title", "Flash Cards"))
    cards = [f for f in guide.get("flashcards", []) if isinstance(f, dict)]
    if not cards:
        return _page_shell(title, "", "<p style='text-align:center;color:#4a5f80;padding:3rem'>No flash cards available.</p>")

    cards_json = json.dumps(cards)
    css = """
  .fc-counter{text-align:center;color:#8aa0c8;font-size:.88rem;margin-bottom:1.2rem}
  .card{perspective:900px;height:220px;cursor:pointer;margin-bottom:1.5rem}
  .inner{position:relative;width:100%;height:100%;
         transition:transform .45s cubic-bezier(.4,0,.2,1);transform-style:preserve-3d}
  .card.flipped .inner{transform:rotateY(180deg)}
  .front,.back{position:absolute;inset:0;border-radius:14px;padding:1.5rem;
               backface-visibility:hidden;display:flex;flex-direction:column;justify-content:center}
  .front{background:linear-gradient(135deg,#1a3a6e,#2e5ca8);border:1px solid rgba(79,142,247,.3)}
  .back{background:rgba(255,255,255,.05);border:1px solid rgba(79,142,247,.2);
        transform:rotateY(180deg)}
  .front .q{font-size:1rem;font-weight:600;color:#e8f0ff;text-align:center}
  .front .hint{font-size:.75rem;color:#8aa0c8;margin-top:.75rem;text-align:center}
  .back .a{font-size:.95rem;color:#e8f0ff;line-height:1.6}
  .nav{display:flex;gap:.75rem;justify-content:center;margin-top:.5rem}
  .btn{padding:.55rem 1.4rem;border-radius:9px;border:1px solid rgba(79,142,247,.3);
       background:rgba(79,142,247,.1);color:#4f8ef7;font-size:.88rem;font-weight:600;
       cursor:pointer;font-family:inherit;transition:all .2s}
  .btn:hover{background:rgba(79,142,247,.2)}
  .btn.correct{background:rgba(34,197,94,.12);border-color:rgba(34,197,94,.4);color:#22c55e}
  .btn.wrong{background:rgba(239,68,68,.1);border-color:rgba(239,68,68,.35);color:#ef4444}
  .progress{height:4px;background:rgba(79,142,247,.15);border-radius:99px;margin-bottom:1.5rem;overflow:hidden}
  .progress-bar{height:100%;background:linear-gradient(90deg,#4f8ef7,#a78bfa);
                border-radius:99px;transition:width .4s}
  .done{text-align:center;padding:2rem;color:#22c55e;font-size:1.1rem;font-weight:600}
"""
    html = """<div class="fc-counter" id="ctr"></div>
<div class="progress"><div class="progress-bar" id="pb"></div></div>
<div id="cardWrap"></div>
<div class="nav">
  <button class="btn wrong" onclick="mark(false)">✗ Don't know</button>
  <button class="btn" onclick="flip()">Flip</button>
  <button class="btn correct" onclick="mark(true)">✓ Know it</button>
</div>"""
    script = f"""
const cards = {cards_json};
const esc = s => String(s ?? '').replace(/[&<>"']/g, m => ({{'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}}[m]));
let idx = 0, known = 0;
function render() {{
  if (idx >= cards.length) {{
    document.getElementById('cardWrap').innerHTML = '<div class="done">🎉 Done! ' + known + '/' + cards.length + ' known</div>';
    document.querySelector('.nav').style.display = 'none';
    document.getElementById('ctr').textContent = 'Complete';
    return;
  }}
  const c = cards[idx];
  document.getElementById('ctr').textContent = (idx+1) + ' / ' + cards.length;
  document.getElementById('pb').style.width = (idx/cards.length*100) + '%';
  document.getElementById('cardWrap').innerHTML = `
    <div class="card" id="card" onclick="flip()">
      <div class="inner">
        <div class="front"><div class="q">${{esc(c.q)}}</div><div class="hint">Click to reveal answer</div></div>
        <div class="back"><div class="a">${{esc(c.a)}}</div></div>
      </div>
    </div>`;
}}
function flip() {{ document.getElementById('card').classList.toggle('flipped'); }}
function mark(k) {{ if(k) known++; idx++; render(); }}
render();
"""
    return _page_shell(f"Flash Cards — {title}", css, html, script)


@app.route("/api/view/quiz/<job_id>")
def view_quiz(job_id):
    if not _valid_job(job_id):
        return "<h2 style='font-family:sans-serif;padding:2rem'>Invalid job ID</h2>", 400
    with _jobs_lock:
        job = get_job(job_id)
    if not job:
        return "<h2 style='font-family:sans-serif;padding:2rem'>Guide not found or expired</h2>", 404
    guide = job["guide"]
    title = _he(guide.get("title", "Quiz"))
    mcqs = [m for m in guide.get("mcqs", []) if isinstance(m, dict)]
    if not mcqs:
        return _page_shell(title, "", "<p style='text-align:center;color:#4a5f80;padding:3rem'>No quiz questions available.</p>")

    mcqs_json = json.dumps(mcqs)
    css = """
  .q-num{color:#8aa0c8;font-size:.82rem;margin-bottom:.4rem}
  .q-text{font-size:1rem;font-weight:600;color:#e8f0ff;margin-bottom:1rem;line-height:1.5}
  .opt{display:flex;align-items:center;gap:.6rem;width:100%;text-align:left;
       padding:.6rem .9rem;border-radius:9px;border:1px solid rgba(79,142,247,.2);
       background:rgba(255,255,255,.04);color:#8aa0c8;font-size:.88rem;cursor:pointer;
       font-family:inherit;margin-bottom:.45rem;transition:all .2s}
  .opt:hover:not(:disabled){border-color:#4f8ef7;color:#e8f0ff}
  .opt.correct{border-color:#22c55e;color:#22c55e;background:rgba(34,197,94,.1)}
  .opt.wrong{border-color:#ef4444;color:#ef4444;background:rgba(239,68,68,.08)}
  .explanation{margin-top:.75rem;padding:.75rem;border-radius:9px;
               background:rgba(79,142,247,.08);border:1px solid rgba(79,142,247,.2);
               font-size:.85rem;color:#8aa0c8;display:none}
  .nav{margin-top:1.25rem;display:flex;justify-content:flex-end}
  .next-btn{padding:.55rem 1.4rem;border-radius:9px;border:none;
            background:linear-gradient(135deg,#4f8ef7,#2e5ca8);color:#fff;
            font-size:.88rem;font-weight:600;cursor:pointer;font-family:inherit;display:none}
  .score-box{text-align:center;padding:2.5rem;background:rgba(79,142,247,.08);
             border:1px solid rgba(79,142,247,.2);border-radius:16px}
  .score-big{font-size:3rem;font-weight:700;color:#4f8ef7}
  .progress{height:4px;background:rgba(79,142,247,.15);border-radius:99px;margin-bottom:1.5rem;overflow:hidden}
  .progress-bar{height:100%;background:linear-gradient(90deg,#4f8ef7,#a78bfa);
                border-radius:99px;transition:width .4s}
"""
    html = """<div class="progress"><div class="progress-bar" id="pb"></div></div>
<div id="qWrap"></div>"""
    script = f"""
const qs = {mcqs_json};
const esc = s => String(s ?? '').replace(/[&<>"']/g, m => ({{'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}}[m]));
let idx=0, score=0, answered=false;
function render() {{
  if(idx>=qs.length){{
    document.getElementById('qWrap').innerHTML=`<div class="score-box">
      <div class="score-big">${{score}}/${{qs.length}}</div>
      <div style="color:#8aa0c8;margin-top:.5rem">Quiz complete!</div>
    </div>`;
    document.getElementById('pb').style.width='100%';
    return;
  }}
  const q=qs[idx]; answered=false;
  document.getElementById('pb').style.width=(idx/qs.length*100)+'%';
  const opts=q.options.map((o,i)=>`<button class="opt" id="o${{i}}" data-l="${{esc(o[0])}}" onclick="pick(this.dataset.l,this)">${{esc(o)}}</button>`).join('');
  document.getElementById('qWrap').innerHTML=`
    <div class="q-num">Question ${{idx+1}} of ${{qs.length}}</div>
    <div class="q-text">${{esc(q.q)}}</div>
    ${{opts}}
    <div class="explanation" id="exp">${{esc(q.explanation||'')}}</div>
    <div class="nav"><button class="next-btn" id="nxt" onclick="next()">Next →</button></div>`;
}}
function pick(letter,btn) {{
  if(answered) return; answered=true;
  const q=qs[idx];
  document.querySelectorAll('.opt').forEach(b=>b.disabled=true);
  if(letter===q.answer){{ btn.classList.add('correct'); score++; }}
  else {{ btn.classList.add('wrong');
    document.querySelectorAll('.opt').forEach(b=>{{if(b.textContent.startsWith(q.answer)) b.classList.add('correct');}});
  }}
  document.getElementById('exp').style.display='block';
  document.getElementById('nxt').style.display='inline-block';
}}
function next(){{ idx++; render(); }}
render();
"""
    return _page_shell(f"Quiz — {title}", css, html, script)


@app.route("/api/status")
def status():
    running = ollama_running()
    models  = ollama_models() if running else []
    active  = next((m for m in models if OLLAMA_MODEL in m), models[0] if models else None)
    return jsonify({"ollama": running, "model": active, "models": models})


# ── YouTube transcript endpoint ───────────────────────────────────────────────

def _parse_caption_xml(xml_text):
    """Parse YouTube caption XML and return joined transcript string."""
    import xml.etree.ElementTree as ET
    try:
        root = ET.fromstring(xml_text)
    except Exception:
        return None
    texts = []
    for elem in root.iter("text"):
        t = (elem.text or "").strip()
        if t:
            t = t.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">").replace("&#39;", "'").replace("&quot;", '"')
            texts.append(t)
    return " ".join(texts) if texts else None


def _fetch_captions(video_id):
    """Try multiple methods to get transcript from YouTube without downloading audio."""
    html = None

    # Method 1 — captionTracks from page source
    try:
        r = http.get(f"https://www.youtube.com/watch?v={video_id}", timeout=15,
                     headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"})
        r.raise_for_status()
        html = r.text
    except Exception:
        pass

    if html:
        m = re.search(r'"captionTracks":\s*(\[.*?\])', html)
        if m:
            try:
                tracks = json.loads(m.group(1))
                # Prefer English ASR, then English, then any
                def track_priority(t):
                    lc = t.get("languageCode", "")
                    kind = t.get("kind", "")
                    if lc.startswith("en") and kind == "asr": return 0
                    if lc.startswith("en"): return 1
                    return 2
                tracks.sort(key=track_priority)
                for track in tracks:
                    track_url = track.get("baseUrl")
                    if not track_url:
                        continue
                    try:
                        tr = http.get(track_url, timeout=15)
                        result = _parse_caption_xml(tr.text)
                        if result:
                            return result
                    except Exception:
                        continue
            except Exception:
                pass

    # Method 2 — timedtext API (works for many videos without captionTracks in HTML)
    for kind in ("asr", ""):
        for lang in ("en", "en-US", "en-GB"):
            params = {"v": video_id, "lang": lang, "fmt": "srv3"}
            if kind:
                params["kind"] = kind
            try:
                tr = http.get("https://www.youtube.com/api/timedtext", params=params, timeout=10)
                if tr.status_code == 200 and tr.text.strip():
                    result = _parse_caption_xml(tr.text)
                    if result:
                        return result
            except Exception:
                continue

    raise ValueError("no_captions")


def _transcribe_with_whisper(video_id):
    """Download audio to /tmp, transcribe with Groq Whisper, delete immediately."""
    import glob
    try:
        import yt_dlp
    except ImportError:
        raise ValueError("yt-dlp not installed — cannot transcribe audio.")

    prefix = os.path.join("/tmp", f"yt_{video_id}_{uuid.uuid4().hex[:8]}")
    try:
        ydl_opts = {
            # Prefer smallest audio: opus<96k > m4a < 96k > any audio
            "format": "bestaudio[abr<=96][ext=webm]/bestaudio[abr<=96][ext=m4a]/bestaudio[ext=m4a]/bestaudio[ext=webm]/bestaudio",
            "outtmpl": prefix + ".%(ext)s",
            "quiet": True,
            "no_warnings": True,
            "noplaylist": True,
            "socket_timeout": 30,
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([f"https://www.youtube.com/watch?v={video_id}"])

        files = glob.glob(prefix + ".*")
        if not files:
            raise ValueError("Audio download produced no file.")
        audio_path = files[0]

        size = os.path.getsize(audio_path)
        if size > 20 * 1024 * 1024:
            raise ValueError("Video audio exceeds 20 MB — try a shorter video (under ~15 minutes).")

        ext = os.path.splitext(audio_path)[1].lstrip(".")
        mime = "audio/mp4" if ext in ("m4a", "mp4") else "audio/webm"

        with open(audio_path, "rb") as af:
            r = http.post(
                "https://api.groq.com/openai/v1/audio/transcriptions",
                headers={"Authorization": f"Bearer {GROQ_API_KEY}"},
                files={"file": (f"audio.{ext}", af, mime)},
                data={"model": "whisper-large-v3-turbo", "response_format": "text"},
                timeout=300,
            )
        if r.status_code == 413:
            raise ValueError("Audio file too large for Whisper — try a shorter video (under ~15 minutes).")
        r.raise_for_status()
        return r.text.strip()

    finally:
        for f in glob.glob(prefix + ".*"):
            try:
                os.remove(f)
            except Exception:
                pass


def _fetch_youtube_transcript(video_id):
    """Try captions first; fall back to Whisper audio transcription."""
    try:
        return _fetch_captions(video_id)
    except ValueError as e:
        if str(e) != "no_captions":
            raise
    if not GROQ_API_KEY:
        raise ValueError("No captions found and GROQ_API_KEY not set for Whisper fallback.")
    return _transcribe_with_whisper(video_id)


def _text_to_slides(text, chunk_size=500):
    """Split plain text into slide-like dicts."""
    words = text.split()
    chunks = []
    current = []
    current_len = 0
    for w in words:
        current.append(w)
        current_len += len(w) + 1
        if current_len >= chunk_size:
            chunks.append(" ".join(current))
            current = []
            current_len = 0
    if current:
        chunks.append(" ".join(current))
    return [
        {"slide_num": i + 1, "title": f"Segment {i + 1}", "content": c}
        for i, c in enumerate(chunks)
    ]


def _extract_video_id(url):
    """Extract YouTube video ID from various URL formats."""
    patterns = [
        r'(?:v=|youtu\.be/|/embed/|/shorts/)([a-zA-Z0-9_-]{11})',
    ]
    for p in patterns:
        m = re.search(p, url)
        if m:
            return m.group(1)
    return None


def _stream_text_as_sse(text, language, out_name, job_source):
    """Shared SSE generator for YouTube/text endpoints."""
    def generate():
        try:
            yield _sse({"step": "extract", "msg": "Preparing content…"})
            slides = _text_to_slides(text)
            if not slides:
                yield _sse({"error": "No content extracted"}); return
            yield _sse({"step": "extract", "msg": f"Split into {len(slides)} segments — analysing…"})

            dcfg = DETAIL["standard"]

            yield _sse({"step": "overview", "msg": "Analysing structure and keywords…"})
            overview = pass1_overview(slides, language, dcfg)

            raw_sections = overview.get("sections", [])
            overview["sections"] = [
                s if isinstance(s, dict) else {"title": str(s), "slide_nums": []}
                for s in (raw_sections if isinstance(raw_sections, list) else [])
            ]
            raw_keywords = overview.get("keywords", [])
            overview["keywords"] = [
                k if isinstance(k, dict) else {"term": str(k), "definition": ""}
                for k in (raw_keywords if isinstance(raw_keywords, list) else [])
            ]
            if not overview["sections"]:
                overview["sections"] = [{
                    "title": overview.get("title", "Content Overview"),
                    "slide_nums": [s["slide_num"] for s in slides]
                }]

            sections = overview["sections"]
            for evt in _sections_parallel(sections, slides, language, dcfg):
                yield _sse(evt)

            for evt in _flashcards_mcq_parallel(overview, language, dcfg):
                yield _sse(evt)

            yield _sse({"step": "pdf", "msg": "Building PDF & Markdown…"})
            pdf_buf = build_pdf(overview, language, out_name)
            pdf_bytes = pdf_buf.read()
            md_text = build_markdown(overview)

            job_id = uuid.uuid4().hex
            store_job(job_id, pdf_bytes, md_text, overview, slides, f"{out_name}_study_guide.pdf")

            yield _sse({"step": "done", "job_id": job_id,
                        "sections":   len(sections),
                        "keywords":   len(overview.get("keywords",   [])),
                        "flashcards": len(overview.get("flashcards", [])),
                        "mcqs":       len(overview.get("mcqs",       []))})

        except Exception as e:
            _log.error("STREAM_TEXT_ERROR: %s\n%s", e, _tb.format_exc())
            yield _sse({"error": _safe_err(e)})
    return generate


@app.route("/api/youtube", methods=["POST"])
def youtube_transcript():
    uid, err = _auth_check(request)
    if err:
        return err
    ok, tok_left, reason = _consume_token(uid)
    if not ok:
        return jsonify({
            "error": "You have no tokens left. Upgrade to continue.",
            "code": "no_tokens", "tokens_remaining": 0
        }), 402

    data = request.json or {}
    url = (data.get("url") or "").strip()
    lang_param = data.get("language", "auto")
    if not url:
        return jsonify({"error": "No URL provided"}), 400
    if not ollama_running():
        return jsonify({"error": "AI service is not configured. Set GROQ_API_KEY."}), 503

    video_id = _extract_video_id(url)
    if not video_id:
        return jsonify({"error": "Could not extract video ID from URL"}), 400

    out_name = _safe_name(f"youtube_{video_id}")

    def generate():
        try:
            yield _sse({"step": "transcript", "msg": "Looking for captions…"})
            try:
                transcript_text = _fetch_captions(video_id)
                yield _sse({"step": "transcript", "msg": "Captions found — processing…"})
            except ValueError as e:
                if str(e) != "no_captions":
                    yield _sse({"error": _safe_err(e)}); return
                if not GROQ_API_KEY:
                    yield _sse({"error": "No captions found and GROQ_API_KEY not set."}); return
                yield _sse({"step": "transcript", "msg": "No captions — downloading audio for Whisper transcription…"})
                try:
                    transcript_text = _transcribe_with_whisper(video_id)
                    yield _sse({"step": "transcript", "msg": "Audio transcribed — processing…"})
                except ValueError as we:
                    yield _sse({"error": str(we)}); return

            language = lang_param if lang_param in ("ar", "en") else _detect_language(transcript_text)
            lang_label = "Arabic" if language == "ar" else "English"
            yield _sse({"step": "transcript", "msg": f"Transcript ready ({lang_label}) — building study guide…", "language": language})
            for event in _stream_text_as_sse(transcript_text, language, out_name, "youtube")():
                yield event
        except Exception as ex:
            _log.error("youtube SSE error: %s", ex)
            yield _sse({"error": str(ex)})

    return Response(
        stream_with_context(generate()),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"}
    )


# ── URL / pasted text endpoint ─────────────────────────────────────────────────

def _fetch_url_text(url):
    """Fetch a PUBLIC webpage and extract readable text from p/h/li tags.
    SSRF guard: public http(s) only, no private/internal addresses, no
    redirects, 5 MB response cap."""
    from urllib.parse import urlparse
    import socket, ipaddress
    p = urlparse(url)
    if p.scheme not in ("http", "https") or not p.hostname:
        raise ValueError("Only public http(s) URLs are supported.")
    try:
        infos = socket.getaddrinfo(p.hostname, p.port or (443 if p.scheme == "https" else 80))
    except OSError:
        raise ValueError("Could not resolve URL host.")
    for info in infos:
        addr = ipaddress.ip_address(info[4][0])
        if (addr.is_private or addr.is_loopback or addr.is_link_local or
                addr.is_reserved or addr.is_multicast or addr.is_unspecified):
            raise ValueError("URL points to a private/internal address — not allowed.")
    try:
        r = http.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"},
                     allow_redirects=False, stream=True)
        if 300 <= r.status_code < 400:
            raise ValueError("URL redirects are not supported — paste the final URL.")
        r.raise_for_status()
        content = r.raw.read(5 * 1024 * 1024 + 1, decode_content=True)
        if len(content) > 5 * 1024 * 1024:
            raise ValueError("Page too large (max 5 MB).")
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"Could not fetch URL: {e}")
    html = content.decode(r.encoding or "utf-8", "replace")
    # Remove script/style blocks
    html = re.sub(r'<(script|style)[^>]*>.*?</\1>', ' ', html, flags=re.DOTALL | re.IGNORECASE)
    # Extract text from meaningful tags
    parts = []
    for m in re.finditer(r'<(h[1-6]|p|li)[^>]*>(.*?)</\1>', html, re.DOTALL | re.IGNORECASE):
        inner = re.sub(r'<[^>]+>', ' ', m.group(2))
        inner = inner.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>').replace('&nbsp;', ' ').replace('&#39;', "'").replace('&quot;', '"')
        inner = ' '.join(inner.split())
        if len(inner) > 20:
            parts.append(inner)
    return ' '.join(parts) if parts else ' '.join(re.sub(r'<[^>]+>', ' ', html).split())


@app.route("/api/summarize-text", methods=["POST"])
def summarize_text():
    uid, err = _auth_check(request)
    if err:
        return err
    ok, tok_left, reason = _consume_token(uid)
    if not ok:
        return jsonify({
            "error": "You have no tokens left. Upgrade to continue.",
            "code": "no_tokens", "tokens_remaining": 0
        }), 402

    data = request.json or {}
    text = (data.get("text") or "").strip()
    url  = (data.get("url")  or "").strip()
    lang_param = data.get("language", "auto")
    filename   = _safe_name(data.get("filename") or "pasted_text")

    if not ollama_running():
        return jsonify({"error": "AI service is not configured. Set GROQ_API_KEY."}), 503

    if not text and url:
        try:
            text = _fetch_url_text(url)
        except ValueError as e:
            return jsonify({"error": str(e)}), 400

    if not text:
        return jsonify({"error": "No text or URL provided"}), 400

    language = lang_param if lang_param in ("ar", "en") else _detect_language(text)
    gen = _stream_text_as_sse(text, language, filename, "text")
    return Response(
        stream_with_context(gen()),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"}
    )


# ── Anki export endpoint ───────────────────────────────────────────────────────

@app.route("/api/export/anki/<job_id>")
def export_anki(job_id):
    if not _valid_job(job_id):
        return jsonify({"error": "Invalid job ID"}), 400
    job = get_job(job_id)
    if not job:
        return jsonify({"error": "Job not found or expired"}), 404
    guide = job.get("guide", {})
    flashcards = [f for f in guide.get("flashcards", []) if isinstance(f, dict)]
    if not flashcards:
        return jsonify({"error": "No flashcards available for this job"}), 404

    import csv
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(["front", "back"])
    def _csv_safe(v):
        v = str(v)
        return "'" + v if v[:1] in ("=", "+", "-", "@", "\t") else v
    for fc in flashcards:
        writer.writerow([_csv_safe(fc.get("q", "")), _csv_safe(fc.get("a", ""))])

    csv_bytes = buf.getvalue().encode("utf-8")
    base = job.get("filename", "study_guide").replace(".pdf", "")
    download_name = f"{_safe_name(base)}_anki.csv"
    return send_file(
        io.BytesIO(csv_bytes),
        mimetype="text/csv",
        as_attachment=True,
        download_name=download_name
    )


@app.route("/", defaults={"path": ""})
@app.route("/<path:path>")
def serve(path):
    if path and os.path.exists(os.path.join(DIST, path)):
        return send_from_directory(DIST, path)
    return send_from_directory(DIST, "index.html")


if __name__ == "__main__":
    print(f"  Model : {OLLAMA_MODEL}")
    print(f"  Ollama: {'running' if ollama_running() else 'NOT running — run: ollama serve'}")
    app.run(debug=False, host="127.0.0.1", port=5000, threaded=True)
