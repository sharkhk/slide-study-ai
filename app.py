import io
import re
import sys
import uuid
import json
import os
import time
import secrets
import threading
import traceback as _tb
import logging
import requests as http
from concurrent.futures import ThreadPoolExecutor, as_completed
from collections import OrderedDict

# ── Logger ─────────────────────────────────────────────────────────────────────
# Stream to stdout so the hosting platform (Render) captures app errors/warnings
# in its live log feed. Also keep a best-effort file log for local debugging.
_LOG_FILE = os.path.join(os.path.dirname(__file__), "debug.log")
_log_handlers = [logging.StreamHandler(sys.stdout)]
try:
    _log_handlers.append(logging.FileHandler(_LOG_FILE))
except Exception:
    pass  # read-only FS — stdout handler is enough
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(message)s",
    handlers=_log_handlers,
)
_log = logging.getLogger("app")
from flask import Flask, request, jsonify, send_file, send_from_directory, Response, stream_with_context
from flask_cors import CORS
from pptx import Presentation
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.platypus import (
    BaseDocTemplate, Frame, PageTemplate,
    Paragraph, Spacer, Table, TableStyle,
    HRFlowable, KeepTogether
)
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.pdfbase import pdfmetrics as _pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont as _TTFont

DIST         = os.path.join(os.path.dirname(__file__), "dist")
GROQ_API_KEY = os.environ.get("GROQ_API_KEY", "")
# Primary model — gpt-oss-120b is far more faithful to the source (fewer
# hallucinations / changed names) and stronger in Arabic than the 20b. Needs the
# Groq Developer tier for its rate limits (which this account has).
GROQ_MODEL   = os.environ.get("GROQ_MODEL", "openai/gpt-oss-120b")
# Fallbacks tried automatically if the primary model is unavailable to the key.
GROQ_FALLBACK_MODELS = [
    m.strip() for m in os.environ.get(
        "GROQ_FALLBACK_MODELS", "openai/gpt-oss-20b,llama-3.3-70b-versatile"
    ).split(",") if m.strip()
]
OLLAMA_URL   = os.environ.get("OLLAMA_URL", "http://localhost:11434")
OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "mistral")

# ── Supabase / Auth ────────────────────────────────────────────────────────────
SUPABASE_URL              = os.environ.get("SUPABASE_URL", "")
SUPABASE_ANON_KEY         = os.environ.get("SUPABASE_ANON_KEY", "")
SUPABASE_SERVICE_ROLE_KEY = os.environ.get("SUPABASE_SERVICE_ROLE_KEY", "")
SUPABASE_JWT_SECRET       = os.environ.get("SUPABASE_JWT_SECRET", "")

# ── Stripe ─────────────────────────────────────────────────────────────────────
STRIPE_SECRET_KEY    = os.environ.get("STRIPE_SECRET_KEY", "")
STRIPE_PUBLISHABLE_KEY = os.environ.get("STRIPE_PUBLISHABLE_KEY", "")
STRIPE_WEBHOOK_SECRET = os.environ.get("STRIPE_WEBHOOK_SECRET", "")
STRIPE_PRICE_ID      = os.environ.get("STRIPE_PRICE_ID", "")
# Monthly Pro price in USD — used only to estimate MRR on the admin dashboard.
# Override with the PRO_MONTHLY_USD env var if the price changes.
try:
    PRO_MONTHLY_USD = float(os.environ.get("PRO_MONTHLY_USD", "2.99"))
except (TypeError, ValueError):
    PRO_MONTHLY_USD = 2.99
APP_URL              = os.environ.get("APP_URL", "https://slide-study-ai.onrender.com")
# Shared secret Cloudflare injects (via a Transform Rule adding header
# X-Origin-Verify) so the origin can tell real Cloudflare traffic from requests
# sent straight to the Render origin. When set, CF-Connecting-IP is only trusted
# on verified requests. Leave unset to keep the old (trust-CF) behaviour.
CF_ORIGIN_SECRET     = os.environ.get("CF_ORIGIN_SECRET", "")

_AUTH_ENABLED = bool(SUPABASE_URL)  # verify via JWKS (asymmetric) or HS256 shared secret

# FAIL CLOSED: with auth disabled every caller becomes a "dev" user with 999
# tokens and an active subscription. That is fine for local dev, but if it ever
# happened in production (SUPABASE_URL unset on a deploy) the whole app — paid
# features included — would be wide open. On Render, refuse to start instead.
_IS_PRODUCTION = bool(os.environ.get("RENDER") or os.environ.get("PRODUCTION"))
if _IS_PRODUCTION and not _AUTH_ENABLED:
    raise RuntimeError(
        "SUPABASE_URL is not set but this is a production environment — refusing "
        "to start in open (no-auth) dev mode. Set SUPABASE_URL and the Supabase keys."
    )

DETAIL = {
    "brief":    {"slide_chars": 400,  "max_slides": 30,  "keywords": "8-10",  "bullets": "2-4",  "n_flash": 6,  "n_mcq": 5,  "num_predict": 1200},
    "standard": {"slide_chars": 700,  "max_slides": 60,  "keywords": "18-25", "bullets": "3-8",  "n_flash": 14, "n_mcq": 10, "num_predict": 2200},
    "detailed": {"slide_chars": 1200, "max_slides": 120, "keywords": "25-35", "bullets": "5-12", "n_flash": 20, "n_mcq": 15, "num_predict": 3200},
}

# ── Arabic PDF support ─────────────────────────────────────────────────────────
_ARABIC_FONT      = "NotoNaskhArabic"
_ARABIC_FONT_PATH = "/tmp/NotoNaskhArabic.ttf"
_arabic_font_ok   = False
_arabic_font_lock = threading.Lock()

def _ensure_arabic_font():
    global _arabic_font_ok
    with _arabic_font_lock:
        if _arabic_font_ok:
            return True
        try:
            if not os.path.exists(_ARABIC_FONT_PATH):
                r = http.get(
                    "https://github.com/googlefonts/noto-fonts/raw/main/hinted/ttf/NotoNaskhArabic/NotoNaskhArabic-Regular.ttf",
                    timeout=30, allow_redirects=True
                )
                r.raise_for_status()
                with open(_ARABIC_FONT_PATH, "wb") as fh:
                    fh.write(r.content)
            _pdfmetrics.registerFont(_TTFont(_ARABIC_FONT, _ARABIC_FONT_PATH))
            _arabic_font_ok = True
            return True
        except Exception as exc:
            print(f"Arabic font error: {exc}", flush=True)
            return False

def _ar(text):
    """Reshape + bidi-flip Arabic for correct visual display in LTR PDF renderer."""
    if not text:
        return text
    try:
        import arabic_reshaper
        from bidi.algorithm import get_display
        return get_display(arabic_reshaper.reshape(str(text)))
    except ImportError:
        return str(text)

# ── Supabase client (lazy, uses service-role key → bypasses RLS) ───────────────
_sb_client = None
_sb_lock   = threading.Lock()

def _get_sb():
    global _sb_client
    if _sb_client:
        return _sb_client
    with _sb_lock:
        if _sb_client:
            return _sb_client
        if SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY:
            try:
                from supabase import create_client
                _sb_client = create_client(SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY)
            except Exception as e:
                print(f"Supabase init error: {e}", flush=True)
    return _sb_client

# ── Stripe setup ───────────────────────────────────────────────────────────────
if STRIPE_SECRET_KEY:
    try:
        import stripe as _stripe
        _stripe.api_key = STRIPE_SECRET_KEY
    except ImportError:
        _stripe = None
else:
    _stripe = None

# ── JWT helpers ────────────────────────────────────────────────────────────────
def _get_bearer(req):
    auth = req.headers.get("Authorization", "")
    return auth[7:] if auth.startswith("Bearer ") else None

_jwks_client = None
def _get_jwks_client():
    global _jwks_client
    if _jwks_client is None and SUPABASE_URL:
        import jwt as _pyjwt
        _jwks_client = _pyjwt.PyJWKClient(SUPABASE_URL.rstrip("/") + "/auth/v1/.well-known/jwks.json")
    return _jwks_client

def _jwt_payload(token):
    """Verify a Supabase-issued JWT and return the full decoded payload dict,
    or None. Asymmetric (ES/RS/PS/Ed) tokens verify against the project JWKS;
    HS256 against the shared secret."""
    if not token:
        return None
    try:
        import jwt as _pyjwt
        alg = _pyjwt.get_unverified_header(token).get("alg", "")
        if alg.startswith(("ES", "RS", "PS", "Ed")):
            client = _get_jwks_client()
            if client is None:
                return None
            key = client.get_signing_key_from_jwt(token).key
            return _pyjwt.decode(token, key, algorithms=[alg], audience="authenticated")
        elif SUPABASE_JWT_SECRET:
            return _pyjwt.decode(token, SUPABASE_JWT_SECRET, algorithms=["HS256"], audience="authenticated")
        return None
    except Exception as _e:
        _log.warning("JWT verify failed: %s", type(_e).__name__)
        return None

def _verify_jwt(token):
    """Return the verified user_id (str) or None."""
    payload = _jwt_payload(token)
    return payload.get("sub") if payload else None

def _identity_from_jwt(token):
    """Pull email / display name / avatar from a Google (or other) OAuth JWT.
    Supabase puts these under user_metadata (full_name/name, avatar_url/picture)."""
    p = _jwt_payload(token) or {}
    meta = p.get("user_metadata") or {}
    return {
        "email":  p.get("email") or meta.get("email") or "",
        "name":   meta.get("full_name") or meta.get("name") or "",
        "avatar": meta.get("avatar_url") or meta.get("picture") or "",
    }

# ── Token helpers ──────────────────────────────────────────────────────────────
def _consume_token(user_id):
    """
    Atomically consume one token via Supabase RPC.
    Returns (ok: bool, tokens_remaining: int, reason: str)
    """
    sb = _get_sb()
    if not sb:
        return True, 999, ""   # Supabase not configured → dev mode, always allow
    try:
        r = sb.rpc("consume_token", {"p_user_id": user_id}).execute()
        d = r.data if isinstance(r.data, dict) else {}
        ok = d.get("success", False)
        return ok, d.get("tokens_remaining", 0), d.get("reason", "")
    except Exception as exc:
        print(f"consume_token error: {exc}", flush=True)
        return False, 0, "db_error"

# Flips to False if the add_tokens RPC isn't installed yet (see migrations/).
_add_tokens_rpc = True

def _add_tokens(sb, user_id, delta):
    """Atomically add (or subtract) `delta` tokens and return the new balance.
    Uses the add_tokens SQL RPC when available — a plain read-modify-write here
    races with the atomic consume_token RPC and loses concurrent decrements.
    Falls back to read-modify-write if the RPC is not installed."""
    global _add_tokens_rpc
    if not sb or not user_id:
        return None
    if _add_tokens_rpc:
        try:
            r = sb.rpc("add_tokens", {"p_user_id": user_id, "p_delta": delta}).execute()
            if isinstance(r.data, int):
                return r.data
            if isinstance(r.data, list) and r.data:
                return r.data[0]
            return r.data
        except Exception as exc:
            _add_tokens_rpc = False
            _log.warning("add_tokens RPC unavailable — using read-modify-write: %s", exc)
    try:
        row = sb.table("users").select("tokens_remaining").eq("id", user_id).single().execute()
        cur = (row.data or {}).get("tokens_remaining")
        if cur is None:
            return None
        new_bal = cur + delta
        sb.table("users").update({"tokens_remaining": new_bal}).eq("id", user_id).execute()
        return new_bal
    except Exception as exc:
        _log.error("_add_tokens fallback error: %s", exc)
        return None

def _refund_token(user_id):
    """Best-effort: give a token back when a job fails after consuming it,
    so users are only charged for a study guide they actually receive."""
    sb = _get_sb()
    if not sb or not user_id:
        return
    new_bal = _add_tokens(sb, user_id, 1)
    if new_bal is not None:
        _log.info("Refunded 1 token to %s (-> %s)", user_id, new_bal)

def _get_user(user_id):
    """Fetch full user row from Supabase."""
    sb = _get_sb()
    if not sb:
        return None
    try:
        r = sb.table("users").select("*").eq("id", user_id).single().execute()
        return r.data
    except Exception:
        return None

def _get_or_create_referral_code(user_id):
    """Return this user's referral code, generating one if not yet set."""
    import hashlib
    sb = _get_sb()
    if not sb:
        return None
    try:
        r = sb.table("users").select("referral_code").eq("id", user_id).single().execute()
        code = (r.data or {}).get("referral_code")
        if code:
            return code
        code = hashlib.sha256(user_id.encode()).hexdigest()[:8].upper()
        sb.table("users").update({"referral_code": code}).eq("id", user_id).execute()
        return code
    except Exception:
        return None

def _award_referral(new_subscriber_id, sb):
    """Award 10 tokens to the referrer when their referee first subscribes."""
    try:
        r = sb.table("users").select("referred_by, referral_paid").eq("id", new_subscriber_id).single().execute()
        if not r.data:
            return
        referrer_id  = r.data.get("referred_by")
        already_paid = r.data.get("referral_paid", False)
        if not referrer_id or already_paid:
            return
        # Mark paid FIRST so a redelivered/duplicate event can't double-award,
        # then grant atomically.
        sb.table("users").update({"referral_paid": True}).eq("id", new_subscriber_id).execute()
        _add_tokens(sb, referrer_id, 10)
        _log.info(f"Referral reward: 10 tokens → {referrer_id} (subscriber={new_subscriber_id})")
    except Exception as exc:
        _log.error(f"_award_referral error: {exc}")

def _auth_check(req):
    """
    Extract + verify JWT from request.
    Returns (user_id, error_response_tuple | None).
    When _AUTH_ENABLED is False (local dev), always returns ('dev', None).
    """
    if not _AUTH_ENABLED:
        return "dev", None
    tok = _get_bearer(req)
    uid = _verify_jwt(tok)
    if not uid:
        return None, (jsonify({"error": "Sign in required", "code": "auth_required"}), 401)
    return uid, None

def _auth_optional(req):
    """Return a verified user id, or None for anonymous callers (no / invalid token).
    Used by endpoints that also allow signed-out users a small free quota."""
    if not _AUTH_ENABLED:
        return "dev"
    tok = _get_bearer(req)
    return _verify_jwt(tok) if tok else None

# ── Anonymous (no-login) free credits ──────────────────────────────────────────
# Let visitors try the product a few times without an account. Tracked per-IP in
# memory over a rolling window; tune with ANON_FREE_LIMIT (0 disables anon use).
ANON_FREE_LIMIT = int(os.environ.get("ANON_FREE_LIMIT", "2"))
_ANON_WINDOW    = int(os.environ.get("ANON_WINDOW_SEC", str(24 * 3600)))
_anon_lock      = threading.Lock()
_anon_usage     = {}  # ip -> [timestamps]

def _anon_remaining(ip):
    now = time.time()
    with _anon_lock:
        times = [t for t in _anon_usage.get(ip, []) if now - t < _ANON_WINDOW]
        _anon_usage[ip] = times
        return max(0, ANON_FREE_LIMIT - len(times))

def _anon_consume(ip):
    """Consume one anonymous free credit. Returns (ok, remaining)."""
    now = time.time()
    with _anon_lock:
        times = [t for t in _anon_usage.get(ip, []) if now - t < _ANON_WINDOW]
        if len(times) >= ANON_FREE_LIMIT:
            _anon_usage[ip] = times
            return False, 0
        times.append(now)
        _anon_usage[ip] = times
        return True, max(0, ANON_FREE_LIMIT - len(times))

def _anon_refund(ip):
    with _anon_lock:
        times = _anon_usage.get(ip)
        if times:
            times.pop()
            _anon_usage[ip] = times

def _refund_credit(uid, ip):
    """Refund one credit to whoever was charged — signed-in user or anon IP."""
    if uid:
        _refund_token(uid)
    elif ip:
        _anon_refund(ip)

app = Flask(__name__, static_folder=DIST, static_url_path="")
app.config['SECRET_KEY'] = secrets.token_hex(32)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50 MB upload limit
CORS(app, origins=[APP_URL, "http://localhost:5173", "http://127.0.0.1:5173"])

# Security headers on every response
@app.after_request
def _security_headers(resp):
    resp.headers.setdefault("X-Content-Type-Options", "nosniff")
    resp.headers.setdefault("X-Frame-Options", "DENY")
    resp.headers.setdefault("Referrer-Policy", "strict-origin-when-cross-origin")
    resp.headers.setdefault("Permissions-Policy", "camera=(), microphone=(), geolocation=()")
    if request.is_secure or request.headers.get("X-Forwarded-Proto", "").lower() == "https":
        resp.headers.setdefault("Strict-Transport-Security", "max-age=31536000; includeSubDomains")
    if request.path.startswith("/api/"):
        resp.headers.setdefault("Cache-Control", "no-store")
    # Content-Security-Policy — defence-in-depth backstop behind the output
    # escaping. The React SPA loads an external bundle (no inline JS), so it gets
    # a strict script-src; the server-rendered pages (/admin, /api/view/*) use
    # inline <script>/onclick, so only those routes relax script-src.
    _p = request.path
    _inline_html = _p == "/admin" or _p.startswith("/admin/") or _p.startswith("/api/view/")
    _script_src = "script-src 'self' 'unsafe-inline'" if _inline_html else "script-src 'self'"
    resp.headers.setdefault("Content-Security-Policy", (
        "default-src 'self'; "
        f"{_script_src}; "
        "style-src 'self' 'unsafe-inline' https://fonts.googleapis.com; "
        "font-src 'self' https://fonts.gstatic.com; "
        "img-src 'self' data: https:; "
        "connect-src 'self' https://*.supabase.co; "
        "frame-ancestors 'none'; base-uri 'none'; object-src 'none'; form-action 'self'"
    ))
    return resp

# Per-IP rate limiting
_rate_limit_lock = threading.Lock()
_rate_limit      = {}
_RATE_WINDOW     = 60
_RATE_MAX        = 5

def _check_rate_limit(ip, scope="main", limit=_RATE_MAX):
    # Exempt ONLY genuine loopback (local dev). A blanket _is_private() exemption
    # let an attacker send CF-Connecting-IP: 10.0.0.1 to disable rate limiting
    # entirely — private-but-not-loopback client IPs must still be limited.
    if ip in ("127.0.0.1", "::1"):
        return True
    key = f"{scope}:{ip}"
    now = time.time()
    with _rate_limit_lock:
        times = [t for t in _rate_limit.get(key, []) if now - t < _RATE_WINDOW]
        if len(times) >= limit:
            _rate_limit[key] = times
            return False
        times.append(now)
        _rate_limit[key] = times
    return True

# ── Visitor tracking ──────────────────────────────────────────────────────────
ADMIN_TOKEN  = os.environ.get("ADMIN_TOKEN") or secrets.token_urlsafe(24)
_visitors    = []
_vis_lock    = threading.Lock()
_blocked_ips = set()   # IPs that are blocked from using the app

# Geo-enrichment cache + concurrency guard: without this, every non-asset
# request spawns a thread doing a 5s HTTP call to ip-api.com (which caps at
# 45/min), so a burst would exhaust threads and hammer the API.
_geo_cache        = {}
_geo_cache_order  = []
_geo_inflight     = set()
_geo_lock         = threading.Lock()
_GEO_MAX_INFLIGHT = 6

_PRIVATE_RANGES = (
    "127.", "::1", "10.", "192.168.", "172.16.", "172.17.", "172.18.",
    "172.19.", "172.20.", "172.21.", "172.22.", "172.23.", "172.24.",
    "172.25.", "172.26.", "172.27.", "172.28.", "172.29.", "172.30.",
    "172.31.", "0.0.0.0",
)

def _is_private(ip):
    return any(ip.startswith(p) for p in _PRIVATE_RANGES)

def _admin_ok():
    supplied = request.headers.get("X-Admin-Token") or request.args.get("token") or ""
    return secrets.compare_digest(supplied, ADMIN_TOKEN)

def _safe_err(e):
    if isinstance(e, ValueError):
        return str(e)
    _log.error("internal error: %s\n%s", e, _tb.format_exc())
    return "Processing failed — please try again."

def _via_cloudflare():
    """True if we can trust this request's CF-Connecting-IP header. When
    CF_ORIGIN_SECRET is configured, a Cloudflare Transform Rule stamps it as
    X-Origin-Verify; a request that reaches the Render origin directly (bypassing
    Cloudflare) won't have it, so its CF-Connecting-IP must NOT be trusted.
    Without the secret configured we can't tell, so we trust it as before."""
    if not CF_ORIGIN_SECRET:
        return True
    return secrets.compare_digest(request.headers.get("X-Origin-Verify", ""), CF_ORIGIN_SECRET)

def _client_ip():
    ra = request.remote_addr or "unknown"
    # CF-Connecting-IP is set by Cloudflare's edge, but ONLY trustworthy for
    # traffic that actually transited Cloudflare (see _via_cloudflare) — the
    # Render origin is reachable directly, where the header is attacker-controlled.
    cf = (request.headers.get("CF-Connecting-IP") or "").strip()
    if cf and _via_cloudflare():
        return cf
    if not _is_private(ra):
        return ra
    # No trusted CF header and remote_addr is a private proxy hop. Fall back to
    # the LAST X-Forwarded-For entry (appended by the nearest trusted proxy) —
    # NOT the first, which is client-supplied and trivially spoofable.
    parts = [p.strip() for p in request.headers.get("X-Forwarded-For", "").split(",") if p.strip()]
    return parts[-1] if parts else ra

def _cache_geo(ip, geo):
    with _geo_lock:
        _geo_cache[ip] = geo
        _geo_cache_order.append(ip)
        if len(_geo_cache_order) > 2000:
            _geo_cache.pop(_geo_cache_order.pop(0), None)
        _geo_inflight.discard(ip)

def _enrich_geo(entry, ip):
    """Background thread: full geolocation via ip-api.com. Result is cached so
    repeat visitors don't re-spawn a thread / re-hit the API."""
    if _is_private(ip):
        geo = {"country": "Local / LAN", "region": "", "city": "localhost",
               "isp": "private network", "lat": "", "lon": ""}
        with _vis_lock:
            entry.update(geo)
        _cache_geo(ip, geo)
        return
    try:
        r = http.get(
            f"http://ip-api.com/json/{ip}"
            f"?fields=status,country,countryCode,regionName,city,isp,org,lat,lon",
            timeout=5
        )
        d = r.json()
        if d.get("status") != "success":
            with _geo_lock:
                _geo_inflight.discard(ip)
            return
        geo = {
            "country": d.get("country", ""),
            "region":  d.get("regionName", ""),
            "city":    d.get("city", ""),
            "isp":     d.get("org") or d.get("isp", ""),
            "lat":     d.get("lat", ""),
            "lon":     d.get("lon", ""),
        }
        with _vis_lock:
            entry.update(geo)
        _cache_geo(ip, geo)
    except Exception:
        with _geo_lock:
            _geo_inflight.discard(ip)

@app.before_request
def track_visitor():
    skip = ("/assets/", "/favicon", "/admin")
    if any(request.path.startswith(s) for s in skip):
        # still enforce block on non-admin paths
        pass
    ip = _client_ip()

    # Block check — return 403 immediately for blocked IPs (except admin itself)
    if ip in _blocked_ips and not request.path.startswith("/admin"):
        from flask import abort
        abort(403)

    if any(request.path.startswith(s) for s in ("/assets/", "/favicon")):
        return

    entry = {
        "time":    time.strftime("%Y-%m-%d %H:%M:%S UTC", time.gmtime()),
        "ip":      ip,
        "country": request.headers.get("CF-IPCountry", ""),
        "region":  "",
        "city":    "",
        "isp":     "",
        "lat":     "",
        "lon":     "",
        "path":    request.path,
        "method":  request.method,
        "ua":      (request.headers.get("User-Agent") or "")[:160],
    }
    with _vis_lock:
        _visitors.insert(0, entry)
        if len(_visitors) > 1000:
            _visitors.pop()
    # Enrich geo (CF headers don't give city/ISP/coords), but reuse the cache,
    # dedupe in-flight lookups per IP, and cap concurrent threads so a burst
    # can't exhaust threads or exceed ip-api.com's rate limit.
    do_enrich = False
    with _geo_lock:
        cached = _geo_cache.get(ip)
        if cached is not None:
            entry.update(cached)
        elif ip not in _geo_inflight and len(_geo_inflight) < _GEO_MAX_INFLIGHT:
            _geo_inflight.add(ip)
            do_enrich = True
    if do_enrich:
        threading.Thread(target=_enrich_geo, args=(entry, ip), daemon=True).start()


@app.route("/admin/block", methods=["POST"])
def admin_block():
    if not _admin_ok():
        return jsonify({"error": "Unauthorized"}), 401
    ip = request.json.get("ip", "").strip()
    if not ip:
        return jsonify({"error": "No IP"}), 400
    with _vis_lock:
        _blocked_ips.add(ip)
    return jsonify({"ok": True, "blocked": ip})


@app.route("/admin/unblock", methods=["POST"])
def admin_unblock():
    if not _admin_ok():
        return jsonify({"error": "Unauthorized"}), 401
    ip = request.json.get("ip", "").strip()
    with _vis_lock:
        _blocked_ips.discard(ip)
    return jsonify({"ok": True, "unblocked": ip})


@app.route("/admin/clear", methods=["POST"])
def admin_clear_log():
    if not _admin_ok():
        return jsonify({"error": "Unauthorized"}), 401
    with _vis_lock:
        _visitors.clear()
    return jsonify({"ok": True})


# Effective monthly price (USD) of the Pro plan, read from the live Stripe price
# and cached for an hour so the admin dashboard doesn't hit Stripe on every load.
# Any recurring interval (year/month/week/day + interval_count) is normalised to
# a monthly figure. Falls back to PRO_MONTHLY_USD when Stripe is unavailable.
_price_cache = {"amount": None, "ts": 0.0}
_price_lock  = threading.Lock()

def _monthly_price_usd():
    now = time.time()
    with _price_lock:
        if _price_cache["amount"] is not None and now - _price_cache["ts"] < 3600:
            return _price_cache["amount"]
    amount = PRO_MONTHLY_USD  # fallback
    if STRIPE_PRICE_ID and STRIPE_SECRET_KEY:
        try:
            import stripe as _stripe
            _stripe.api_key = STRIPE_SECRET_KEY
            p = _stripe.Price.retrieve(STRIPE_PRICE_ID).to_dict()
            cents = p.get("unit_amount")
            if cents is not None:
                val = cents / 100.0
                rec = p.get("recurring") or {}
                interval = rec.get("interval")
                count = rec.get("interval_count", 1) or 1
                if interval == "year":
                    val = val / (12.0 * count)
                elif interval == "week":
                    val = val * (52.0 / 12.0) / count
                elif interval == "day":
                    val = val * (365.0 / 12.0) / count
                elif interval == "month":
                    val = val / count
                amount = round(val, 2)
        except Exception as exc:
            _log.warning("Stripe price fetch failed, using fallback $%.2f: %s", amount, exc)
    with _price_lock:
        _price_cache["amount"] = amount
        _price_cache["ts"] = now
    return amount


# Matches a canonical UUID (the Supabase user id). Used to validate admin
# action targets so a caller can't inject arbitrary values.
_UUID_RE = re.compile(r'^[0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{12}$')


@app.route("/admin/user/grant", methods=["POST"])
def admin_user_grant():
    """Grant (or deduct) tokens for one user. Admin-only, header-authenticated."""
    if not _admin_ok():
        return jsonify({"error": "Unauthorized"}), 401
    data = request.get_json(silent=True) or {}
    uid = str(data.get("user_id", "")).strip()
    try:
        amount = int(data.get("amount", 0))
    except (TypeError, ValueError):
        return jsonify({"error": "amount must be an integer"}), 400
    if not _UUID_RE.match(uid):
        return jsonify({"error": "invalid user_id"}), 400
    if amount == 0 or amount < -1000 or amount > 1000:
        return jsonify({"error": "amount must be a non-zero integer in [-1000, 1000]"}), 400
    sb = _get_sb()
    if sb is None:
        return jsonify({"error": "Supabase not configured"}), 500
    new_bal = _add_tokens(sb, uid, amount)
    if new_bal is None:
        return jsonify({"error": "user not found or update failed"}), 404
    _log.info("ADMIN grant %+d tokens -> user %s (new balance %s)", amount, uid, new_bal)
    return jsonify({"ok": True, "tokens_remaining": new_bal})


@app.route("/admin/user/cancel", methods=["POST"])
def admin_user_cancel():
    """Cancel a user's subscription. If they have a Stripe subscription it is set
    to cancel at period end (they keep access until then; the webhook finalises
    the status). Otherwise the local status is downgraded to free. Admin-only."""
    if not _admin_ok():
        return jsonify({"error": "Unauthorized"}), 401
    data = request.get_json(silent=True) or {}
    uid = str(data.get("user_id", "")).strip()
    if not _UUID_RE.match(uid):
        return jsonify({"error": "invalid user_id"}), 400
    sb = _get_sb()
    if sb is None:
        return jsonify({"error": "Supabase not configured"}), 500
    try:
        row = sb.table("users").select("subscription_id").eq("id", uid).single().execute()
    except Exception as exc:
        return jsonify({"error": f"lookup failed: {exc}"}), 500
    sub_id = (row.data or {}).get("subscription_id") or ""
    if sub_id and STRIPE_SECRET_KEY:
        try:
            import stripe as _stripe
            _stripe.api_key = STRIPE_SECRET_KEY
            _stripe.Subscription.modify(sub_id, cancel_at_period_end=True)
        except Exception as exc:
            _log.error("ADMIN cancel: Stripe error for %s: %s", uid, exc)
            return jsonify({"error": f"Stripe cancel failed: {exc}"}), 502
        _log.info("ADMIN cancel: Stripe sub %s set to cancel at period end (user %s)", sub_id, uid)
        return jsonify({"ok": True, "canceled_at_period_end": True})
    # No Stripe subscription on file — downgrade locally.
    try:
        sb.table("users").update({"subscription_status": "free",
                                  "subscription_period_end": None}).eq("id", uid).execute()
    except Exception as exc:
        return jsonify({"error": f"update failed: {exc}"}), 500
    _log.info("ADMIN cancel: user %s downgraded to free (no Stripe sub)", uid)
    return jsonify({"ok": True, "canceled_at_period_end": False})


# Token gate for /admin. Lets you enter the admin token once; the browser then
# remembers it (localStorage) and auto-opens the dashboard on later visits, so
# you never retype it. The token value is only ever entered by you.
ADMIN_GATE_HTML = """<!DOCTYPE html><html><head><meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>Admin — Alimne</title>
<style>
 *{box-sizing:border-box;margin:0;padding:0}
 body{font-family:'Segoe UI',system-ui,sans-serif;background:#050d1a;color:#e8f0ff;
      min-height:100vh;display:flex;align-items:center;justify-content:center;padding:1rem}
 .card{background:#0a1628;border:1px solid #1a3a6e;border-radius:14px;padding:1.8rem;max-width:400px;width:100%}
 h2{font-size:1.1rem;margin:0 0 .4rem;display:flex;gap:.5rem;align-items:center}
 p{color:#8aa0c8;font-size:.88rem;margin:0 0 1.1rem;line-height:1.5}
 input[type=password]{width:100%;padding:.7rem;border-radius:8px;border:1px solid #1a3a6e;
      background:#050d1a;color:#e8f0ff;font-size:.9rem;margin-bottom:.75rem}
 button{width:100%;padding:.72rem;border-radius:8px;border:none;background:#4f8ef7;color:#fff;
      font-weight:600;font-size:.9rem;cursor:pointer}
 button:hover{opacity:.9}
 .err{color:#f87171;font-size:.85rem;margin-bottom:.75rem;display:none}
 label{display:flex;gap:.5rem;align-items:center;color:#8aa0c8;font-size:.82rem;margin:0 0 1rem;cursor:pointer}
</style></head><body>
<div class="card">
  <h2>🔒 Alimne Admin</h2>
  <p>Enter your admin token once. This browser will remember it, so you won't be asked again.</p>
  <div class="err" id="err">That token was rejected — check it and try again.</div>
  <input id="tok" type="password" placeholder="Admin token" autocomplete="off" autofocus>
  <label><input type="checkbox" id="remember" checked> Remember on this device</label>
  <button id="go">Open dashboard</button>
</div>
<script>
 var KEY = "alimne_admin_token";
 var tried = new URLSearchParams(location.search).get("token");
 if (tried) {
   // A token was supplied but we still landed on this gate => it was invalid.
   try { localStorage.removeItem(KEY); } catch (e) {}
   document.getElementById("err").style.display = "block";
 } else {
   var saved = null; try { saved = localStorage.getItem(KEY); } catch (e) {}
   if (saved) { location.replace("/admin?token=" + encodeURIComponent(saved)); }
 }
 function go() {
   var v = document.getElementById("tok").value.trim();
   if (!v) return;
   try {
     if (document.getElementById("remember").checked) localStorage.setItem(KEY, v);
     else localStorage.removeItem(KEY);
   } catch (e) {}
   location.href = "/admin?token=" + encodeURIComponent(v);
 }
 document.getElementById("go").addEventListener("click", go);
 document.getElementById("tok").addEventListener("keydown", function (e) { if (e.key === "Enter") go(); });
</script>
</body></html>"""


@app.route("/admin")
def admin_page():
    if not _admin_ok():
        return ADMIN_GATE_HTML, 401
    token = ADMIN_TOKEN
    with _vis_lock:
        vis_copy     = list(_visitors)
        blocked_copy = set(_blocked_ips)

    def flag(cc):
        # Convert 2-letter country code to emoji flag
        if not cc or len(cc) != 2:
            return ""
        return "".join(chr(0x1F1E6 + ord(c) - ord('A')) for c in cc.upper())

    # JS-string-safe escaper for values dropped into an onclick='...(\'x\')'
    def _js(s):
        return _he(s).replace("\\", "\\\\").replace("'", "\\'")

    rows = ""
    for v in vis_copy:
        blocked = v["ip"] in blocked_copy
        # All visitor fields below are attacker-controlled (UA/path/headers, or
        # geo derived from a spoofable IP) — escape every one to prevent stored XSS.
        ip_h    = _he(v.get("ip", ""))
        ip_js   = _js(v.get("ip", ""))
        loc_parts = [p for p in [v.get("city",""), v.get("region",""), v.get("country","")] if p]
        location  = _he(", ".join(loc_parts)) if loc_parts else "—"
        map_link  = ""
        if v.get("lat") and v.get("lon"):
            map_link = f'<a href="https://maps.google.com/?q={_he(v["lat"])},{_he(v["lon"])}" target="_blank" rel="noopener noreferrer" style="color:#4f8ef7;font-size:11px">📍 map</a>'
        block_btn = (
            f'<button onclick="unblock(\'{ip_js}\')" '
            f'style="background:#16a34a;color:#fff;border:none;padding:3px 10px;border-radius:6px;cursor:pointer;font-size:12px">✓ Unblock</button>'
            if blocked else
            f'<button onclick="blockIp(\'{ip_js}\')" '
            f'style="background:#dc2626;color:#fff;border:none;padding:3px 10px;border-radius:6px;cursor:pointer;font-size:12px">⛔ Block</button>'
        )
        row_style = "background:#1a0a0a" if blocked else ""
        rows += f"""
          <tr style="{row_style}">
            <td style="color:#8aa0c8;white-space:nowrap">{_he(v.get('time',''))}</td>
            <td><b style="{'color:#f87171' if blocked else ''}">{ip_h}</b>
                {'<span style="background:#7f1d1d;color:#fca5a5;padding:1px 6px;border-radius:4px;font-size:10px;margin-left:4px">BLOCKED</span>' if blocked else ''}
            </td>
            <td>{_he(v.get('country','—'))}</td>
            <td>{location} {map_link}</td>
            <td style="max-width:180px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;color:#8aa0c8">{_he(v.get('isp','—'))}</td>
            <td style="color:#6b7fa8">{_he(v.get('path',''))}</td>
            <td style="max-width:200px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;font-size:11px;color:#4a5f80">{_he(v.get('ua',''))}</td>
            <td>{block_btn}</td>
          </tr>"""

    blocked_section = ""
    if blocked_copy:
        blocked_rows = "".join(
            f'<tr><td style="color:#f87171;padding:6px 12px">{_he(ip)}</td>'
            f'<td><button onclick="unblock(\'{_js(ip)}\')" style="background:#16a34a;color:#fff;border:none;padding:2px 10px;border-radius:6px;cursor:pointer;font-size:12px">Unblock</button></td></tr>'
            for ip in sorted(blocked_copy)
        )
        blocked_section = f"""
        <div style="margin:1.5rem 2rem;background:#1a0a0a;border:1px solid #7f1d1d;border-radius:10px;padding:1rem">
          <h3 style="margin:0 0 0.75rem;color:#f87171;font-size:0.95rem">⛔ Blocked IPs ({len(blocked_copy)})</h3>
          <table style="border-collapse:collapse;font-size:13px"><tbody>{blocked_rows}</tbody></table>
        </div>"""

    # ── Subscribers (from the Supabase `users` table) ────────────────────────────
    subs_rows = ""
    subs_total = 0
    subs_active = 0
    new_this_week = 0
    subs_error = ""
    try:
        _sb = _get_sb()
        if _sb is None:
            subs_error = "Supabase is not configured on this server (dev mode)."
        else:
            _res = _sb.table("users").select("*").limit(2000).execute()
            _users = _res.data or []
            _users.sort(key=lambda u: str(u.get("created_at") or ""), reverse=True)
            subs_total = len(_users)
            # Users created in the last 7 days (ISO date-prefix compare — TZ-safe enough).
            _week_ago = time.strftime("%Y-%m-%d", time.gmtime(time.time() - 7 * 86400))
            new_this_week = sum(1 for u in _users if str(u.get("created_at") or "")[:10] >= _week_ago)
            # Referral tallies: how many people each user invited, and how many paid.
            _invited, _invited_paid = {}, {}
            _id_to_email = {}
            for u in _users:
                _id_to_email[str(u.get("id") or "")] = u.get("email") or ""
                rb = u.get("referred_by")
                if rb:
                    _invited[rb] = _invited.get(rb, 0) + 1
                    if u.get("referral_paid"):
                        _invited_paid[rb] = _invited_paid.get(rb, 0) + 1
            for u in _users:
                uid_u  = str(u.get("id") or "")
                status = str(u.get("subscription_status") or "free").lower()
                active = status == "active"
                if active:
                    subs_active += 1
                email  = u.get("email") or "—"
                name   = u.get("name") or "—"
                try:
                    toks_i = int(u.get("tokens_remaining", 0) or 0)
                except (TypeError, ValueError):
                    toks_i = 0
                renews = str(u.get("subscription_period_end") or "")[:10] or "—"
                joined = str(u.get("created_at") or "")[:10] or "—"
                code   = u.get("referral_code") or "—"
                inv    = _invited.get(uid_u, 0)
                inv_p  = _invited_paid.get(uid_u, 0)
                refby  = _id_to_email.get(str(u.get("referred_by") or ""), "") or "—"
                if active:
                    badge = '<span style="background:#065f46;color:#6ee7b7;padding:2px 9px;border-radius:5px;font-size:11px;font-weight:600">● active</span>'
                elif status == "canceling":
                    badge = '<span style="background:#78350f;color:#fcd34d;padding:2px 9px;border-radius:5px;font-size:11px">● canceling</span>'
                else:
                    badge = f'<span style="background:#16233f;color:#8aa0c8;padding:2px 9px;border-radius:5px;font-size:11px">{_he(status)}</span>'
                ref_html = f'<code style="color:#7cc4ff">{_he(code)}</code>'
                if inv:
                    ref_html += f' · <span style="color:#6ee7b7">{inv} invited</span>'
                    if inv_p:
                        ref_html += f' <span style="color:#8aa0c8">({inv_p} paid)</span>'
                actions = (f'<button onclick="grantTokens(\'{_js(uid_u)}\',\'{_js(str(email))}\')" '
                           'style="background:#1e3a5f;color:#cfe0ff;border:none;padding:3px 9px;border-radius:6px;cursor:pointer;font-size:12px">＋ Tokens</button>')
                if active or status == "canceling" or u.get("subscription_id"):
                    actions += (f' <button onclick="cancelSub(\'{_js(uid_u)}\',\'{_js(str(email))}\')" '
                                'style="background:#7f1d1d;color:#fecaca;border:none;padding:3px 9px;border-radius:6px;cursor:pointer;font-size:12px">Cancel</button>')
                subs_rows += f"""
                  <tr class="subrow" data-email="{_he(str(email).lower())}" data-name="{_he(str(name).lower())}" data-active="{1 if active else 0}" data-status="{_he(status)}" data-tokens="{toks_i}" data-refby="{_he(str(refby).lower())}" data-renews="{_he(renews)}" data-joined="{_he(joined)}">
                    <td><b>{_he(email)}</b></td>
                    <td style="color:#b9c9e6">{_he(name)}</td>
                    <td>{badge}</td>
                    <td style="text-align:center;font-weight:600">{toks_i}</td>
                    <td style="font-size:12px">{ref_html}</td>
                    <td style="color:#8aa0c8;font-size:12px">{_he(refby)}</td>
                    <td style="color:#8aa0c8;white-space:nowrap">{_he(renews)}</td>
                    <td style="color:#8aa0c8;white-space:nowrap">{_he(joined)}</td>
                    <td style="white-space:nowrap">{actions}</td>
                  </tr>"""
    except Exception as _e:
        subs_error = str(_e)

    monthly_price = _monthly_price_usd()  # live Stripe price (cached ~1h), USD/month
    mrr = subs_active * monthly_price
    arr = mrr * 12

    subs_th = ("padding:10px 12px;text-align:left;font-weight:600;color:#8aa0c8;"
               "background:#0f2040;border-bottom:1px solid #1a3a6e;white-space:nowrap")
    def _sth(label, idx, center=False):
        c = ";text-align:center" if center else ""
        return (f'<th onclick="subSort({idx})" style="{subs_th};cursor:pointer;user-select:none{c}">'
                f'{label} <span style="opacity:.35;font-size:10px">⇅</span></th>')
    subs_section = f"""
    <div style="margin:1.5rem 2rem">
      <h3 style="margin:0 0 .6rem;color:#e8f0ff;font-size:1rem;display:flex;align-items:center;gap:.5rem;flex-wrap:wrap">
        👥 Subscribers
        <span style="background:#4f8ef7;color:#fff;padding:2px 10px;border-radius:20px;font-size:12px">{subs_total} users</span>
        <span style="background:#065f46;color:#6ee7b7;padding:2px 10px;border-radius:20px;font-size:12px">{subs_active} active</span>
      </h3>
      {f'<p style="color:#f87171;font-size:.85rem;margin-bottom:.5rem">Could not load subscribers: {_he(subs_error)}</p>' if subs_error else ''}
      <div style="display:flex;gap:.75rem;margin-bottom:.85rem;flex-wrap:wrap">
        <div style="background:linear-gradient(135deg,#0a2f3f,#07213a);border:1px solid #16556b;border-radius:12px;padding:.75rem 1.1rem;min-width:190px">
          <div style="color:#6ee7b7;font-size:11px;font-weight:700;letter-spacing:.6px;text-transform:uppercase">MRR (est.)</div>
          <div style="font-size:1.7rem;font-weight:800;color:#e8f0ff;line-height:1.2">${mrr:,.2f}</div>
          <div style="color:#8aa0c8;font-size:11px">{subs_active} active × ${monthly_price:.2f}/mo · ARR ${arr:,.0f}</div>
        </div>
        <div style="background:#0a1628;border:1px solid #1a3a6e;border-radius:12px;padding:.75rem 1.1rem;min-width:130px">
          <div style="color:#8aa0c8;font-size:11px;font-weight:700;letter-spacing:.6px;text-transform:uppercase">Active</div>
          <div style="font-size:1.7rem;font-weight:800;color:#6ee7b7;line-height:1.2">{subs_active}</div>
          <div style="color:#8aa0c8;font-size:11px">of {subs_total} users</div>
        </div>
        <div style="background:#0a1628;border:1px solid #1a3a6e;border-radius:12px;padding:.75rem 1.1rem;min-width:130px">
          <div style="color:#8aa0c8;font-size:11px;font-weight:700;letter-spacing:.6px;text-transform:uppercase">New this week</div>
          <div style="font-size:1.7rem;font-weight:800;color:#7cc4ff;line-height:1.2">{new_this_week}</div>
          <div style="color:#8aa0c8;font-size:11px">joined in last 7 days</div>
        </div>
      </div>
      <div style="display:flex;gap:.6rem;margin-bottom:.6rem;flex-wrap:wrap;align-items:center">
        <input id="subSearch" placeholder="🔎 Search email or name…" oninput="subFilter()"
               style="background:#050d1a;border:1px solid #1a3a6e;color:#e8f0ff;padding:.45rem .7rem;border-radius:8px;font-size:13px;min-width:220px">
        <label style="display:flex;gap:.4rem;align-items:center;color:#8aa0c8;font-size:13px;cursor:pointer">
          <input type="checkbox" id="payingOnly" onchange="subFilter()"> Paying only</label>
        <button class="btn btn-gray" onclick="subExportCSV()">⬇ Export CSV</button>
        <span id="subShown" style="color:#4a5f80;font-size:12px"></span>
      </div>
      <div style="overflow-x:auto;border:1px solid #16233f;border-radius:10px">
        <table id="subTable" style="width:100%;border-collapse:collapse;font-size:13px">
          <thead><tr>
            {_sth("Email", 0)}{_sth("Name", 1)}{_sth("Status", 2)}{_sth("Tokens", 3, True)}
            <th style="{subs_th}">Referral</th>{_sth("Referred by", 5)}{_sth("Renews", 6)}{_sth("Joined", 7)}
            <th style="{subs_th}">Actions</th>
          </tr></thead>
          <tbody id="subBody">{subs_rows if subs_rows else '<tr><td colspan="9" style="padding:2rem;text-align:center;color:#4a5f80">No users yet.</td></tr>'}</tbody>
        </table>
      </div>
    </div>"""

    return f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8">
<title>Admin — Alimne</title>
<style>
  *{{box-sizing:border-box;margin:0;padding:0}}
  body{{font-family:'Segoe UI',system-ui,sans-serif;background:#050d1a;color:#e8f0ff;min-height:100vh}}
  .topbar{{display:flex;align-items:center;justify-content:space-between;padding:1rem 2rem;
           background:#0a1628;border-bottom:1px solid #1a3a6e;position:sticky;top:0;z-index:10}}
  .topbar h1{{font-size:1rem;font-weight:700;display:flex;align-items:center;gap:.6rem}}
  .badge{{background:#4f8ef7;color:#fff;padding:2px 10px;border-radius:20px;font-size:12px}}
  .badge.red{{background:#dc2626}}
  .actions{{display:flex;gap:.5rem}}
  .btn{{padding:6px 14px;border-radius:8px;border:none;cursor:pointer;font-size:13px;font-weight:600;transition:opacity .2s}}
  .btn:hover{{opacity:.8}}
  .btn-red{{background:#dc2626;color:#fff}}
  .btn-green{{background:#16a34a;color:#fff}}
  .btn-gray{{background:#1e3a5f;color:#8aa0c8}}
  .wrap{{overflow-x:auto}}
  table{{width:100%;border-collapse:collapse;font-size:13px}}
  th{{background:#0f2040;padding:10px 12px;text-align:left;font-weight:600;color:#8aa0c8;
      border-bottom:1px solid #1a3a6e;white-space:nowrap}}
  td{{padding:9px 12px;border-bottom:1px solid #0d1e35;vertical-align:middle}}
  tr:hover td{{background:#0a1e38}}
  .empty{{padding:3rem;text-align:center;color:#4a5f80;font-size:0.9rem}}
  #toast{{position:fixed;bottom:1.5rem;right:1.5rem;background:#16a34a;color:#fff;
          padding:.6rem 1.2rem;border-radius:8px;font-size:13px;display:none;z-index:100}}
</style></head>
<body>
<div class="topbar">
  <h1>📊 Alimne — Admin
    <span class="badge">{len(vis_copy)} visits</span>
    {f'<span class="badge red">⛔ {len(blocked_copy)} blocked</span>' if blocked_copy else ''}
  </h1>
  <div class="actions">
    <button class="btn btn-gray" onclick="location.reload()">↻ Refresh</button>
    <button class="btn btn-red" onclick="clearLog()">🗑 Clear Log</button>
    <button class="btn btn-gray" onclick="forgetToken()" title="Forget the saved admin token on this device">⎋ Sign out</button>
  </div>
</div>
{subs_section}
{blocked_section}
<h3 style="margin:1.5rem 2rem .5rem;color:#8aa0c8;font-size:.95rem">🌐 Recent visitors</h3>
<div class="wrap">
<table>
  <thead><tr>
    <th>Time (UTC)</th><th>IP Address</th><th>Country</th>
    <th>Location</th><th>ISP / Org</th><th>Path</th><th>User Agent</th><th>Action</th>
  </tr></thead>
  <tbody>
  {rows if rows else '<tr><td colspan="8" class="empty">No visitors recorded yet.</td></tr>'}
  </tbody>
</table>
</div>
<div id="toast"></div>
<script>
const TOKEN = "{str(token).replace(chr(92), chr(92)*2).replace(chr(34), chr(92)+chr(34))}";
// Remember the token on this device so /admin auto-opens next time, and strip
// the token out of the address bar / history now that it's saved.
try {{ localStorage.setItem("alimne_admin_token", TOKEN); if (location.search) history.replaceState(null, "", "/admin"); }} catch (e) {{}}
function toast(msg, color="#16a34a"){{
  const t = document.getElementById("toast");
  t.textContent = msg; t.style.background = color; t.style.display = "block";
  setTimeout(()=>t.style.display="none", 2500);
}}
// Send the admin token in the X-Admin-Token HEADER, never the URL query string
// (query strings are recorded in access logs; the header is not).
async function blockIp(ip){{
  if(!confirm("Block " + ip + "?\\nThis will 403 all their requests immediately.")) return;
  const r = await fetch("/admin/block", {{
    method:"POST", headers:{{"Content-Type":"application/json","X-Admin-Token":TOKEN}},
    body: JSON.stringify({{ip}})
  }});
  if(r.ok){{ toast("⛔ Blocked: " + ip, "#dc2626"); setTimeout(()=>location.reload(),1200); }}
}}
async function unblock(ip){{
  const r = await fetch("/admin/unblock", {{
    method:"POST", headers:{{"Content-Type":"application/json","X-Admin-Token":TOKEN}},
    body: JSON.stringify({{ip}})
  }});
  if(r.ok){{ toast("✓ Unblocked: " + ip); setTimeout(()=>location.reload(),1200); }}
}}
async function clearLog(){{
  if(!confirm("Clear all visitor log entries?")) return;
  const r = await fetch("/admin/clear", {{method:"POST", headers:{{"X-Admin-Token":TOKEN}}}});
  if(r.ok){{ toast("🗑 Log cleared"); setTimeout(()=>location.reload(),1200); }}
}}

// ── Subscribers: search / paying-only filter / sort / CSV / actions ──────────
function subFilter(){{
  var q = (document.getElementById("subSearch").value || "").toLowerCase().trim();
  var payingOnly = document.getElementById("payingOnly").checked;
  var rows = document.querySelectorAll("#subBody tr.subrow");
  var shown = 0;
  rows.forEach(function(r){{
    var okQ = !q || r.dataset.email.indexOf(q) >= 0 || r.dataset.name.indexOf(q) >= 0 || (r.dataset.refby||"").indexOf(q) >= 0;
    var okP = !payingOnly || r.dataset.active === "1";
    var vis = okQ && okP;
    r.style.display = vis ? "" : "none";
    if (vis) shown++;
  }});
  var el = document.getElementById("subShown");
  if (el) el.textContent = shown + " shown";
}}
var _subSort = {{}};
var _SUBCOLS = {{0:["email",0], 1:["name",0], 2:["status",0], 3:["tokens",1], 5:["refby",0], 6:["renews",0], 7:["joined",0]}};
function subSort(idx){{
  var spec = _SUBCOLS[idx]; if(!spec) return;
  var key = spec[0], numeric = spec[1];
  var dir = _subSort[idx] === 1 ? -1 : 1; _subSort = {{}}; _subSort[idx] = dir;
  var body = document.getElementById("subBody");
  var rows = Array.prototype.slice.call(body.querySelectorAll("tr.subrow"));
  rows.sort(function(a,b){{
    var va = a.dataset[key] || "", vb = b.dataset[key] || "";
    if (numeric) return ((parseFloat(va)||0) - (parseFloat(vb)||0)) * dir;
    return (va < vb ? -1 : (va > vb ? 1 : 0)) * dir;
  }});
  rows.forEach(function(r){{ body.appendChild(r); }});
}}
function subExportCSV(){{
  var rows = document.querySelectorAll("#subBody tr.subrow");
  var out = [["Email","Name","Status","Tokens","Referred by","Renews","Joined"]];
  rows.forEach(function(r){{
    if (r.style.display === "none") return;
    out.push([r.dataset.email, r.dataset.name, r.dataset.status, r.dataset.tokens, r.dataset.refby||"", r.dataset.renews||"", r.dataset.joined||""]);
  }});
  var csv = out.map(function(row){{ return row.map(function(c){{ return '"' + String(c).replace(/"/g,'""') + '"'; }}).join(","); }}).join("\\n");
  var blob = new Blob([csv], {{type:"text/csv"}});
  var a = document.createElement("a");
  a.href = URL.createObjectURL(blob); a.download = "alimne-subscribers.csv";
  document.body.appendChild(a); a.click(); a.remove();
  toast("⬇ Exported " + (out.length - 1) + " rows");
}}
async function grantTokens(uid, email){{
  var v = prompt("Grant tokens to " + email + "\\n(use a negative number to deduct):", "30");
  if (v === null) return;
  var amount = parseInt(v, 10);
  if (!amount) {{ toast("Enter a non-zero number", "#dc2626"); return; }}
  var r = await fetch("/admin/user/grant", {{method:"POST", headers:{{"Content-Type":"application/json","X-Admin-Token":TOKEN}}, body: JSON.stringify({{user_id: uid, amount: amount}})}});
  var d = await r.json().catch(function(){{ return {{}}; }});
  if (r.ok) {{ toast("✓ " + email + ": " + d.tokens_remaining + " tokens"); setTimeout(function(){{ location.reload(); }}, 900); }}
  else {{ toast("✗ " + (d.error || "failed"), "#dc2626"); }}
}}
async function cancelSub(uid, email){{
  if (!confirm("Cancel subscription for " + email + "?\\nStripe subscriptions cancel at period end (they keep access until then).")) return;
  var r = await fetch("/admin/user/cancel", {{method:"POST", headers:{{"Content-Type":"application/json","X-Admin-Token":TOKEN}}, body: JSON.stringify({{user_id: uid}})}});
  var d = await r.json().catch(function(){{ return {{}}; }});
  if (r.ok) {{ toast(d.canceled_at_period_end ? "✓ Cancels at period end" : "✓ Set to free"); setTimeout(function(){{ location.reload(); }}, 900); }}
  else {{ toast("✗ " + (d.error || "failed"), "#dc2626"); }}
}}
function forgetToken(){{
  try {{ localStorage.removeItem("alimne_admin_token"); }} catch(e) {{}}
  location.href = "/admin";
}}
subFilter();
</script>
</body></html>"""

_SAFE_NAME = re.compile(r'[^\w\-. ]')
_JOB_ID_RE = re.compile(r'^[0-9a-f]{32}$')

def _safe_name(s, maxlen=80):
    return _SAFE_NAME.sub('_', str(s))[:maxlen]

def _valid_job(job_id):
    return bool(_JOB_ID_RE.match(str(job_id)))

def _he(s):
    return (str(s).replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')
            .replace('"','&quot;').replace("'","&#39;"))

# ── Job store (memory ONLY — uploads and guides are never written to disk)
_JOB_TTL   = 900  # 15 minutes

# Purge any job files left on disk by older versions that persisted to /tmp
import shutil as _shutil
_shutil.rmtree(os.path.join(os.sep, "tmp", "slide-study-jobs"), ignore_errors=True)

_jobs      = OrderedDict()
_jobs_lock = threading.RLock()  # reentrant — get_job may acquire while route holds it

# Ollama can only run one inference at a time locally. Serialise all calls so
# the ThreadPoolExecutor doesn't flood it, which causes truncated JSON output.
_ollama_sem      = threading.Semaphore(1)
_DEBUG_RAW = os.environ.get("DEBUG_RAW") == "1"
_ollama_raw_lock = threading.Lock()  # protect concurrent debug-file writes

# ── Groq token-per-minute pacer ────────────────────────────────────────────────
# Groq's free tier caps tokens-per-minute (default 8000). Rather than fire calls
# and eat 30s "retry-after" waits on every 429, proactively pace: track tokens
# used in a rolling 60s window and sleep just enough to stay under the limit.
# Raise GROQ_TPM_LIMIT if you upgrade your Groq tier (Dev tier is ~250k).
GROQ_TPM_LIMIT   = int(os.environ.get("GROQ_TPM_LIMIT", "7000"))  # headroom under 8000
_groq_tpm_lock   = threading.Lock()
_groq_tpm_events = []  # list of (timestamp, tokens)

def _groq_pace(est_tokens):
    """Block until sending `est_tokens` keeps the rolling-minute total under the
    limit. A single call larger than the limit is allowed through on its own."""
    for _ in range(120):  # safety bound (~2 min max wait)
        with _groq_tpm_lock:
            now = time.time()
            while _groq_tpm_events and now - _groq_tpm_events[0][0] > 60:
                _groq_tpm_events.pop(0)
            used = sum(t for _, t in _groq_tpm_events)
            if not _groq_tpm_events or used + est_tokens <= GROQ_TPM_LIMIT:
                _groq_tpm_events.append((now, est_tokens))
                return
            wait = 60 - (now - _groq_tpm_events[0][0]) + 0.5
        time.sleep(max(1.0, min(wait, 35)))

def store_job(job_id, pdf_bytes, md_text, guide, slides, filename):
    ts = time.time()
    with _jobs_lock:
        _jobs[job_id] = {
            "pdf": pdf_bytes, "md": md_text,
            "guide": guide,   "slides": slides,
            "filename": filename, "ts": ts
        }
        stale = [k for k, v in list(_jobs.items()) if time.time() - v["ts"] > _JOB_TTL]
        for k in stale:
            del _jobs[k]

def get_job(job_id):
    with _jobs_lock:
        job = _jobs.get(job_id)
        if job and time.time() - job["ts"] > _JOB_TTL:
            del _jobs[job_id]
            return None
    return job

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY        = colors.HexColor('#0a1628')
NAVY_MID    = colors.HexColor('#1a3a6e')
NAVY_LIGHT  = colors.HexColor('#2e5ca8')
ACCENT      = colors.HexColor('#4f8ef7')
KW_BG       = colors.HexColor('#eef3ff')
ROW_ALT     = colors.HexColor('#f4f7ff')
BORDER      = colors.HexColor('#c5d8ff')
TEXT        = colors.HexColor('#0d1b2e')
TEXT_LIGHT  = colors.HexColor('#4a5f80')
WHITE       = colors.white
CARD_Q      = colors.HexColor('#1a3a6e')
CARD_A      = colors.HexColor('#f8faff')
GREEN       = colors.HexColor('#16a34a')
SECTION_BG  = colors.HexColor('#f0f5ff')


# ── Ollama helpers ─────────────────────────────────────────────────────────────

def ollama_running():
    if GROQ_API_KEY:
        return True
    try:
        return http.get(f"{OLLAMA_URL}/api/tags", timeout=3).status_code == 200
    except Exception:
        return False

def ollama_models():
    if GROQ_API_KEY:
        return [GROQ_MODEL]
    try:
        return [m["name"] for m in http.get(f"{OLLAMA_URL}/api/tags", timeout=3).json().get("models", [])]
    except Exception:
        return []

def _extract_json(text):
    text = re.sub(r'^```(?:json)?\s*', '', text.strip(), flags=re.IGNORECASE)
    text = re.sub(r'\s*```$', '', text).strip()

    # ── Pass 1: try the whole text as-is ────────────────────────────────────
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass

    # ── Find the start of the first JSON object or array ───────────────────
    obj_pos = text.find('{')
    arr_pos = text.find('[')
    if obj_pos == -1 and arr_pos == -1:
        raise ValueError("No JSON found in model output")
    if arr_pos != -1 and (obj_pos == -1 or arr_pos < obj_pos):
        start, open_c, close_c = arr_pos, '[', ']'
    else:
        start, open_c, close_c = obj_pos, '{', '}'

    # ── Pass 2: string-aware bracket scan ──────────────────────────────────
    # Tracks whether we're inside a JSON string so { } inside strings are ignored
    depth, in_str, escaped, end = 0, False, False, -1
    stack = []          # track every opening bracket for repair below
    for i, ch in enumerate(text[start:], start):
        if escaped:
            escaped = False
            continue
        if ch == '\\' and in_str:
            escaped = True
            continue
        if ch == '"':
            in_str = not in_str
            continue
        if in_str:
            continue
        if ch in ('{', '['):
            depth += 1
            stack.append(ch)
        elif ch in ('}', ']'):
            depth -= 1
            if stack:
                stack.pop()
            if depth == 0:
                end = i
                break

    if end != -1:
        try:
            return json.loads(text[start:end + 1])
        except json.JSONDecodeError:
            pass

    # ── Pass 3: JSON is truncated — repair by closing open brackets ─────────
    # Re-scan to find the deepest valid position and close everything open
    pairs = {'{': '}', '[': ']'}
    depth2, in_str2, escaped2 = 0, False, False
    open_stack = []
    last_good_pos = start  # last position where depth was 0 between top-level items
    for i, ch in enumerate(text[start:], start):
        if escaped2:
            escaped2 = False
            continue
        if ch == '\\' and in_str2:
            escaped2 = True
            continue
        if ch == '"':
            in_str2 = not in_str2
            continue
        if in_str2:
            continue
        if ch in ('{', '['):
            depth2 += 1
            open_stack.append(ch)
        elif ch in ('}', ']'):
            depth2 -= 1
            if open_stack:
                open_stack.pop()

    # Trim trailing incomplete entry: remove everything after the last comma
    # at depth==1 so the partial last entry is dropped
    snippet = text[start:].rstrip()
    # Remove trailing comma + anything after it (the cut-off entry)
    snippet = re.sub(r',\s*[^,\]\}]*$', '', snippet)
    # Close all still-open brackets
    closing = ''.join(pairs[c] for c in reversed(open_stack))
    # Make sure we haven't over-trimmed: if snippet is just the opening char, add empty body
    repaired = snippet + closing
    try:
        return json.loads(repaired)
    except json.JSONDecodeError:
        pass

    # ── Pass 4: nuclear fallback — try every possible closing combination ──
    for suffix in ['}', ']}', '}}', '"]}}', '"]}', '"}']:
        try:
            return json.loads(snippet + suffix)
        except json.JSONDecodeError:
            continue

    raise ValueError(f"Cannot parse model output (length={len(text)}). "
                     "Model may have returned incomplete JSON.")

def _call_ollama(prompt, retries=3, num_predict=4096):
    if GROQ_API_KEY:
        return _call_groq(prompt, retries=retries, max_tokens=num_predict)
    payload = {
        "model": OLLAMA_MODEL,
        "format": "json",
        "messages": [
            {"role": "system", "content": "You output only valid JSON. No markdown, no explanation."},
            {"role": "user",   "content": prompt},
        ],
        "stream": False,
        "options": {"temperature": 0, "num_predict": num_predict},
    }
    last_err = None
    for attempt in range(retries):
        try:
            with _ollama_sem:
                r = http.post(f"{OLLAMA_URL}/api/chat", json=payload, timeout=300)
                r.raise_for_status()
                raw_content = r.json()["message"]["content"]

            if _DEBUG_RAW:
                with _ollama_raw_lock:
                    with open(os.path.join(os.path.dirname(__file__), "ollama_raw.txt"), "w", encoding="utf-8") as _f:
                        _f.write(raw_content)

            return _extract_json(raw_content)
        except Exception as e:
            last_err = e
            _log.warning("OLLAMA attempt %d/%d failed: %s", attempt + 1, retries, e)
            if attempt < retries - 1:
                time.sleep(1.5 * (attempt + 1))
    raise last_err


def _call_groq(prompt, retries=5, max_tokens=2048):
    if not GROQ_API_KEY:
        raise ValueError("GROQ_API_KEY is not set — configure it in Render environment variables.")
    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json",
    }
    # Try the primary model, then fall back to alternates if the model is
    # unavailable to this key (deprecated / re-tiered → 404/400 model_not_found).
    models = [GROQ_MODEL] + [m for m in GROQ_FALLBACK_MODELS if m != GROQ_MODEL]
    last_err = None
    for model in models:
        payload = {
            "model": model,
            "messages": [
                {"role": "system", "content":
                    "You are a precise study-guide generator. Output ONLY valid JSON — no markdown, no commentary. "
                    "Use ONLY the information the user provides. Never invent, guess, add, or rename facts, people, "
                    "places, terms, acronyms, symbols or numbers. Copy every name and technical term exactly as it "
                    "appears in the source; do not translate or alter names."},
                {"role": "user",   "content": prompt},
            ],
            "temperature": 0,
            "max_tokens": max_tokens,
        }
        # Estimate this call's token cost (prompt ≈ chars/4, plus the reserved
        # output) and pace to stay under the per-minute limit before sending.
        est_tokens = len(prompt) // 4 + max_tokens + 120
        for attempt in range(retries):
            try:
                _groq_pace(est_tokens)
                r = http.post(
                    "https://api.groq.com/openai/v1/chat/completions",
                    json=payload, headers=headers, timeout=120
                )
                if r.status_code == 429:
                    wait = int(r.headers.get("retry-after", 10))
                    _log.warning("GROQ rate limited (429) — waiting %ds. body=%s", wait, r.text[:300])
                    last_err = RuntimeError(f"GROQ rate limited (429): {r.text[:200]}")
                    time.sleep(wait)
                    continue
                if r.status_code in (400, 404):
                    # Model-level problem (unknown/deprecated/no-access) — don't
                    # burn retries; move on to the next fallback model.
                    _log.error("GROQ HTTP %s (model=%s) — trying next model. body=%s",
                               r.status_code, model, r.text[:400])
                    last_err = RuntimeError(f"GROQ {r.status_code} for model {model}: {r.text[:200]}")
                    break
                if r.status_code >= 400:
                    _log.error("GROQ HTTP %s (model=%s): %s", r.status_code, model, r.text[:400])
                r.raise_for_status()
                raw_content = r.json()["choices"][0]["message"]["content"]
                if model != GROQ_MODEL:
                    _log.warning("GROQ served by fallback model %s", model)
                return _extract_json(raw_content)
            except Exception as e:
                last_err = e
                _log.warning("GROQ attempt %d/%d (model=%s) failed: %s", attempt + 1, retries, model, e)
                if attempt < retries - 1:
                    time.sleep(2 * (attempt + 1))
    if last_err is None:
        last_err = RuntimeError("GROQ call failed — all models/retries exhausted")
    raise last_err


# ── Three-pass AI processing ──────────────────────────────────────────────────

def _lang_rules(language):
    """Faithfulness + language rules injected into every prompt so the model
    stays true to the source and outputs cleanly in the requested language."""
    if language == "ar":
        lang_line = ("Write every text value in Modern Standard Arabic (العربية). "
                     "Use correct, natural Arabic grammar and spelling.")
    else:
        lang_line = "Write every text value in clear English."
    return (
        "- " + lang_line + " Keep all JSON keys in English exactly as shown.\n"
        "- ACCURACY IS CRITICAL: use ONLY facts that appear in the source above. "
        "Do NOT invent, assume, or add anything that is not stated in the source.\n"
        "- Copy every name, person, place, acronym, symbol (e.g. ATP, CO2, DNA) and number "
        "EXACTLY as written in the source. Never rename, translate, transliterate, or alter a "
        "proper noun, technical term, or formula — even when writing in Arabic.\n"
        "- If the source does not cover something, leave it out rather than making it up."
    )


def _as_dict(result, list_key=None):
    """If the model returned a list instead of a dict, coerce it."""
    if isinstance(result, dict):
        return result
    if isinstance(result, list):
        # If list contains dicts, pick the first one
        first_dict = next((x for x in result if isinstance(x, dict)), None)
        if first_dict:
            return first_dict
        # Otherwise wrap under the given key
        if list_key:
            return {list_key: result}
    return {}


def _detect_language(content):
    """Detect 'ar' or 'en' from slides list or plain text. Based on Arabic char ratio."""
    if isinstance(content, list):
        sample = " ".join(
            f"{s.get('title', '')} {s.get('content', '')}" for s in content[:20]
        )
    else:
        sample = str(content)[:6000]
    arabic = sum(1 for c in sample if '؀' <= c <= 'ۿ')
    alpha  = sum(1 for c in sample if c.isalpha())
    return "ar" if alpha > 0 and arabic / alpha > 0.25 else "en"


def pass1_overview(slides, language, dcfg=None):
    """Get title, objectives, section groupings, and keywords."""
    dcfg = dcfg or DETAIL["standard"]
    lang = "in Arabic" if language == "ar" else "in English"
    # Include ALL slides for grouping; use title-only for slides beyond max_slides
    # to keep the LLM input within token limits for large presentations.
    outline = "".join(
        f"Slide {s['slide_num']}: {s['title']}\n{s['content'][:dcfg['slide_chars']]}\n\n"
        if i < dcfg["max_slides"] else
        f"Slide {s['slide_num']}: {s['title']}\n"
        for i, s in enumerate(slides)
    )
    result = _call_ollama(f"""Create a study guide overview {lang} from this source material.

{outline}

Return JSON:
{{
  "title": "TOPIC IN CAPS",
  "subtitle": "course/session description",
  "objectives": ["objective 1", "objective 2", "objective 3"],
  "sections": [
    {{"title": "SECTION NAME", "slide_nums": [1, 2, 3]}}
  ],
  "keywords": [
    {{"term": "Term", "definition": "Specific 1-sentence definition or role description"}}
  ]
}}

Rules:
- Group related slides into 3-7 logical sections
- objectives: take from a learning-objectives section if present, otherwise summarise the actual content (do not invent goals)
- keywords: {dcfg['keywords']} terms — only concepts, acronyms, roles and processes that actually appear in the source
{_lang_rules(language)}
- Output JSON only""", num_predict=dcfg["num_predict"])
    return _as_dict(result)


def pass2_section(title, section_slides, language, dcfg=None):
    """Get detailed bullets + optional comparison table for one section."""
    dcfg = dcfg or DETAIL["standard"]
    lang = "in Arabic" if language == "ar" else "in English"
    # Cap per-slide and total content so a crafted file can't blow up the Groq
    # token/cost bill (pass 1 caps slide count, but pass 2 sent full content).
    _per = dcfg.get("slide_chars", 700) * 3
    content = "".join(
        f"Slide {s['slide_num']}: {s['title']}\n{s['content'][:_per]}\n\n"
        for s in section_slides
    )[:60_000]
    result = _call_ollama(f"""Extract detailed exam study notes {lang} for the section "{title}".

{content}

Return JSON:
{{
  "bullets": [
    "Full specific fact, definition, or step from the content",
    "Another detailed point — include names, numbers, roles, processes"
  ],
  "table": {{
    "headers": ["Column 1", "Column 2"],
    "rows": [["value", "value"]]
  }}
}}

Rules:
- bullets: {dcfg['bullets']} specific, exam-worthy facts taken directly from the content above
- table: include ONLY if content has roles/comparisons/structured lists; otherwise omit the table field entirely
{_lang_rules(language)}
- Output JSON only""", num_predict=dcfg["num_predict"])
    result = _as_dict(result, list_key="bullets")
    # If model returned a flat list of strings under "bullets", normalise each element
    bullets = result.get("bullets", [])
    result["bullets"] = [str(b) for b in bullets if b]
    return result


def pass3_flashcards(guide, language, dcfg=None):
    """Generate Q&A flash cards from the guide content."""
    dcfg = dcfg or DETAIL["standard"]
    lang = "in Arabic" if language == "ar" else "in English"
    kw_text  = ", ".join(k["term"] for k in guide.get("keywords", [])[:20] if isinstance(k, dict))
    sec_text = " | ".join(s["title"] for s in guide.get("sections", []) if isinstance(s, dict))
    bullets_ctx = ""
    for sec in guide.get("sections", []):
        if isinstance(sec, dict) and sec.get("bullets"):
            bullets_ctx += f"\n{sec.get('title','')}:\n" + "\n".join(f"- {b}" for b in sec["bullets"])
    result = _call_ollama(f"""Create exam flash cards {lang} for a study guide about: {guide.get('title', '')}.

Topics: {sec_text}
Key terms: {kw_text}
Study content:{bullets_ctx}

Return JSON:
{{
  "flashcards": [
    {{"q": "Question that tests a key concept?", "a": "Clear, concise answer"}},
    {{"q": "Define [term]?", "a": "Definition"}},
    {{"q": "What are the responsibilities of [role]?", "a": "Specific responsibilities"}}
  ]
}}

Rules:
- Create exactly {dcfg['n_flash']} flash cards
- Base EVERY question and answer directly on the study content provided above — no generic or invented questions
- Mix definition questions, "what is" questions, "name the" questions, and role/responsibility questions
- Answers must be specific and factually match the content
{_lang_rules(language)}
- Output JSON only""", num_predict=dcfg["num_predict"])
    # Model may return the array directly instead of wrapping it
    if isinstance(result, list):
        return {"flashcards": [x for x in result if isinstance(x, dict)]}
    return _as_dict(result, list_key="flashcards")


def pass4_mcq(guide, language, dcfg=None):
    """Generate multiple-choice quiz questions."""
    dcfg   = dcfg or DETAIL["standard"]
    lang   = "in Arabic" if language == "ar" else "in English"
    n      = dcfg["n_mcq"]
    kw_text  = ", ".join(k["term"] for k in guide.get("keywords", [])[:20] if isinstance(k, dict))
    sec_text = " | ".join(s["title"] for s in guide.get("sections", []) if isinstance(s, dict))
    bullets_ctx = ""
    for sec in guide.get("sections", []):
        if isinstance(sec, dict) and sec.get("bullets"):
            bullets_ctx += f"\n{sec.get('title','')}:\n" + "\n".join(f"- {b}" for b in sec["bullets"])
    result = _call_ollama(f"""Create {n} multiple-choice exam questions {lang} for: {guide.get('title','')}.

Topics: {sec_text}
Key terms: {kw_text}
Study content:{bullets_ctx}

Return JSON:
{{
  "mcqs": [
    {{
      "q": "Question text?",
      "options": ["A. option one", "B. option two", "C. option three", "D. option four"],
      "answer": "A",
      "explanation": "Why A is correct"
    }}
  ]
}}

Rules:
- Exactly {n} questions, 4 options each (A B C D)
- Every question, the correct option, and the explanation must be grounded in the study content above — do not invent facts or use outside knowledge
- answer: just the letter
- Mix easy and hard questions
{_lang_rules(language)}
- JSON only""", num_predict=dcfg["num_predict"])
    if isinstance(result, list):
        return {"mcqs": [x for x in result if isinstance(x, dict)]}
    return _as_dict(result, list_key="mcqs")


def _sections_parallel(sections, content_slides, language, dcfg):
    """Process sections sequentially (Groq rate limits prevent safe concurrency).
    Yields plain dicts. Flashcards+MCQ are parallelized separately."""
    n = len(sections)
    for i, sec in enumerate(sections):
        yield {"step": "section", "msg": f"Section {i+1}/{n}: {sec.get('title', '')}…"}
        nums = set(sec.get("slide_nums", []))
        sl = [s for s in content_slides if s["slide_num"] in nums]
        if not sl:
            chunk = max(1, len(content_slides) // n)
            start = i * chunk
            end   = start + chunk if i < n - 1 else len(content_slides)
            sl    = content_slides[start:end] or content_slides
        try:
            det = pass2_section(sec.get("title", ""), sl, language, dcfg)
            sec["bullets"] = det.get("bullets", [])
            if isinstance(det.get("table"), dict):
                sec["table"] = det["table"]
        except Exception:
            _log.error("pass2 [%s] error:\n%s", sec.get("title", "?"), _tb.format_exc())
            sec["bullets"] = []
        yield {"step": "section", "msg": f"Sections: {i+1}/{n} done…"}


def _flashcards_mcq_parallel(overview, language, dcfg, include_quiz=True):
    """Run pass3 then pass4 sequentially (Groq free tier rate limits concurrent calls).
    Yields plain dicts. When include_quiz is False, skip flashcards + quiz (summary-only)."""
    if not include_quiz:
        overview["flashcards"] = []
        overview["mcqs"] = []
        yield {"step": "summary", "msg": "Summary-only mode — skipping flash cards and quiz…"}
        return
    yield {"step": "flashcards", "msg": "Generating flash cards…"}
    try:
        overview["flashcards"] = pass3_flashcards(overview, language, dcfg).get("flashcards", [])
    except Exception:
        _log.error("pass3 error:\n%s", _tb.format_exc())
        overview["flashcards"] = []
    yield {"step": "flashcards", "msg": "Flash cards ready…"}

    yield {"step": "mcq", "msg": "Generating quiz…"}
    try:
        overview["mcqs"] = pass4_mcq(overview, language, dcfg).get("mcqs", [])
    except Exception:
        _log.error("pass4 error:\n%s", _tb.format_exc())
        overview["mcqs"] = []
    yield {"step": "mcq", "msg": "Quiz ready…"}


def build_markdown(guide):
    """Convert guide dict to a Markdown string (labels follow the guide language)."""
    is_ar = guide.get("language") == "ar"
    L = {
        "title":      "دليل الدراسة" if is_ar else "Study Guide",
        "objectives": "الأهداف التعليمية" if is_ar else "Learning Objectives",
        "keywords":   "قاموس المصطلحات" if is_ar else "Keywords Cheatsheet",
        "flashcards": "بطاقات المراجعة" if is_ar else "Flash Cards",
        "quiz":       "أسئلة الاختيار من متعدد" if is_ar else "Multiple Choice Questions",
        "q":          "سؤال" if is_ar else "Q",
        "a":          "الإجابة" if is_ar else "A",
    }
    lines = [f"# {guide.get('title', L['title'])}", ""]
    if guide.get("subtitle"):
        lines += [f"*{guide['subtitle']}*", ""]
    objs = [o for o in guide.get("objectives", []) if isinstance(o, str)]
    if objs:
        lines += [f"## {L['objectives']}", ""]
        for o in objs: lines.append(f"- {o}")
        lines.append("")
    for sec in guide.get("sections", []):
        if not isinstance(sec, dict): continue
        lines += [f"## {sec.get('title', '')}", ""]
        for b in sec.get("bullets", []): lines.append(f"- {b}")
        tbl = sec.get("table")
        if isinstance(tbl, dict) and tbl.get("headers") and tbl.get("rows"):
            lines.append("")
            lines.append("| " + " | ".join(tbl["headers"]) + " |")
            lines.append("|" + "|".join(["---"] * len(tbl["headers"])) + "|")
            for row in tbl["rows"]:
                lines.append("| " + " | ".join(str(c) for c in row) + " |")
        lines.append("")
    kws = [k for k in guide.get("keywords", []) if isinstance(k, dict)]
    if kws:
        lines += [f"## {L['keywords']}", ""]
        for k in kws: lines.append(f"**{k.get('term','')}** — {k.get('definition','')}")
        lines.append("")
    fcs = [f for f in guide.get("flashcards", []) if isinstance(f, dict)]
    if fcs:
        lines += [f"## {L['flashcards']}", ""]
        for i, fc in enumerate(fcs, 1):
            lines += [f"**{L['q']}{i}:** {fc.get('q','')}", f"**{L['a']}:** {fc.get('a','')}", ""]
    mcqs = [m for m in guide.get("mcqs", []) if isinstance(m, dict)]
    if mcqs:
        lines += [f"## {L['quiz']}", ""]
        for i, m in enumerate(mcqs, 1):
            lines.append(f"**{i}. {m.get('q','')}**")
            for opt in m.get("options", []): lines.append(f"   {opt}")
            lines += [f"   ✓ **{m.get('answer','')}** — {m.get('explanation','')}", ""]
    return "\n".join(lines)


def ask_ollama(slides, language, progress_cb=None):
    """Orchestrate three-pass processing with optional progress callback."""
    content_slides = [s for s in slides if s["content"].strip() or s["title"].strip()]

    if progress_cb: progress_cb("overview", f"Analysing structure ({len(content_slides)} slides)…")
    overview = pass1_overview(content_slides, language)

    # Normalize sections/keywords (mistral may return plain strings)
    raw_sec = overview.get("sections", [])
    overview["sections"] = [
        s if isinstance(s, dict) else {"title": str(s), "slide_nums": []}
        for s in (raw_sec if isinstance(raw_sec, list) else [])
    ]
    raw_kw = overview.get("keywords", [])
    overview["keywords"] = [
        k if isinstance(k, dict) else {"term": str(k), "definition": ""}
        for k in (raw_kw if isinstance(raw_kw, list) else [])
    ]
    if not overview["sections"]:
        overview["sections"] = [{"title": overview.get("title", "Overview"), "slide_nums": [s["slide_num"] for s in content_slides]}]

    sections = overview["sections"]
    n = len(sections)
    for i, sec in enumerate(sections):
        if progress_cb: progress_cb("section", f"Building section {i+1} of {n}: {sec.get('title','')}…")
        nums = set(sec.get("slide_nums", []))
        sl = [s for s in content_slides if s["slide_num"] in nums]
        if not sl:
            chunk = max(1, len(content_slides) // n)
            start = i * chunk
            end   = start + chunk if i < n - 1 else len(content_slides)
            sl    = content_slides[start:end] or content_slides
        detail = pass2_section(sec.get("title", ""), sl, language)
        sec["bullets"] = detail.get("bullets", [])
        if isinstance(detail.get("table"), dict):
            sec["table"] = detail["table"]

    if progress_cb: progress_cb("flashcards", "Generating flash cards…")
    try:
        fc = pass3_flashcards(overview, language)
        overview["flashcards"] = fc.get("flashcards", [])
    except Exception:
        overview["flashcards"] = []

    return overview


# ── File extraction — PPTX, PPT binary, PDF ───────────────────────────────────

def _extract_pptx_raw(raw):
    """Open as pptx; on relationship errors fall back to raw ZIP XML walk."""
    try:
        prs = Presentation(io.BytesIO(raw))
    except Exception as e:
        if "officeDocument" not in str(e) and "relationship" not in str(e):
            raise ValueError(f"Cannot open PowerPoint file: {e}")
        # Fallback: read slide XML directly from the ZIP
        import zipfile, xml.etree.ElementTree as ET
        slides = []
        with zipfile.ZipFile(io.BytesIO(raw)) as z:
            slide_files = sorted(
                [n for n in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml', n)],
                key=lambda x: int(re.findall(r'\d+', x)[-1])
            )
            for idx, sf in enumerate(slide_files, 1):
                parts = [
                    t.text.strip()
                    for t in ET.fromstring(z.read(sf)).iter(
                        '{http://schemas.openxmlformats.org/drawingml/2006/main}t'
                    )
                    if t.text and t.text.strip()
                ]
                if parts:
                    slides.append({"slide_num": idx, "title": parts[0][:80], "content": "\n".join(parts)})
        if not slides:
            raise ValueError("No readable content found in this file")
        return slides

    if len(prs.slides) > _MAX_PAGES:
        raise ValueError(f"Presentation has too many slides (max {_MAX_PAGES}).")
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        ts    = slide.shapes.title
        title = ts.text.strip() if ts and ts.text.strip() else f"Slide {i}"
        body  = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape != ts and shape.text.strip()
        ]
        slides.append({"slide_num": i, "title": title, "content": "\n".join(body)})
    return slides


def _extract_ppt_binary(raw):
    """Extract text from old binary OLE2 .ppt via record-level parsing."""
    import struct
    try:
        import olefile
    except ImportError:
        raise ValueError("olefile not installed — run: pip install olefile")

    ole = olefile.OleFileIO(io.BytesIO(raw))
    if not ole.exists('PowerPoint Document'):
        raise ValueError("Not a valid binary PowerPoint file")

    stream = ole.openstream('PowerPoint Document').read()
    texts  = []

    def parse(buf, start, end):
        i = start
        while i + 8 <= end:
            rec_ver  = struct.unpack_from('<H', buf, i)[0] & 0x0F
            rec_type = struct.unpack_from('<H', buf, i + 2)[0]
            rec_len  = struct.unpack_from('<I', buf, i + 4)[0]
            data_end = i + 8 + rec_len
            if data_end > end:
                break
            if rec_type == 0x0FA0:       # TextCharsAtom — UTF-16LE
                t = buf[i+8:data_end].decode('utf-16-le', errors='ignore').strip()
                if t: texts.append(t)
            elif rec_type == 0x0FA8:     # TextBytesAtom — Latin-1
                t = buf[i+8:data_end].decode('latin-1', errors='ignore').strip()
                if t: texts.append(t)
            elif rec_ver == 0xF:         # Container — recurse into children
                parse(buf, i + 8, data_end)
            i = data_end

    parse(stream, 0, len(stream))

    if not texts:
        raise ValueError("No readable text found in the binary .ppt file")

    # Group into pseudo-slides of ~5 text blocks
    slides = []
    for n, chunk in enumerate([texts[j:j+5] for j in range(0, len(texts), 5)], 1):
        slides.append({"slide_num": n, "title": chunk[0][:80], "content": "\n".join(chunk)})
    return slides


def _extract_pdf(raw):
    """Extract text from PDF, one entry per page."""
    try:
        import pypdf
    except ImportError:
        raise ValueError("pypdf not installed — run: pip install pypdf")

    reader = pypdf.PdfReader(io.BytesIO(raw))
    if len(reader.pages) > _MAX_PAGES:
        raise ValueError(f"PDF has too many pages (max {_MAX_PAGES}).")
    slides = []
    for i, page in enumerate(reader.pages, 1):
        text  = (page.extract_text() or "").strip()
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        if lines:
            slides.append({"slide_num": i, "title": lines[0][:80], "content": "\n".join(lines)})

    if not slides:
        raise ValueError(
            "No readable text found in the PDF — it may be a scanned/image-based file."
        )
    return slides


def _extract_docx(raw):
    """Extract text from .docx (Open XML Word) files, grouped by headings."""
    try:
        from docx import Document
    except ImportError:
        raise ValueError("python-docx not installed — run: pip install python-docx")

    doc = Document(io.BytesIO(raw))
    slides, current_title, current_lines, slide_num = [], "Document", [], 1
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style.name.startswith('Heading'):
            if current_lines:
                slides.append({"slide_num": slide_num, "title": current_title, "content": "\n".join(current_lines)})
                slide_num += 1
                current_lines = []
            current_title = text[:80]
        else:
            current_lines.append(text)
    if current_lines:
        slides.append({"slide_num": slide_num, "title": current_title, "content": "\n".join(current_lines)})
    if not slides:
        raise ValueError("No readable text found in the Word document.")
    return slides


def _extract_doc_ole(raw):
    """Extract text from old binary .doc via OLE2 stream decoding."""
    import olefile
    ole = olefile.OleFileIO(io.BytesIO(raw))
    if not ole.exists('WordDocument'):
        raise ValueError("Not a valid binary Word (.doc) file.")
    stream = ole.openstream('WordDocument').read()
    # Word stores main text as UTF-16-LE; extract printable runs
    try:
        text = stream.decode('utf-16-le', errors='ignore')
    except Exception:
        text = stream.decode('latin-1', errors='ignore')
    # Keep only printable chars + Arabic/newlines, strip control chars
    text = re.sub(r'[^\x20-\x7E؀-ۿÀ-ɏ\n\r\t]+', ' ', text)
    text = re.sub(r'[ \t]{3,}', '  ', text).strip()
    lines = [l.strip() for l in text.splitlines() if len(l.strip()) > 3]
    if not lines:
        raise ValueError(
            "Could not extract readable text from this .doc file. "
            "Try saving it as .docx and uploading again."
        )
    # Group lines into pseudo-slides of 15 lines each
    slides = []
    for n, chunk in enumerate([lines[i:i+15] for i in range(0, len(lines), 15)], 1):
        slides.append({"slide_num": n, "title": chunk[0][:80], "content": "\n".join(chunk)})
    return slides


def _extract_txt(raw):
    """Extract text from a plain-text file (.txt)."""
    for enc in ('utf-8', 'utf-16', 'latin-1', 'cp1256'):
        try:
            text = raw.decode(enc, errors='strict')
            break
        except (UnicodeDecodeError, LookupError):
            continue
    else:
        text = raw.decode('utf-8', errors='replace')
    text = text.strip()
    if not text:
        raise ValueError("The text file appears to be empty.")
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    # Group into pseudo-slides of 20 lines
    slides = []
    for n, chunk in enumerate([lines[i:i+20] for i in range(0, len(lines), 20)], 1):
        slides.append({"slide_num": n, "title": chunk[0][:80], "content": "\n".join(chunk)})
    return slides


# ── Upload-parsing DoS guards ──────────────────────────────────────────────────
_MAX_UNCOMPRESSED = 300 * 1024 * 1024   # total decompressed bytes (zip-bomb guard)
_MAX_ZIP_RATIO    = 200                 # per-entry compression-ratio ceiling
_MAX_PAGES        = 1200                # hard cap on PDF pages / PPTX slides parsed

def _check_zip_bomb(raw):
    """Reject OOXML/zip inputs that would decompress to an unreasonable size.
    MAX_CONTENT_LENGTH only bounds the COMPRESSED upload; a 50 MB zip bomb can
    expand to many GB and OOM the instance."""
    import zipfile
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as z:
            total = 0
            for zi in z.infolist():
                total += zi.file_size
                if total > _MAX_UNCOMPRESSED:
                    raise ValueError("File is too large when decompressed — refusing to process.")
                if (zi.compress_size > 0 and zi.file_size > 1_000_000
                        and (zi.file_size / zi.compress_size) > _MAX_ZIP_RATIO):
                    raise ValueError("File has a suspicious compression ratio — refusing to process.")
    except zipfile.BadZipFile:
        pass  # not a real zip; the downstream parser will reject it


def extract_slides(file_stream, filename=""):
    """Detect format from magic bytes (and filename for .txt) and dispatch."""
    raw = file_stream.read()
    fname = (filename or "").lower()

    if raw[:4] == b'PK\x03\x04':                         # ZIP-based (pptx or docx)
        _check_zip_bomb(raw)
        import zipfile
        try:
            with zipfile.ZipFile(io.BytesIO(raw)) as z:
                names = z.namelist()
            if any(n.startswith('word/') for n in names):
                return _extract_docx(raw)
        except Exception:
            pass
        return _extract_pptx_raw(raw)

    if raw[:8] == b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1':  # OLE2 (ppt or doc)
        try:
            import olefile
            with olefile.OleFileIO(io.BytesIO(raw)) as ole:
                if ole.exists('WordDocument'):
                    return _extract_doc_ole(raw)
        except Exception:
            pass
        return _extract_ppt_binary(raw)

    if raw[:4] == b'%PDF':                                # PDF
        return _extract_pdf(raw)

    if fname.endswith('.txt') or fname.endswith('.md'):   # Plain text
        return _extract_txt(raw)

    # Last resort: try decoding as plain text
    try:
        decoded = raw.decode('utf-8', errors='strict')
        if len(decoded.strip()) > 50:
            return _extract_txt(raw)
    except UnicodeDecodeError:
        pass

    raise ValueError(
        "Unrecognised file format. Supported: .pptx, .ppt, .pdf, .docx, .doc, .txt"
    )


# ── PDF builder ────────────────────────────────────────────────────────────────

def _st(name, parent, **kw):
    s = ParagraphStyle(name, parent=parent)
    for k, v in kw.items(): setattr(s, k, v)
    return s

def _page_num(canvas, doc):
    canvas.saveState()
    canvas.setFont('Helvetica', 7.5)
    canvas.setFillColor(TEXT_LIGHT)
    n = canvas.getPageNumber()
    canvas.drawRightString(doc.width + doc.leftMargin, 0.55*cm, f"{n}")
    canvas.drawString(doc.leftMargin, 0.55*cm, doc.title_str if hasattr(doc, 'title_str') else '')
    canvas.restoreState()


def build_pdf(guide, language, out_filename="study_guide"):
    if not isinstance(guide, dict):
        guide = {}
    guide["sections"]   = [s for s in guide.get("sections",   []) if isinstance(s, dict)]
    guide["keywords"]   = [k for k in guide.get("keywords",   []) if isinstance(k, dict)]
    guide["flashcards"] = [f for f in guide.get("flashcards", []) if isinstance(f, dict)]
    guide["objectives"] = [o for o in guide.get("objectives", []) if isinstance(o, str)]

    is_ar = (language == "ar")
    ar_ok = is_ar and _ensure_arabic_font()
    AF    = _ARABIC_FONT if ar_ok else "Helvetica"
    AFB   = _ARABIC_FONT if ar_ok else "Helvetica-Bold"
    ALIGN = TA_RIGHT if is_ar else TA_LEFT

    def _xesc(s):
        # reportlab Paragraph parses an intra-paragraph mini-markup (<b>, <font>,
        # <img src=…>). Model/source text must be XML-escaped or an injected
        # <img src="http://169.254.169.254/…"> would make the PDF builder itself
        # perform a blind SSRF, and any literal "&" (e.g. "R&D") would break the
        # whole build. All app-added markup is added OUTSIDE this function.
        return str(s).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    def T(text):
        s = _ar(str(text)) if ar_ok else str(text)
        return _xesc(s)

    L = {
        "objectives": T("الأهداف التعليمية") if is_ar else "◆  LEARNING OBJECTIVES",
        "obj_bullet": "• " if is_ar else "◆  ",
        "contents":   T("المحتويات") if is_ar else "CONTENTS",
        "kw_head":    T("قاموس المصطلحات") if is_ar else "■  KEYWORDS CHEATSHEET",
        "kw_append":  [(T("قاموس المصطلحات"), ""), (T("بطاقات المراجعة"), "")] if is_ar
                      else [("KEYWORDS CHEATSHEET", ""), ("FLASH CARDS", "")],
        "fc_head":    T("بطاقات المراجعة") if is_ar else "■  FLASH CARDS",
        "sec_bullet": "" if is_ar else "■  ",
        "bul_bullet": "• " if is_ar else "▸  ",
        "q_pre":      T("سؤال: ") if is_ar else "Q:  ",
        "guide":      T("دليل الدراسة بالذكاء الاصطناعي") if is_ar else "AI Exam Study Guide",
        "luck":       T("حظ سعيد!") if is_ar else "Good luck!",
    }

    buf = io.BytesIO()
    W   = 17.4*cm

    class Doc(BaseDocTemplate):
        pass

    doc = Doc(buf, pagesize=A4,
              leftMargin=1.8*cm, rightMargin=1.8*cm,
              topMargin=1.8*cm, bottomMargin=1.5*cm)
    doc.title_str = guide.get("title", "")

    frame = Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height - 0.3*cm,
                  id='main', leftPadding=0, rightPadding=0, topPadding=0, bottomPadding=0)
    doc.addPageTemplates([PageTemplate(id='main', frames=[frame], onPage=_page_num)])

    base = getSampleStyleSheet()["Normal"]
    ST = {
        "h_title":  _st("HT",  base, fontSize=20, fontName=AFB, textColor=WHITE,        alignment=TA_CENTER, leading=26),
        "h_sub":    _st("HS",  base, fontSize=9,  fontName=AF,  textColor=colors.HexColor("#a8c4e8"), alignment=TA_CENTER),
        "obj_head": _st("OH",  base, fontSize=11, fontName=AFB, textColor=NAVY,         spaceBefore=4, spaceAfter=4, alignment=ALIGN),
        "obj_item": _st("OI",  base, fontSize=9.5,fontName=AF,  textColor=TEXT,         leftIndent=0 if is_ar else 14, rightIndent=14 if is_ar else 0, spaceAfter=3, leading=16, alignment=ALIGN),
        "toc_title":_st("TOT", base, fontSize=11, fontName=AFB, textColor=NAVY,         spaceAfter=6, alignment=ALIGN),
        "toc_item": _st("TOI", base, fontSize=9.5,fontName=AF,  textColor=TEXT,         leftIndent=0 if is_ar else 10, rightIndent=10 if is_ar else 0, spaceAfter=2, alignment=ALIGN),
        "sec_title":_st("SCT", base, fontSize=11, fontName=AFB, textColor=WHITE,        alignment=ALIGN),
        "bullet":   _st("BL",  base, fontSize=9.5,fontName=AF,  textColor=TEXT,         leftIndent=0 if is_ar else 12, rightIndent=12 if is_ar else 0, spaceAfter=3, leading=16, alignment=ALIGN),
        "tbl_hdr":  _st("TH",  base, fontSize=9,  fontName=AFB, textColor=WHITE,        alignment=ALIGN),
        "tbl_cell": _st("TC",  base, fontSize=9,  fontName=AF,  textColor=TEXT,         leading=14, alignment=ALIGN),
        "kw_term":  _st("KT",  base, fontSize=9,  fontName=AFB, textColor=NAVY_MID,     alignment=ALIGN),
        "kw_def":   _st("KD",  base, fontSize=9,  fontName=AF,  textColor=TEXT,         leading=14, alignment=ALIGN),
        "kw_head":  _st("KH",  base, fontSize=11, fontName=AFB, textColor=WHITE,        alignment=TA_CENTER),
        "fc_q":     _st("FCQ", base, fontSize=9,  fontName=AFB, textColor=WHITE,        leading=14, alignment=ALIGN),
        "fc_a":     _st("FCA", base, fontSize=9,  fontName=AF,  textColor=TEXT,         leading=14, alignment=ALIGN),
        "fc_head":  _st("FCH", base, fontSize=11, fontName=AFB, textColor=WHITE,        alignment=TA_CENTER),
        "footer":   _st("FT",  base, fontSize=8,  fontName=AF,  textColor=TEXT_LIGHT,   alignment=TA_CENTER),
    }

    elems = []

    # ── Header ────────────────────────────────────────────────────────────────
    raw_title = guide.get("title", "Study Guide")
    title    = T(raw_title.upper() if not is_ar else raw_title)
    subtitle = T(guide.get("subtitle", "Exam Study Guide"))
    hdr = Table([
        [Paragraph(title, ST["h_title"])],
        [Paragraph(f"{subtitle}  ·  {OLLAMA_MODEL}", ST["h_sub"])],
    ], colWidths=[W])
    hdr.setStyle(TableStyle([
        ("BACKGROUND",    (0,0), (-1,-1), NAVY),
        ("ALIGN",         (0,0), (-1,-1), "CENTER"),
        ("TOPPADDING",    (0,0), (0,0),   16),
        ("BOTTOMPADDING", (0,1), (0,1),   14),
        ("TOPPADDING",    (0,1), (0,1),   2),
        ("LEFTPADDING",   (0,0), (-1,-1), 12),
        ("RIGHTPADDING",  (0,0), (-1,-1), 12),
    ]))
    elems.append(hdr)
    elems.append(Spacer(1, 0.3*cm))

    # ── Learning Objectives ───────────────────────────────────────────────────
    objectives = guide.get("objectives", [])
    if objectives:
        rows = [[Paragraph(L["objectives"], ST["obj_head"])]]
        for o in objectives:
            rows.append([Paragraph(f"{L['obj_bullet']}{T(o)}", ST["obj_item"])])
        t = Table(rows, colWidths=[W])
        t.setStyle(TableStyle([
            ("BACKGROUND",    (0,0), (-1,0),  SECTION_BG),
            ("BOX",           (0,0), (-1,-1), 0.8, BORDER),
            ("LINEBELOW",     (0,0), (-1,0),  0.8, BORDER),
            ("TOPPADDING",    (0,0), (-1,-1), 5),
            ("BOTTOMPADDING", (0,0), (-1,-1), 5),
            ("LEFTPADDING",   (0,0), (-1,-1), 10),
            ("RIGHTPADDING",  (0,0), (-1,-1), 10),
        ]))
        elems.append(t)
        elems.append(Spacer(1, 0.3*cm))

    # ── Table of Contents ─────────────────────────────────────────────────────
    sections = guide.get("sections", [])
    if sections:
        toc_rows = [[Paragraph(L["contents"], ST["toc_title"])]]
        all_items = [(T(s.get("title", "")), "") for s in sections] + L["kw_append"]
        for i, (name, _) in enumerate(all_items, 1):
            dot_row = Table(
                [[Paragraph(f"{i}.  {name}", ST["toc_item"]), Paragraph("", ST["toc_item"])]],
                colWidths=[W*0.85, W*0.15]
            )
            dot_row.setStyle(TableStyle([
                ("LINEBELOW",     (0,0), (-1,-1), 0.3, BORDER),
                ("TOPPADDING",    (0,0), (-1,-1), 3),
                ("BOTTOMPADDING", (0,0), (-1,-1), 3),
                ("LEFTPADDING",   (0,0), (-1,-1), 6),
                ("RIGHTPADDING",  (0,0), (-1,-1), 6),
            ]))
            toc_rows.append([dot_row])
        toc = Table(toc_rows, colWidths=[W])
        toc.setStyle(TableStyle([
            ("BACKGROUND",    (0,0), (-1,0),  SECTION_BG),
            ("BOX",           (0,0), (-1,-1), 0.8, BORDER),
            ("TOPPADDING",    (0,0), (-1,-1), 4),
            ("BOTTOMPADDING", (0,0), (-1,-1), 4),
            ("LEFTPADDING",   (0,0), (-1,-1), 0),
        ]))
        elems.append(toc)
        elems.append(Spacer(1, 0.35*cm))

    # ── Sections ──────────────────────────────────────────────────────────────
    for idx, sec in enumerate(sections, 1):
        block = []

        sec_title_text = T(sec.get("title", ""))
        sec_hdr = Table(
            [[Paragraph(
                f"{L['sec_bullet']}{idx} · {sec_title_text if is_ar else sec_title_text.upper()}",
                ST["sec_title"]
            )]],
            colWidths=[W]
        )
        sec_hdr.setStyle(TableStyle([
            ("BACKGROUND",    (0,0), (-1,-1), NAVY_MID),
            ("TOPPADDING",    (0,0), (-1,-1), 7),
            ("BOTTOMPADDING", (0,0), (-1,-1), 7),
            ("LEFTPADDING",   (0,0), (-1,-1), 10),
            ("RIGHTPADDING",  (0,0), (-1,-1), 10),
        ]))
        block.append(sec_hdr)

        bullets = sec.get("bullets", [])
        if bullets:
            bdata = [[Paragraph(f"{L['bul_bullet']}{T(b)}", ST["bullet"])] for b in bullets]
            bt = Table(bdata, colWidths=[W])
            bt.setStyle(TableStyle([
                ("TOPPADDING",    (0,0), (-1,-1), 5),
                ("BOTTOMPADDING", (0,0), (-1,-1), 4),
                ("LEFTPADDING",   (0,0), (-1,-1), 10),
                ("RIGHTPADDING",  (0,0), (-1,-1), 10),
                ("LINEBELOW",     (0,0), (-1,-2), 0.3, colors.HexColor("#dde8ff")),
                ("BOX",           (0,0), (-1,-1), 0.5, BORDER),
            ]))
            block.append(bt)

        tbl = sec.get("table")
        if isinstance(tbl, dict) and tbl.get("headers") and tbl.get("rows"):
            headers = tbl["headers"]
            n_cols  = len(headers)
            col_w   = W / n_cols
            tbl_rows = [[Paragraph(T(h), ST["tbl_hdr"]) for h in headers]]
            for ri, row in enumerate(tbl["rows"]):
                padded = (list(row) + [""] * n_cols)[:n_cols]
                tbl_rows.append([Paragraph(T(str(c)), ST["tbl_cell"]) for c in padded])
            inner = Table(tbl_rows, colWidths=[col_w]*n_cols)
            ts = [
                ("BACKGROUND",    (0,0), (-1,0),  NAVY_LIGHT),
                ("TOPPADDING",    (0,0), (-1,-1), 5),
                ("BOTTOMPADDING", (0,0), (-1,-1), 5),
                ("LEFTPADDING",   (0,0), (-1,-1), 7),
                ("GRID",          (0,0), (-1,-1), 0.4, BORDER),
            ]
            for ri in range(1, len(tbl_rows)):
                if ri % 2 == 0:
                    ts.append(("BACKGROUND", (0,ri), (-1,ri), ROW_ALT))
            inner.setStyle(TableStyle(ts))
            block.append(Spacer(1, 0.15*cm))
            block.append(inner)

        block.append(Spacer(1, 0.3*cm))
        elems.append(KeepTogether(block))

    # ── Keywords Cheatsheet ───────────────────────────────────────────────────
    keywords = guide.get("keywords", [])
    if keywords:
        kw_hdr = Table(
            [[Paragraph(L["kw_head"], ST["kw_head"])]],
            colWidths=[W]
        )
        kw_hdr.setStyle(TableStyle([
            ("BACKGROUND",    (0,0), (-1,-1), NAVY),
            ("TOPPADDING",    (0,0), (-1,-1), 8),
            ("BOTTOMPADDING", (0,0), (-1,-1), 8),
            ("LEFTPADDING",   (0,0), (-1,-1), 10),
        ]))
        elems.append(kw_hdr)

        if is_ar:
            # Arabic: definition left, term right (visual RTL order)
            lc, rc = W*0.60, W*0.40
            kw_rows = [[
                Paragraph(T(k.get("definition", "")), ST["kw_def"]),
                Paragraph(T(k.get("term", "")),       ST["kw_term"]),
            ] for k in keywords]
        else:
            lc, rc = W*0.27, W*0.73
            kw_rows = [[
                Paragraph(_xesc(k.get("term", "")),       ST["kw_term"]),
                Paragraph(_xesc(k.get("definition", "")), ST["kw_def"]),
            ] for k in keywords]
        kw_t = Table(kw_rows, colWidths=[lc, rc])
        kts = [
            ("VALIGN",        (0,0), (-1,-1), "TOP"),
            ("TOPPADDING",    (0,0), (-1,-1), 5),
            ("BOTTOMPADDING", (0,0), (-1,-1), 5),
            ("LEFTPADDING",   (0,0), (-1,-1), 8),
            ("RIGHTPADDING",  (0,0), (-1,-1), 8),
            ("LINEBELOW",     (0,0), (-1,-2), 0.3, BORDER),
            ("BOX",           (0,0), (-1,-1), 0.5, BORDER),
        ]
        for ri in range(len(kw_rows)):
            if ri % 2 == 0:
                kts.append(("BACKGROUND", (0,ri), (-1,ri), KW_BG))
        kw_t.setStyle(TableStyle(kts))
        elems.append(kw_t)
        elems.append(Spacer(1, 0.35*cm))

    # ── Flash Cards ───────────────────────────────────────────────────────────
    flashcards = guide.get("flashcards", [])
    if flashcards:
        fc_hdr = Table(
            [[Paragraph(L["fc_head"], ST["fc_head"])]],
            colWidths=[W]
        )
        fc_hdr.setStyle(TableStyle([
            ("BACKGROUND",    (0,0), (-1,-1), NAVY),
            ("TOPPADDING",    (0,0), (-1,-1), 8),
            ("BOTTOMPADDING", (0,0), (-1,-1), 8),
            ("LEFTPADDING",   (0,0), (-1,-1), 10),
        ]))
        elems.append(fc_hdr)
        elems.append(Spacer(1, 0.2*cm))

        cw = (W - 0.3*cm) / 2
        pairs = [flashcards[i:i+2] for i in range(0, len(flashcards), 2)]
        for pair in pairs:
            row_cells = []
            for fc in pair:
                card = Table([
                    [Paragraph(f"{L['q_pre']}{T(fc.get('q',''))}", ST["fc_q"])],
                    [Paragraph(T(fc.get('a','')), ST["fc_a"])],
                ], colWidths=[cw])
                card.setStyle(TableStyle([
                    ("BACKGROUND",    (0,0), (-1,0),  CARD_Q),
                    ("BACKGROUND",    (0,1), (-1,1),  CARD_A),
                    ("TOPPADDING",    (0,0), (-1,-1), 7),
                    ("BOTTOMPADDING", (0,0), (-1,-1), 7),
                    ("LEFTPADDING",   (0,0), (-1,-1), 9),
                    ("RIGHTPADDING",  (0,0), (-1,-1), 9),
                    ("BOX",           (0,0), (-1,-1), 0.8, BORDER),
                    ("LINEBELOW",     (0,0), (-1,0),  0.5, BORDER),
                ]))
                row_cells.append(card)
            if len(row_cells) == 1:
                row_cells.append(Spacer(cw, 1))
            grid_row = Table([row_cells], colWidths=[cw, cw], hAlign='LEFT')
            grid_row.setStyle(TableStyle([("LEFTPADDING",(0,0),(-1,-1),0),("RIGHTPADDING",(0,0),(-1,-1),0)]))
            elems.append(grid_row)
            elems.append(Spacer(1, 0.2*cm))

    # ── Footer ────────────────────────────────────────────────────────────────
    elems.append(HRFlowable(width="100%", thickness=0.5, color=BORDER))
    elems.append(Spacer(1, 0.1*cm))
    elems.append(Paragraph(
        f"{T(guide.get('title',''))}  ·  {L['guide']}  ·  {OLLAMA_MODEL}  ·  {L['luck']}",
        ST["footer"]
    ))

    doc.multiBuild(elems)
    buf.seek(0)
    return buf


# ── Public config endpoint ────────────────────────────────────────────────────
@app.route("/api/config")
def api_config():
    """Return public keys the frontend needs to initialise Supabase and Stripe."""
    return jsonify({
        "supabase_url":          SUPABASE_URL,
        "supabase_anon_key":     SUPABASE_ANON_KEY,
        "stripe_publishable_key": STRIPE_PUBLISHABLE_KEY,
        "auth_enabled":          _AUTH_ENABLED,
        "anon_free_limit":       ANON_FREE_LIMIT,
        "anon_remaining":        _anon_remaining(_client_ip()),
    })


# ── Auth — current user ────────────────────────────────────────────────────────
@app.route("/api/auth/me")
def auth_me():
    uid, err = _auth_check(request)
    if err:
        return err
    if uid == "dev":
        return jsonify({"email": "dev@local", "name": "Dev", "tokens_remaining": 999,
                        "subscription_status": "active", "referral_code": "DEVLOCAL"})
    ident = _identity_from_jwt(_get_bearer(request))
    user  = _get_user(uid)
    sb    = _get_sb()
    if not user:
        # Row missing (signup trigger didn't run) — create it from the token.
        if sb:
            try:
                sb.table("users").upsert({
                    "id": uid, "email": ident["email"],
                    "name": ident["name"], "avatar_url": ident["avatar"],
                }).execute()
            except Exception as exc:
                _log.error("auth_me create-user failed: %s", exc)
        user = _get_user(uid) or {"id": uid, "email": ident["email"],
                                  "name": ident["name"], "avatar_url": ident["avatar"]}
    elif sb:
        # Backfill name/avatar/email from Google if we don't have them yet.
        patch = {}
        if ident["name"]   and not user.get("name"):       patch["name"]       = ident["name"]
        if ident["avatar"] and not user.get("avatar_url"): patch["avatar_url"] = ident["avatar"]
        if ident["email"]  and not user.get("email"):      patch["email"]      = ident["email"]
        if patch:
            try:
                sb.table("users").update(patch).eq("id", uid).execute()
                user.update(patch)
            except Exception as exc:
                _log.error("auth_me backfill failed: %s", exc)
    ref_code = user.get("referral_code") or _get_or_create_referral_code(uid)
    return jsonify({
        "id":                      user["id"],
        "email":                   user["email"],
        "name":                    user.get("name") or "",
        "avatar_url":              user.get("avatar_url") or "",
        "tokens_remaining":        user.get("tokens_remaining", 0),
        "subscription_status":     user.get("subscription_status", "free"),
        "subscription_period_end": str(user.get("subscription_period_end") or ""),
        "referral_code":           ref_code or "",
    })


# ── Referral — apply code ─────────────────────────────────────────────────────
@app.route("/api/referral/apply", methods=["POST"])
def referral_apply():
    uid, err = _auth_check(request)
    if err:
        return err
    if uid == "dev":
        return jsonify({"success": False, "reason": "dev_mode"})
    code = (request.get_json(silent=True) or {}).get("code", "").strip().upper()
    if not code:
        return jsonify({"success": False, "reason": "no_code"}), 400
    sb = _get_sb()
    if not sb:
        return jsonify({"success": False, "reason": "unavailable"})
    try:
        # Find referrer by code (cannot self-refer)
        ref = sb.table("users").select("id").eq("referral_code", code).neq("id", uid).execute()
        if not ref.data:
            return jsonify({"success": False, "reason": "invalid_code"})
        referrer_id = ref.data[0]["id"]
        # Apply only if not already referred
        sb.table("users").update({"referred_by": referrer_id}).eq("id", uid).is_("referred_by", "null").execute()
        # Remove stored code from client regardless (avoid re-tries)
        return jsonify({"success": True})
    except Exception as exc:
        _log.error(f"referral_apply error: {exc}")
        return jsonify({"success": False, "reason": "db_error"}), 500


# ── Referral — stats ───────────────────────────────────────────────────────────
@app.route("/api/referral/stats")
def referral_stats():
    uid, err = _auth_check(request)
    if err:
        return err
    if uid == "dev":
        return jsonify({"total": 0, "paid": 0, "tokens_earned": 0})
    sb = _get_sb()
    if not sb:
        return jsonify({"total": 0, "paid": 0, "tokens_earned": 0})
    try:
        rows = sb.table("users").select("referral_paid").eq("referred_by", uid).execute()
        total = len(rows.data) if rows.data else 0
        paid  = sum(1 for r in (rows.data or []) if r.get("referral_paid"))
        return jsonify({"total": total, "paid": paid, "tokens_earned": paid * 10})
    except Exception:
        return jsonify({"total": 0, "paid": 0, "tokens_earned": 0})


# ── Stripe — create checkout session ──────────────────────────────────────────
@app.route("/api/stripe/checkout", methods=["POST"])
def stripe_checkout():
    uid, err = _auth_check(request)
    if err:
        return err
    if not _stripe or not STRIPE_PRICE_ID:
        return jsonify({"error": "Payments not configured"}), 503

    user = _get_user(uid)
    if not user:
        return jsonify({"error": "User not found"}), 404

    try:
        # Reuse existing Stripe customer or create a new one
        customer_id = user.get("stripe_customer_id")
        if not customer_id:
            cust = _stripe.Customer.create(
                email=user["email"],
                name=user.get("name", ""),
                metadata={"supabase_id": uid},
            )
            customer_id = cust["id"]
            _get_sb().table("users").update(
                {"stripe_customer_id": customer_id}
            ).eq("id", uid).execute()

        session = _stripe.checkout.Session.create(
            customer=customer_id,
            payment_method_types=["card"],
            line_items=[{"price": STRIPE_PRICE_ID, "quantity": 1}],
            mode="subscription",
            # Collect address + phone: gives Stripe's fraud engine more legit
            # signals (AVS, verified contact) and cuts false-positive blocks.
            billing_address_collection="required",
            phone_number_collection={"enabled": True},
            customer_update={"address": "auto", "name": "auto"},
            success_url=f"{APP_URL}/?sub=success",
            cancel_url=f"{APP_URL}/?sub=canceled",
            metadata={"user_id": uid},
        )
        return jsonify({"url": session.url})
    except Exception as exc:
        return jsonify({"error": str(exc)}), 500


# ── Stripe — customer billing portal ──────────────────────────────────────────
@app.route("/api/stripe/portal", methods=["POST"])
def stripe_portal():
    uid, err = _auth_check(request)
    if err:
        return err
    if not _stripe:
        return jsonify({"error": "Payments not configured"}), 503

    user = _get_user(uid)
    if not user or not user.get("stripe_customer_id"):
        return jsonify({"error": "No billing account found"}), 404

    try:
        portal = _stripe.billing_portal.Session.create(
            customer=user["stripe_customer_id"],
            return_url=APP_URL,
        )
        return jsonify({"url": portal.url})
    except Exception as exc:
        return jsonify({"error": str(exc)}), 500


# ── Stripe — webhook ───────────────────────────────────────────────────────────
# Bounded in-memory idempotency guard: Stripe redelivers events (retries,
# manual resends), and our handlers overwrite token balances — replaying a
# checkout.completed would re-grant tokens the user already spent.
_processed_events       = set()
_processed_events_order = []
_proc_events_lock       = threading.Lock()

def _event_already_processed(eid):
    if not eid:
        return False
    with _proc_events_lock:
        if eid in _processed_events:
            return True
        _processed_events.add(eid)
        _processed_events_order.append(eid)
        if len(_processed_events_order) > 2000:
            _processed_events.discard(_processed_events_order.pop(0))
    return False

@app.route("/api/stripe/webhook", methods=["POST"])
def stripe_webhook():
    if not _stripe:
        return jsonify({"error": "Payments not configured"}), 503

    payload   = request.get_data()
    sig       = request.headers.get("Stripe-Signature", "")
    try:
        event = _stripe.Webhook.construct_event(payload, sig, STRIPE_WEBHOOK_SECRET)
    except ValueError:
        return jsonify({"error": "Invalid payload"}), 400
    except _stripe.error.SignatureVerificationError:
        return jsonify({"error": "Invalid signature"}), 400

    # Ack redelivered events without re-running side effects.
    if _event_already_processed(event.get("id")):
        return jsonify({"ok": True, "deduped": True})

    sb  = _get_sb()
    typ = event["type"]

    if typ == "checkout.session.completed":
        sess    = event["data"]["object"]
        user_id = (sess.get("metadata") or {}).get("user_id")
        if user_id and sb:
            sb.table("users").update({
                "subscription_status": "active",
                "tokens_remaining":    30,
                "tokens_month":        time.strftime("%Y-%m"),
            }).eq("id", user_id).execute()
            _award_referral(user_id, sb)   # reward referrer if applicable

    elif typ in ("customer.subscription.updated", "customer.subscription.deleted"):
        sub         = event["data"]["object"]
        cust_id     = sub.get("customer")
        status      = sub.get("status", "")
        # current_period_end moved from the subscription top level onto its
        # items in recent Stripe API versions — fall back to the first item.
        period_end  = sub.get("current_period_end")
        if not period_end:
            try:
                _items = (sub.get("items") or {}).get("data") or []
                period_end = _items[0].get("current_period_end") if _items else None
            except Exception:
                period_end = None
        if sb and cust_id:
            update = {"subscription_id": sub.get("id", "")}
            if status == "active":
                update["subscription_status"] = "active"
                if period_end:
                    from datetime import datetime, timezone
                    update["subscription_period_end"] = datetime.fromtimestamp(
                        period_end, tz=timezone.utc
                    ).isoformat()
            elif status in ("canceled", "unpaid", "past_due"):
                update["subscription_status"] = status
            sb.table("users").update(update).eq("stripe_customer_id", cust_id).execute()

    elif typ == "invoice.payment_succeeded":
        inv     = event["data"]["object"]
        cust_id = inv.get("customer")
        if inv.get("billing_reason") == "subscription_cycle" and sb and cust_id:
            sb.table("users").update({
                "tokens_remaining": 30,
                "tokens_month":     time.strftime("%Y-%m"),
            }).eq("stripe_customer_id", cust_id).execute()

    return jsonify({"ok": True})


# ── SSE streaming endpoint ─────────────────────────────────────────────────────

def _sse(data):
    return f"data: {json.dumps(data)}\n\n"

@app.route("/api/summarize-stream", methods=["POST"])
def summarize_stream():
    if not _check_rate_limit(_client_ip(), scope="summarize", limit=_RATE_MAX):
        return jsonify({"error": "Too many requests. Please wait a minute before trying again."}), 429
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    # ── Credit gate: signed-in users spend a token; anonymous users get a
    #    small free quota per IP so they can try without an account. ───────────
    uid = _auth_optional(request)
    anon_ip = None
    if uid:
        ok, tok_left, reason = _consume_token(uid)
        if not ok:
            return jsonify({
                "error": "You have no tokens left. Upgrade to continue.",
                "code": "no_tokens", "tokens_remaining": 0
            }), 402
    else:
        anon_ip = _client_ip()
        ok, tok_left = _anon_consume(anon_ip)
        if not ok:
            return jsonify({
                "error": "You've used your free previews. Sign up free to get more.",
                "code": "signin_for_more", "tokens_remaining": 0
            }), 402

    f            = request.files["file"]
    lang_param   = request.form.get("language", "auto")
    out_name     = _safe_name(request.form.get("filename", f.filename.rsplit(".", 1)[0]))
    detail_level = request.form.get("detail", "standard")
    dcfg         = DETAIL.get(detail_level, DETAIL["standard"])
    include_quiz = request.form.get("mode", "full") != "summary"

    _ALLOWED_EXT = (".pptx", ".ppt", ".pdf", ".docx", ".doc", ".txt")
    if not f.filename.lower().endswith(_ALLOWED_EXT):
        _refund_credit(uid, anon_ip)
        return jsonify({"error": "Unsupported file type. Supported: .pptx, .ppt, .pdf, .docx, .doc, .txt"}), 400

    if not ollama_running():
        _refund_credit(uid, anon_ip)
        return jsonify({"error": "AI service is not configured. Set GROQ_API_KEY."}), 503

    file_bytes = f.read()

    def generate():
        try:
            yield _sse({"step": "extract", "msg": "Extracting content…"})
            slides = extract_slides(io.BytesIO(file_bytes), filename=f.filename)
            if not any(s["content"] or s["title"] for s in slides):
                _refund_credit(uid, anon_ip)
                yield _sse({"error": "No readable content in this file"}); return

            total = len([s for s in slides if s["content"].strip()])

            # Auto-detect language from slide content
            language = lang_param if lang_param in ("ar", "en") else _detect_language(slides)
            lang_label = "Arabic" if language == "ar" else "English"
            yield _sse({"step": "extract", "msg": f"Found {total} content slides ({lang_label}) — analysing…", "language": language})

            # Pass 1
            yield _sse({"step": "overview", "msg": "Analysing structure and keywords…"})
            overview = pass1_overview(
                [s for s in slides if s["content"].strip() or s["title"].strip()],
                language, dcfg
            )

            # Normalize: mistral sometimes returns sections/keywords as plain strings
            raw_sections = overview.get("sections", [])
            overview["sections"] = [
                s if isinstance(s, dict) else {"title": str(s), "slide_nums": []}
                for s in (raw_sections if isinstance(raw_sections, list) else [])
            ]
            raw_keywords = overview.get("keywords", [])
            overview["keywords"] = [
                k if isinstance(k, dict) else {"term": str(k), "definition": ""}
                for k in (raw_keywords if isinstance(raw_keywords, list) else [])
            ]

            content_slides = [s for s in slides if s["content"].strip()]

            # Fallback: if no sections produced, wrap all content into one
            if not overview["sections"]:
                overview["sections"] = [{
                    "title": overview.get("title", "Content Overview"),
                    "slide_nums": [s["slide_num"] for s in content_slides]
                }]

            sections = overview["sections"]

            for evt in _sections_parallel(sections, content_slides, language, dcfg):
                yield _sse(evt)

            for evt in _flashcards_mcq_parallel(overview, language, dcfg, include_quiz):
                yield _sse(evt)

            # Build PDF + Markdown
            overview["language"] = language   # so exports + the in-app viewer localise
            yield _sse({"step": "pdf", "msg": "Building PDF & Markdown…"})
            pdf_buf   = build_pdf(overview, language, out_name)
            pdf_bytes = pdf_buf.read()
            md_text   = build_markdown(overview)

            job_id = uuid.uuid4().hex
            store_job(job_id, pdf_bytes, md_text, overview, slides, f"{out_name}_study_guide.pdf")

            yield _sse({"step": "done", "job_id": job_id,
                        "tokens_remaining": tok_left,
                        "sections":   len(sections),
                        "keywords":   len(overview.get("keywords",   [])),
                        "flashcards": len(overview.get("flashcards", [])),
                        "mcqs":       len(overview.get("mcqs",       []))})

        except Exception as e:
            _log.error("GENERATE_ERROR: %s\n%s", e, _tb.format_exc())
            _refund_credit(uid, anon_ip)
            yield _sse({"error": _safe_err(e)})

    return Response(
        stream_with_context(generate()),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"}
    )


@app.route("/api/download/<job_id>")
def download_job(job_id):
    if not _valid_job(job_id):
        return jsonify({"error": "Invalid job ID"}), 400
    fmt      = request.args.get("format", "pdf")
    filename = _safe_name(request.args.get("filename", "study_guide"))
    job = get_job(job_id)
    if not job:
        return jsonify({"error": "File not found or expired"}), 404
    if fmt == "md":
        content = (job.get("md") or "").encode("utf-8")
        return send_file(io.BytesIO(content), mimetype="text/markdown",
                         as_attachment=True, download_name=f"{filename}.md")
    return send_file(io.BytesIO(job["pdf"]), mimetype="application/pdf",
                     as_attachment=True, download_name=job.get("filename", f"{filename}.pdf"))


@app.route("/api/delete/<job_id>", methods=["POST"])
def delete_job(job_id):
    """Immediately purge a job from memory (user-requested delete-now)."""
    if not _valid_job(job_id):
        return jsonify({"error": "Invalid job ID"}), 400
    with _jobs_lock:
        existed = _jobs.pop(job_id, None) is not None
    return jsonify({"ok": True, "deleted": existed})


@app.route("/api/guide/<job_id>")
def get_guide(job_id):
    if not _valid_job(job_id):
        return jsonify({"error": "Invalid job ID"}), 400
    with _jobs_lock:
        job = get_job(job_id)
    if not job:
        return jsonify({"error": "Session expired — re-upload the file"}), 404
    guide = job.get("guide", {})
    return jsonify({
        "title":      guide.get("title", ""),
        "subtitle":   guide.get("subtitle", ""),
        "sections":   guide.get("sections",   []),
        "flashcards": guide.get("flashcards", []),
        "mcqs":       guide.get("mcqs", []),
        "keywords":   guide.get("keywords",   []),
        "objectives": guide.get("objectives", []),
        "language":   guide.get("language", "en"),
    })


@app.route("/api/chat/<job_id>", methods=["POST"])
def chat_with_slides(job_id):
    if not _check_rate_limit(_client_ip(), scope="chat", limit=20):
        return jsonify({"error": "Too many requests. Please wait a minute."}), 429
    uid, err = _auth_check(request)
    if err:
        return err
    if not _valid_job(job_id):
        return jsonify({"error": "Invalid job ID"}), 400
    data     = request.json or {}
    question = data.get("question", "")[:500].strip()
    language = "ar" if data.get("language") == "ar" else "en"
    if not question:
        return jsonify({"error": "No question provided"}), 400
    with _jobs_lock:
        job = get_job(job_id)
    if not job:
        return jsonify({"error": "Session expired — re-upload the file"}), 404

    # Build rich context from guide (sections + keywords) rather than raw slide chunks
    guide   = job.get("guide") or {}
    context_parts = []
    if guide.get("title"):
        context_parts.append(f"Title: {guide['title']}")
    if guide.get("objectives"):
        context_parts.append("Objectives:\n" + "\n".join(f"- {o}" for o in guide["objectives"]))
    for sec in (guide.get("sections") or []):
        bullets = sec.get("bullets") or []
        if bullets:
            context_parts.append(f"\n[{sec.get('title','')}]\n" + "\n".join(
                f"- {b}" if isinstance(b, str) else f"- {b.get('text') or b.get('fact','')}"
                for b in bullets
            ))
    if guide.get("keywords"):
        kw_lines = [
            f"{k['term']}: {k.get('definition','')}" if isinstance(k, dict) else str(k)
            for k in guide["keywords"][:30]
        ]
        context_parts.append("Key terms:\n" + "\n".join(kw_lines))
    context = "\n\n".join(context_parts) or "No material available."

    lang = "in Arabic" if language == "ar" else "in English"
    try:
        result = _call_ollama(
            f"""Answer this question {lang} using ONLY the study material below.

Material:
{context}

Question: {question}

Return JSON: {{"answer": "your detailed answer"}}

Rules:
- Answer directly and specifically from the material
- If genuinely not covered, say so briefly
- Be helpful and detailed; include facts, definitions, examples from the material
- JSON only""", num_predict=1024)
        return jsonify({"answer": result.get("answer", "No answer found in the material.")})
    except Exception as e:
        return jsonify({"error": _safe_err(e)}), 500


@app.route("/api/download-zip", methods=["POST"])
def download_zip():
    import zipfile as zf
    job_ids = [jid for jid in (request.json or {}).get("job_ids", []) if _valid_job(jid)]
    buf = io.BytesIO()
    with zf.ZipFile(buf, "w", zf.ZIP_DEFLATED) as z:
        with _jobs_lock:
            for jid in job_ids:
                job = _jobs.get(jid)
                if job and job.get("pdf"):
                    z.writestr(job["filename"], job["pdf"])
    buf.seek(0)
    return send_file(buf, mimetype="application/zip", as_attachment=True,
                     download_name="study_guides.zip")


@app.route("/api/view/md/<job_id>")
def view_md(job_id):
    if not _valid_job(job_id):
        return "<h2 style='font-family:sans-serif;padding:2rem'>Invalid job ID</h2>", 400
    with _jobs_lock:
        job = get_job(job_id)
    if not job:
        return "<h2 style='font-family:sans-serif;padding:2rem'>Guide not found or expired (10 min TTL)</h2>", 404
    title = _he(job["guide"].get("title", "Study Guide"))
    md = job["md"]

    def _md_to_html(text):
        lines, out = text.split('\n'), []
        i = 0
        while i < len(lines):
            l = lines[i]
            if l.startswith('# '):
                out.append(f'<h1>{_inline(_he(l[2:]))}</h1>')
            elif l.startswith('## '):
                out.append(f'<h2>{_inline(_he(l[3:]))}</h2>')
            elif l.startswith('### '):
                out.append(f'<h3>{_inline(_he(l[4:]))}</h3>')
            elif l.startswith('- ') or l.startswith('* '):
                out.append(f'<li>{_inline(_he(l[2:]))}</li>')
            elif l.startswith('| ') and '|' in l[2:]:
                rows, align = [], l
                while i < len(lines) and lines[i].startswith('|'):
                    if not re.match(r'^\|[-| :]+\|$', lines[i]):
                        cells = [c.strip() for c in lines[i].strip('|').split('|')]
                        tag = 'th' if rows == [] else 'td'
                        out.append('<tr>' + ''.join(f'<{tag}>{_inline(_he(c))}</{tag}>' for c in cells) + '</tr>')
                    i += 1
                i -= 1
            elif l.strip() == '':
                out.append('<br>')
            else:
                out.append(f'<p>{_inline(_he(l))}</p>')
            i += 1
        return '\n'.join(out)

    def _inline(t):
        t = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', t)
        t = re.sub(r'\*(.+?)\*',     r'<em>\1</em>',         t)
        t = re.sub(r'`(.+?)`',       r'<code>\1</code>',     t)
        return t

    body = _md_to_html(md)
    return f"""<!DOCTYPE html><html><head><meta charset="UTF-8">
<title>{title}</title>
<style>
  body{{font-family:'Segoe UI',system-ui,sans-serif;max-width:820px;margin:0 auto;padding:2rem;
       background:#f8faff;color:#0a1628;line-height:1.7}}
  h1{{font-size:1.7rem;color:#1a3a6e;border-bottom:2px solid #c5d8ff;padding-bottom:.5rem}}
  h2{{font-size:1.2rem;color:#2e5ca8;margin-top:1.8rem;border-left:4px solid #4f8ef7;padding-left:.75rem}}
  h3{{font-size:1rem;color:#1a3a6e}}
  li{{margin-bottom:.35rem}}
  table{{border-collapse:collapse;width:100%;margin:1rem 0}}
  th{{background:#1a3a6e;color:#fff;padding:8px 12px;text-align:left}}
  td{{padding:7px 12px;border-bottom:1px solid #dde8ff}}
  tr:nth-child(even) td{{background:#f0f5ff}}
  code{{background:#e8f0ff;padding:1px 5px;border-radius:4px;font-size:.9em}}
  strong{{color:#1a3a6e}}
  @media print{{body{{background:#fff}}}}
</style></head><body>
{body}
<hr style="margin-top:2rem;border-color:#c5d8ff">
<p style="font-size:.8rem;color:#8aa0c8;text-align:center">Generated by Alimne (علّمني) · {OLLAMA_MODEL}</p>
</body></html>"""


def _page_shell(title, body_css, body_html, script=""):
    return f"""<!DOCTYPE html><html><head><meta charset="UTF-8">
<title>{title}</title>
<style>
  *{{box-sizing:border-box;margin:0;padding:0}}
  body{{font-family:'Segoe UI',system-ui,sans-serif;background:#050d1a;color:#e8f0ff;min-height:100vh}}
  .top{{background:#0a1628;border-bottom:1px solid #1a3a6e;padding:1rem 2rem;
        display:flex;align-items:center;justify-content:space-between}}
  .top h1{{font-size:1rem;font-weight:700;color:#e8f0ff}}
  .badge{{background:#4f8ef7;color:#fff;padding:2px 10px;border-radius:20px;font-size:12px}}
  .wrap{{max-width:760px;margin:0 auto;padding:2rem 1.25rem}}
  {body_css}
</style></head><body>
<div class="top"><h1>📖 {title}</h1></div>
<div class="wrap">{body_html}</div>
<script>{script}</script>
</body></html>"""


@app.route("/api/view/cards/<job_id>")
def view_cards(job_id):
    if not _valid_job(job_id):
        return "<h2 style='font-family:sans-serif;padding:2rem'>Invalid job ID</h2>", 400
    with _jobs_lock:
        job = get_job(job_id)
    if not job:
        return "<h2 style='font-family:sans-serif;padding:2rem'>Guide not found or expired</h2>", 404
    guide = job["guide"]
    title = _he(guide.get("title", "Flash Cards"))
    cards = [f for f in guide.get("flashcards", []) if isinstance(f, dict)]
    if not cards:
        return _page_shell(title, "", "<p style='text-align:center;color:#4a5f80;padding:3rem'>No flash cards available.</p>")

    # Escape "<" so model-derived content can't break out of the <script> block
    # (json.dumps does NOT escape "</script>"). < is valid JSON and JS.
    cards_json = json.dumps(cards).replace("<", "\\u003c")
    css = """
  .fc-counter{text-align:center;color:#8aa0c8;font-size:.88rem;margin-bottom:1.2rem}
  .card{perspective:900px;height:220px;cursor:pointer;margin-bottom:1.5rem}
  .inner{position:relative;width:100%;height:100%;
         transition:transform .45s cubic-bezier(.4,0,.2,1);transform-style:preserve-3d}
  .card.flipped .inner{transform:rotateY(180deg)}
  .front,.back{position:absolute;inset:0;border-radius:14px;padding:1.5rem;
               backface-visibility:hidden;display:flex;flex-direction:column;justify-content:center}
  .front{background:linear-gradient(135deg,#1a3a6e,#2e5ca8);border:1px solid rgba(79,142,247,.3)}
  .back{background:rgba(255,255,255,.05);border:1px solid rgba(79,142,247,.2);
        transform:rotateY(180deg)}
  .front .q{font-size:1rem;font-weight:600;color:#e8f0ff;text-align:center}
  .front .hint{font-size:.75rem;color:#8aa0c8;margin-top:.75rem;text-align:center}
  .back .a{font-size:.95rem;color:#e8f0ff;line-height:1.6}
  .nav{display:flex;gap:.75rem;justify-content:center;margin-top:.5rem}
  .btn{padding:.55rem 1.4rem;border-radius:9px;border:1px solid rgba(79,142,247,.3);
       background:rgba(79,142,247,.1);color:#4f8ef7;font-size:.88rem;font-weight:600;
       cursor:pointer;font-family:inherit;transition:all .2s}
  .btn:hover{background:rgba(79,142,247,.2)}
  .btn.correct{background:rgba(34,197,94,.12);border-color:rgba(34,197,94,.4);color:#22c55e}
  .btn.wrong{background:rgba(239,68,68,.1);border-color:rgba(239,68,68,.35);color:#ef4444}
  .progress{height:4px;background:rgba(79,142,247,.15);border-radius:99px;margin-bottom:1.5rem;overflow:hidden}
  .progress-bar{height:100%;background:linear-gradient(90deg,#4f8ef7,#a78bfa);
                border-radius:99px;transition:width .4s}
  .done{text-align:center;padding:2rem;color:#22c55e;font-size:1.1rem;font-weight:600}
"""
    html = """<div class="fc-counter" id="ctr"></div>
<div class="progress"><div class="progress-bar" id="pb"></div></div>
<div id="cardWrap"></div>
<div class="nav">
  <button class="btn wrong" onclick="mark(false)">✗ Don't know</button>
  <button class="btn" onclick="flip()">Flip</button>
  <button class="btn correct" onclick="mark(true)">✓ Know it</button>
</div>"""
    script = f"""
const cards = {cards_json};
const esc = s => String(s ?? '').replace(/[&<>"']/g, m => ({{'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}}[m]));
let idx = 0, known = 0;
function render() {{
  if (idx >= cards.length) {{
    document.getElementById('cardWrap').innerHTML = '<div class="done">🎉 Done! ' + known + '/' + cards.length + ' known</div>';
    document.querySelector('.nav').style.display = 'none';
    document.getElementById('ctr').textContent = 'Complete';
    return;
  }}
  const c = cards[idx];
  document.getElementById('ctr').textContent = (idx+1) + ' / ' + cards.length;
  document.getElementById('pb').style.width = (idx/cards.length*100) + '%';
  document.getElementById('cardWrap').innerHTML = `
    <div class="card" id="card" onclick="flip()">
      <div class="inner">
        <div class="front"><div class="q">${{esc(c.q)}}</div><div class="hint">Click to reveal answer</div></div>
        <div class="back"><div class="a">${{esc(c.a)}}</div></div>
      </div>
    </div>`;
}}
function flip() {{ document.getElementById('card').classList.toggle('flipped'); }}
function mark(k) {{ if(k) known++; idx++; render(); }}
render();
"""
    return _page_shell(f"Flash Cards — {title}", css, html, script)


@app.route("/api/view/quiz/<job_id>")
def view_quiz(job_id):
    if not _valid_job(job_id):
        return "<h2 style='font-family:sans-serif;padding:2rem'>Invalid job ID</h2>", 400
    with _jobs_lock:
        job = get_job(job_id)
    if not job:
        return "<h2 style='font-family:sans-serif;padding:2rem'>Guide not found or expired</h2>", 404
    guide = job["guide"]
    title = _he(guide.get("title", "Quiz"))
    mcqs = [m for m in guide.get("mcqs", []) if isinstance(m, dict)]
    if not mcqs:
        return _page_shell(title, "", "<p style='text-align:center;color:#4a5f80;padding:3rem'>No quiz questions available.</p>")

    # Escape "<" so model-derived content can't break out of the <script> block.
    mcqs_json = json.dumps(mcqs).replace("<", "\\u003c")
    css = """
  .q-num{color:#8aa0c8;font-size:.82rem;margin-bottom:.4rem}
  .q-text{font-size:1rem;font-weight:600;color:#e8f0ff;margin-bottom:1rem;line-height:1.5}
  .opt{display:flex;align-items:center;gap:.6rem;width:100%;text-align:left;
       padding:.6rem .9rem;border-radius:9px;border:1px solid rgba(79,142,247,.2);
       background:rgba(255,255,255,.04);color:#8aa0c8;font-size:.88rem;cursor:pointer;
       font-family:inherit;margin-bottom:.45rem;transition:all .2s}
  .opt:hover:not(:disabled){border-color:#4f8ef7;color:#e8f0ff}
  .opt.correct{border-color:#22c55e;color:#22c55e;background:rgba(34,197,94,.1)}
  .opt.wrong{border-color:#ef4444;color:#ef4444;background:rgba(239,68,68,.08)}
  .explanation{margin-top:.75rem;padding:.75rem;border-radius:9px;
               background:rgba(79,142,247,.08);border:1px solid rgba(79,142,247,.2);
               font-size:.85rem;color:#8aa0c8;display:none}
  .nav{margin-top:1.25rem;display:flex;justify-content:flex-end}
  .next-btn{padding:.55rem 1.4rem;border-radius:9px;border:none;
            background:linear-gradient(135deg,#4f8ef7,#2e5ca8);color:#fff;
            font-size:.88rem;font-weight:600;cursor:pointer;font-family:inherit;display:none}
  .score-box{text-align:center;padding:2.5rem;background:rgba(79,142,247,.08);
             border:1px solid rgba(79,142,247,.2);border-radius:16px}
  .score-big{font-size:3rem;font-weight:700;color:#4f8ef7}
  .progress{height:4px;background:rgba(79,142,247,.15);border-radius:99px;margin-bottom:1.5rem;overflow:hidden}
  .progress-bar{height:100%;background:linear-gradient(90deg,#4f8ef7,#a78bfa);
                border-radius:99px;transition:width .4s}
"""
    html = """<div class="progress"><div class="progress-bar" id="pb"></div></div>
<div id="qWrap"></div>"""
    script = f"""
const qs = {mcqs_json};
const esc = s => String(s ?? '').replace(/[&<>"']/g, m => ({{'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}}[m]));
let idx=0, score=0, answered=false;
function render() {{
  if(idx>=qs.length){{
    document.getElementById('qWrap').innerHTML=`<div class="score-box">
      <div class="score-big">${{score}}/${{qs.length}}</div>
      <div style="color:#8aa0c8;margin-top:.5rem">Quiz complete!</div>
    </div>`;
    document.getElementById('pb').style.width='100%';
    return;
  }}
  const q=qs[idx]; answered=false;
  document.getElementById('pb').style.width=(idx/qs.length*100)+'%';
  const opts=q.options.map((o,i)=>`<button class="opt" id="o${{i}}" data-l="${{esc(o[0])}}" onclick="pick(this.dataset.l,this)">${{esc(o)}}</button>`).join('');
  document.getElementById('qWrap').innerHTML=`
    <div class="q-num">Question ${{idx+1}} of ${{qs.length}}</div>
    <div class="q-text">${{esc(q.q)}}</div>
    ${{opts}}
    <div class="explanation" id="exp">${{esc(q.explanation||'')}}</div>
    <div class="nav"><button class="next-btn" id="nxt" onclick="next()">Next →</button></div>`;
}}
function pick(letter,btn) {{
  if(answered) return; answered=true;
  const q=qs[idx];
  document.querySelectorAll('.opt').forEach(b=>b.disabled=true);
  if(letter===q.answer){{ btn.classList.add('correct'); score++; }}
  else {{ btn.classList.add('wrong');
    document.querySelectorAll('.opt').forEach(b=>{{if(b.textContent.startsWith(q.answer)) b.classList.add('correct');}});
  }}
  document.getElementById('exp').style.display='block';
  document.getElementById('nxt').style.display='inline-block';
}}
function next(){{ idx++; render(); }}
render();
"""
    return _page_shell(f"Quiz — {title}", css, html, script)


@app.route("/api/status")
def status():
    running = ollama_running()
    models  = ollama_models() if running else []
    active  = next((m for m in models if OLLAMA_MODEL in m), models[0] if models else None)
    return jsonify({"ollama": running, "model": active, "models": models})


# ── YouTube transcript endpoint ───────────────────────────────────────────────

def _parse_caption_xml(xml_text):
    """Parse YouTube caption XML and return joined transcript string."""
    import xml.etree.ElementTree as ET
    try:
        root = ET.fromstring(xml_text)
    except Exception:
        return None
    texts = []
    for elem in root.iter("text"):
        t = (elem.text or "").strip()
        if t:
            t = t.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">").replace("&#39;", "'").replace("&quot;", '"')
            texts.append(t)
    return " ".join(texts) if texts else None


def _fetch_captions(video_id):
    """Try multiple methods to get transcript from YouTube without downloading audio."""
    html = None

    # Method 1 — captionTracks from page source
    try:
        r = http.get(f"https://www.youtube.com/watch?v={video_id}", timeout=15,
                     headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"})
        r.raise_for_status()
        html = r.text
    except Exception:
        pass

    if html:
        m = re.search(r'"captionTracks":\s*(\[.*?\])', html)
        if m:
            try:
                tracks = json.loads(m.group(1))
                # Prefer English ASR, then English, then any
                def track_priority(t):
                    lc = t.get("languageCode", "")
                    kind = t.get("kind", "")
                    if lc.startswith("en") and kind == "asr": return 0
                    if lc.startswith("en"): return 1
                    return 2
                tracks.sort(key=track_priority)
                for track in tracks:
                    track_url = track.get("baseUrl")
                    if not track_url:
                        continue
                    try:
                        tr = http.get(track_url, timeout=15)
                        result = _parse_caption_xml(tr.text)
                        if result:
                            return result
                    except Exception:
                        continue
            except Exception:
                pass

    # Method 2 — timedtext API (works for many videos without captionTracks in HTML)
    for kind in ("asr", ""):
        for lang in ("en", "en-US", "en-GB"):
            params = {"v": video_id, "lang": lang, "fmt": "srv3"}
            if kind:
                params["kind"] = kind
            try:
                tr = http.get("https://www.youtube.com/api/timedtext", params=params, timeout=10)
                if tr.status_code == 200 and tr.text.strip():
                    result = _parse_caption_xml(tr.text)
                    if result:
                        return result
            except Exception:
                continue

    raise ValueError("no_captions")


def _transcribe_with_whisper(video_id):
    """Download audio to a private temp dir, transcribe with Groq Whisper, delete
    immediately. A per-request mkdtemp avoids predictable /tmp names and the race
    where two concurrent requests for the same video clobber each other's files."""
    import glob, tempfile, shutil
    try:
        import yt_dlp
    except ImportError:
        raise ValueError("yt-dlp not installed — cannot transcribe audio.")

    workdir = tempfile.mkdtemp(prefix="yt_")
    prefix  = os.path.join(workdir, "audio")
    try:
        ydl_opts = {
            # Prefer smallest audio: opus<96k > m4a < 96k > any audio
            "format": "bestaudio[abr<=96][ext=webm]/bestaudio[abr<=96][ext=m4a]/bestaudio[ext=m4a]/bestaudio[ext=webm]/bestaudio",
            "outtmpl": prefix + ".%(ext)s",
            "quiet": True,
            "no_warnings": True,
            "noplaylist": True,
            "socket_timeout": 30,
            "max_filesize": 22 * 1024 * 1024,  # abort mid-download instead of filling disk
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([f"https://www.youtube.com/watch?v={video_id}"])

        files = glob.glob(prefix + ".*")
        if not files:
            raise ValueError("Audio download produced no file.")
        audio_path = files[0]

        size = os.path.getsize(audio_path)
        if size > 20 * 1024 * 1024:
            raise ValueError("Video audio exceeds 20 MB — try a shorter video (under ~15 minutes).")

        ext = os.path.splitext(audio_path)[1].lstrip(".")
        mime = "audio/mp4" if ext in ("m4a", "mp4") else "audio/webm"

        with open(audio_path, "rb") as af:
            r = http.post(
                "https://api.groq.com/openai/v1/audio/transcriptions",
                headers={"Authorization": f"Bearer {GROQ_API_KEY}"},
                files={"file": (f"audio.{ext}", af, mime)},
                data={"model": "whisper-large-v3-turbo", "response_format": "text"},
                timeout=300,
            )
        if r.status_code == 413:
            raise ValueError("Audio file too large for Whisper — try a shorter video (under ~15 minutes).")
        r.raise_for_status()
        return r.text.strip()

    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def _fetch_youtube_transcript(video_id):
    """Try captions first; fall back to Whisper audio transcription."""
    try:
        return _fetch_captions(video_id)
    except ValueError as e:
        if str(e) != "no_captions":
            raise
    if not GROQ_API_KEY:
        raise ValueError("No captions found and GROQ_API_KEY not set for Whisper fallback.")
    return _transcribe_with_whisper(video_id)


def _text_to_slides(text, chunk_size=500):
    """Split plain text into slide-like dicts."""
    words = text.split()
    chunks = []
    current = []
    current_len = 0
    for w in words:
        current.append(w)
        current_len += len(w) + 1
        if current_len >= chunk_size:
            chunks.append(" ".join(current))
            current = []
            current_len = 0
    if current:
        chunks.append(" ".join(current))
    return [
        {"slide_num": i + 1, "title": f"Segment {i + 1}", "content": c}
        for i, c in enumerate(chunks)
    ]


def _extract_video_id(url):
    """Extract YouTube video ID — only accepts youtube.com and youtu.be hostnames."""
    import urllib.parse
    parsed = urllib.parse.urlparse(url)
    host = (parsed.hostname or "").lower()
    host = host[4:] if host.startswith("www.") else host  # strip prefix, not charset
    if host not in ("youtube.com", "youtu.be", "m.youtube.com"):
        return None
    patterns = [
        r'(?:v=|youtu\.be/|/embed/|/shorts/)([a-zA-Z0-9_-]{11})',
    ]
    for p in patterns:
        m = re.search(p, url)
        if m:
            return m.group(1)
    return None


def _stream_text_as_sse(text, language, out_name, job_source, dcfg=None, tok_left=None, include_quiz=True, uid=None, anon_ip=None):
    """Shared SSE generator for YouTube/text endpoints."""
    _dcfg = dcfg or DETAIL["standard"]
    def generate():
        try:
            yield _sse({"step": "extract", "msg": "Preparing content…"})
            slides = _text_to_slides(text)
            if not slides:
                _refund_credit(uid, anon_ip)
                yield _sse({"error": "No content extracted"}); return
            yield _sse({"step": "extract", "msg": f"Split into {len(slides)} segments — analysing…"})

            dcfg = _dcfg

            yield _sse({"step": "overview", "msg": "Analysing structure and keywords…"})
            overview = pass1_overview(slides, language, dcfg)

            raw_sections = overview.get("sections", [])
            overview["sections"] = [
                s if isinstance(s, dict) else {"title": str(s), "slide_nums": []}
                for s in (raw_sections if isinstance(raw_sections, list) else [])
            ]
            raw_keywords = overview.get("keywords", [])
            overview["keywords"] = [
                k if isinstance(k, dict) else {"term": str(k), "definition": ""}
                for k in (raw_keywords if isinstance(raw_keywords, list) else [])
            ]
            if not overview["sections"]:
                overview["sections"] = [{
                    "title": overview.get("title", "Content Overview"),
                    "slide_nums": [s["slide_num"] for s in slides]
                }]

            sections = overview["sections"]
            for evt in _sections_parallel(sections, slides, language, dcfg):
                yield _sse(evt)

            for evt in _flashcards_mcq_parallel(overview, language, dcfg, include_quiz):
                yield _sse(evt)

            overview["language"] = language   # so exports + the in-app viewer localise
            yield _sse({"step": "pdf", "msg": "Building PDF & Markdown…"})
            pdf_buf = build_pdf(overview, language, out_name)
            pdf_bytes = pdf_buf.read()
            md_text = build_markdown(overview)

            job_id = uuid.uuid4().hex
            store_job(job_id, pdf_bytes, md_text, overview, slides, f"{out_name}_study_guide.pdf")

            done_data = {"step": "done", "job_id": job_id,
                         "sections":   len(sections),
                         "keywords":   len(overview.get("keywords",   [])),
                         "flashcards": len(overview.get("flashcards", [])),
                         "mcqs":       len(overview.get("mcqs",       []))}
            if tok_left is not None:
                done_data["tokens_remaining"] = tok_left
            yield _sse(done_data)

        except Exception as e:
            _log.error("STREAM_TEXT_ERROR: %s\n%s", e, _tb.format_exc())
            _refund_credit(uid, anon_ip)
            yield _sse({"error": _safe_err(e)})
    return generate


def _yt_duration(video_id):
    """Return video duration in seconds (or None if unknown)."""
    try:
        import yt_dlp
        opts = {"quiet": True, "no_warnings": True, "skip_download": True, "socket_timeout": 20}
        with yt_dlp.YoutubeDL(opts) as ydl:
            info = ydl.extract_info(f"https://www.youtube.com/watch?v={video_id}", download=False)
        return info.get("duration")
    except Exception:
        return None

@app.route("/api/youtube", methods=["POST"])
def youtube_transcript():
    # Rate-limit BEFORE the expensive yt-dlp scrape / Whisper path (this endpoint
    # had none, so anon callers could force unbounded yt-dlp + paid transcription).
    if not _check_rate_limit(_client_ip(), scope="youtube", limit=_RATE_MAX):
        return jsonify({"error": "Too many requests — please wait a moment and try again."}), 429
    uid = _auth_optional(request)
    anon_ip = None

    data = request.get_json(silent=True) or {}
    url = (data.get("url") or "").strip()
    lang_param = data.get("language", "auto")
    detail_level = data.get("detail", "standard")
    yt_dcfg = DETAIL.get(detail_level, DETAIL["standard"])
    include_quiz = data.get("mode", "full") != "summary"
    if not url:
        return jsonify({"error": "No URL provided"}), 400
    if not ollama_running():
        return jsonify({"error": "AI service is not configured. Set GROQ_API_KEY."}), 503

    video_id = _extract_video_id(url)
    if not video_id:
        return jsonify({"error": "Could not extract video ID from URL"}), 400

    # Reject videos longer than 50 minutes (checked before consuming a token)
    _dur = _yt_duration(video_id)
    if _dur and _dur > 3000:
        return jsonify({"error": "This video is too long. Maximum supported length is 50 minutes."}), 400

    if uid:
        ok, tok_left, reason = _consume_token(uid)
        if not ok:
            return jsonify({
                "error": "You have no tokens left. Upgrade to continue.",
                "code": "no_tokens", "tokens_remaining": 0
            }), 402
    else:
        anon_ip = _client_ip()
        ok, tok_left = _anon_consume(anon_ip)
        if not ok:
            return jsonify({
                "error": "You've used your free previews. Sign up free to get more.",
                "code": "signin_for_more", "tokens_remaining": 0
            }), 402

    out_name = _safe_name(f"youtube_{video_id}")

    def generate():
        try:
            yield _sse({"step": "transcript", "msg": "Looking for captions…"})
            try:
                transcript_text = _fetch_captions(video_id)
                yield _sse({"step": "transcript", "msg": "Captions found — processing…"})
            except ValueError as e:
                if str(e) != "no_captions":
                    _refund_credit(uid, anon_ip)
                    yield _sse({"error": _safe_err(e)}); return
                if not GROQ_API_KEY:
                    _refund_credit(uid, anon_ip)
                    yield _sse({"error": "No captions found and GROQ_API_KEY not set."}); return
                yield _sse({"step": "transcript", "msg": "No captions — downloading audio for Whisper transcription…"})
                try:
                    transcript_text = _transcribe_with_whisper(video_id)
                    yield _sse({"step": "transcript", "msg": "Audio transcribed — processing…"})
                except ValueError as we:
                    _refund_credit(uid, anon_ip)
                    yield _sse({"error": str(we)}); return

            language = lang_param if lang_param in ("ar", "en") else _detect_language(transcript_text)
            lang_label = "Arabic" if language == "ar" else "English"
            yield _sse({"step": "transcript", "msg": f"Transcript ready ({lang_label}) — building study guide…", "language": language})
            for event in _stream_text_as_sse(transcript_text, language, out_name, "youtube", yt_dcfg, tok_left, include_quiz, uid=uid, anon_ip=anon_ip)():
                yield event
        except Exception as ex:
            _log.error("youtube SSE error: %s", ex)
            _refund_credit(uid, anon_ip)
            yield _sse({"error": str(ex)})

    return Response(
        stream_with_context(generate()),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"}
    )


# ── URL / pasted text endpoint ─────────────────────────────────────────────────

class _PinnedIPAdapter(http.adapters.HTTPAdapter):
    """Pin the socket connection to a pre-validated IP so a DNS rebind can't
    swap in an internal address between our SSRF check and the actual request,
    while preserving TLS SNI + certificate hostname verification for the host."""
    def __init__(self, host, pinned_ip, is_https, *args, **kwargs):
        self._host      = host
        self._pinned_ip = pinned_ip
        self._is_https  = is_https
        super().__init__(*args, **kwargs)

    def init_poolmanager(self, *args, **kwargs):
        if self._is_https:
            kwargs["server_hostname"] = self._host   # SNI = real host
            kwargs["assert_hostname"] = self._host    # verify cert against real host
        return super().init_poolmanager(*args, **kwargs)

    def send(self, request, **kwargs):
        from urllib.parse import urlparse, urlunparse
        parsed = urlparse(request.url)
        request.headers["Host"] = parsed.netloc      # keep the original Host header
        ip = f"[{self._pinned_ip}]" if ":" in self._pinned_ip else self._pinned_ip
        netloc = f"{ip}:{parsed.port}" if parsed.port else ip
        request.url = urlunparse(parsed._replace(netloc=netloc))
        return super().send(request, **kwargs)


def _fetch_url_text(url):
    """Fetch a PUBLIC webpage and extract readable text. SSRF guard: public
    http(s) only, no private/internal addresses, no redirects, 5 MB cap. The
    connection is pinned to the vetted IP to defeat DNS-rebinding attacks."""
    from urllib.parse import urlparse
    import socket, ipaddress
    p = urlparse(url)
    if p.scheme not in ("http", "https") or not p.hostname:
        raise ValueError("Only public http(s) URLs are supported.")
    port = p.port or (443 if p.scheme == "https" else 80)
    try:
        infos = socket.getaddrinfo(p.hostname, port, proto=socket.IPPROTO_TCP)
    except OSError:
        raise ValueError("Could not resolve URL host.")
    pinned_ip = None
    for info in infos:
        addr = ipaddress.ip_address(info[4][0])
        if (addr.is_private or addr.is_loopback or addr.is_link_local or
                addr.is_reserved or addr.is_multicast or addr.is_unspecified):
            raise ValueError("URL points to a private/internal address — not allowed.")
        if pinned_ip is None:
            pinned_ip = str(addr)
    if not pinned_ip:
        raise ValueError("Could not resolve URL host.")
    sess = http.Session()
    sess.mount(f"{p.scheme}://", _PinnedIPAdapter(p.hostname, pinned_ip, p.scheme == "https"))
    try:
        r = sess.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"},
                     allow_redirects=False, stream=True)
        if 300 <= r.status_code < 400:
            raise ValueError("URL redirects are not supported — paste the final URL.")
        r.raise_for_status()
        content = r.raw.read(5 * 1024 * 1024 + 1, decode_content=True)
        if len(content) > 5 * 1024 * 1024:
            raise ValueError("Page too large (max 5 MB).")
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"Could not fetch URL: {e}")
    finally:
        sess.close()
    html = content.decode(r.encoding or "utf-8", "replace")
    # Remove script/style blocks
    html = re.sub(r'<(script|style)[^>]*>.*?</\1>', ' ', html, flags=re.DOTALL | re.IGNORECASE)
    # Extract text from meaningful tags
    parts = []
    for m in re.finditer(r'<(h[1-6]|p|li)[^>]*>(.*?)</\1>', html, re.DOTALL | re.IGNORECASE):
        inner = re.sub(r'<[^>]+>', ' ', m.group(2))
        inner = inner.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>').replace('&nbsp;', ' ').replace('&#39;', "'").replace('&quot;', '"')
        inner = ' '.join(inner.split())
        if len(inner) > 20:
            parts.append(inner)
    return ' '.join(parts) if parts else ' '.join(re.sub(r'<[^>]+>', ' ', html).split())


@app.route("/api/summarize-text", methods=["POST"])
def summarize_text():
    # Rate-limit BEFORE the server-side URL fetch + multi-pass LLM run.
    if not _check_rate_limit(_client_ip(), scope="text", limit=_RATE_MAX):
        return jsonify({"error": "Too many requests — please wait a moment and try again."}), 429
    uid = _auth_optional(request)
    anon_ip = None
    if uid:
        ok, tok_left, reason = _consume_token(uid)
        if not ok:
            return jsonify({
                "error": "You have no tokens left. Upgrade to continue.",
                "code": "no_tokens", "tokens_remaining": 0
            }), 402
    else:
        anon_ip = _client_ip()
        ok, tok_left = _anon_consume(anon_ip)
        if not ok:
            return jsonify({
                "error": "You've used your free previews. Sign up free to get more.",
                "code": "signin_for_more", "tokens_remaining": 0
            }), 402

    # silent=True: a non-JSON body must not raise here (it would 415/500 AFTER
    # the token was already consumed above, with no refund). Empty body → {} →
    # falls through to the "No text or URL" refund path below.
    data = request.get_json(silent=True) or {}
    text = (data.get("text") or "").strip()
    url  = (data.get("url")  or "").strip()
    lang_param   = data.get("language", "auto")
    filename     = _safe_name(data.get("filename") or "pasted_text")
    detail_level = data.get("detail", "standard")
    txt_dcfg     = DETAIL.get(detail_level, DETAIL["standard"])
    include_quiz = data.get("mode", "full") != "summary"

    if not ollama_running():
        _refund_credit(uid, anon_ip)
        return jsonify({"error": "AI service is not configured. Set GROQ_API_KEY."}), 503

    if not text and url:
        try:
            text = _fetch_url_text(url)
        except ValueError as e:
            _refund_credit(uid, anon_ip)
            return jsonify({"error": str(e)}), 400

    if not text:
        _refund_credit(uid, anon_ip)
        return jsonify({"error": "No text or URL provided"}), 400

    # Cap total input so a huge paste / large fetched page can't amplify Groq cost.
    text = text[:500_000]

    language = lang_param if lang_param in ("ar", "en") else _detect_language(text)
    gen = _stream_text_as_sse(text, language, filename, "text", txt_dcfg, tok_left, include_quiz, uid=uid, anon_ip=anon_ip)
    return Response(
        stream_with_context(gen()),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"}
    )


# ── Anki export endpoint ───────────────────────────────────────────────────────

@app.route("/api/export/anki/<job_id>")
def export_anki(job_id):
    if not _valid_job(job_id):
        return jsonify({"error": "Invalid job ID"}), 400
    job = get_job(job_id)
    if not job:
        return jsonify({"error": "Job not found or expired"}), 404
    guide = job.get("guide", {})
    flashcards = [f for f in guide.get("flashcards", []) if isinstance(f, dict)]
    if not flashcards:
        return jsonify({"error": "No flashcards available for this job"}), 404

    import csv
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(["front", "back"])
    def _csv_safe(v):
        v = str(v)
        # Excel/Sheets evaluate a formula even when it's preceded by leading
        # whitespace or a carriage return, so test the first NON-blank char.
        stripped = v.lstrip(" \t\r\n")
        return "'" + v if stripped[:1] in ("=", "+", "-", "@") else v
    for fc in flashcards:
        writer.writerow([_csv_safe(fc.get("q", "")), _csv_safe(fc.get("a", ""))])

    csv_bytes = buf.getvalue().encode("utf-8")
    base = job.get("filename", "study_guide").replace(".pdf", "")
    download_name = f"{_safe_name(base)}_anki.csv"
    return send_file(
        io.BytesIO(csv_bytes),
        mimetype="text/csv",
        as_attachment=True,
        download_name=download_name
    )


@app.route("/api/summarize", methods=["POST"])
def summarize():
    return jsonify({"error": "This endpoint is deprecated. Use /api/summarize-stream instead."}), 410



# \u2500\u2500 Legal pages (privacy / terms) \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
def _legal_shell(title, body):
    return """<!DOCTYPE html><html lang="en"><head><meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>""" + title + """ \u2014 Alimne</title>
<style>
  :root{color-scheme:dark}
  *{box-sizing:border-box}
  body{margin:0;background:#0a1628;color:#dce6f5;font-family:'Segoe UI',system-ui,-apple-system,sans-serif;line-height:1.7}
  .wrap{max-width:820px;margin:0 auto;padding:2.5rem 1.5rem 4rem}
  .brand{display:flex;align-items:center;gap:.6rem;margin-bottom:2rem;text-decoration:none}
  .logo{width:38px;height:38px;border-radius:11px;background:linear-gradient(135deg,#4f8ef7,#a78bfa);
        display:flex;align-items:center;justify-content:center;font-weight:800;color:#fff;font-size:1.1rem;flex-shrink:0}
  .brand span{font-weight:800;font-size:1.15rem;color:#fff}
  .brand small{display:block;font-weight:500;font-size:.7rem;color:#7f93b3;letter-spacing:.02em}
  h1{font-size:1.9rem;color:#fff;margin:.2rem 0 .3rem}
  .updated{color:#7f93b3;font-size:.85rem;margin-bottom:2rem}
  h2{color:#8fb4ff;font-size:1.15rem;margin:2.2rem 0 .6rem}
  p,li{color:#c2d0e6;font-size:.95rem}
  a{color:#6fa8ff}
  ul{padding-left:1.2rem}
  .note{background:rgba(79,142,247,.08);border:1px solid rgba(79,142,247,.2);border-radius:12px;padding:1rem 1.2rem;margin:1.5rem 0}
  footer{margin-top:3rem;padding-top:1.5rem;border-top:1px solid rgba(120,140,180,.18);color:#7f93b3;font-size:.82rem}
  footer a{margin-right:1rem}
</style></head><body><div class="wrap">
<a class="brand" href="/"><div class="logo">A</div><span>Alimne<small>by souc ai</small></span></a>
""" + body + """
<footer>
  <a href="/">Home</a><a href="/privacy">Privacy</a><a href="/terms">Terms</a>
  <div style="margin-top:.6rem">\u00a9 2026 Alimne \u00b7 a souc ai product \u00b7 <a href="mailto:sales@souc.ai">sales@souc.ai</a></div>
</footer></div></body></html>"""


@app.route("/privacy")
def privacy_page():
    body = """
<h1>Privacy Policy</h1>
<div class="updated">Last updated: 26 August 2026</div>
<p>Alimne ("we", "us"), operated by souc ai, turns your slides, documents, pasted text and
YouTube videos into study guides and summaries. Privacy is core to how the product is built.
This policy explains what we handle and why.</p>

<div class="note"><strong>The short version:</strong> the files and text you upload are processed
<strong>in memory only</strong>, are <strong>never written to disk or seen by any human</strong>, and are
<strong>automatically deleted within 15 minutes</strong> \u2014 or immediately when you press
"Delete now". We do not sell your data and we do not show ads.</div>

<h2>1. Study content you submit</h2>
<ul>
<li><strong>Files, pasted text and URLs</strong> are held in server memory only for the length of your
session (maximum 15 minutes) and are purged automatically after that window, or instantly on your request.
They are never persisted to disk, logged in full, or reviewed by a person.</li>
<li>To generate a guide, the extracted text is sent to our AI provider (Groq) for processing. It is used
only to produce your result and is not used to train models by us.</li>
</ul>

<h2>2. Account information</h2>
<ul>
<li>If you sign in, we store your <strong>email address</strong> and a display name/avatar (when provided by
Google) in our authentication database (Supabase) to identify your account and track your monthly token balance.</li>
<li>We use a session cookie / local storage entry to keep you signed in.</li>
</ul>

<h2>3. Payments</h2>
<p>Subscriptions are processed by <strong>Stripe, Inc.</strong> We never receive or store your full card
number \u2014 Stripe handles payment details directly. We store only a Stripe customer reference and your
subscription status.</p>

<h2>4. What we do not do</h2>
<ul>
<li>We do not sell, rent, or trade your personal data.</li>
<li>We do not run advertising or third-party ad trackers.</li>
<li>We do not retain your study material beyond the 90-minute processing window.</li>
</ul>

<h2>5. Data retention &amp; your rights</h2>
<p>Study jobs: deleted within 15 minutes (or on demand). Account data: kept until you ask us to delete it.
You may request access to, or deletion of, your account data at any time by emailing
<a href="mailto:sales@souc.ai">sales@souc.ai</a>.</p>

<h2>6. Third-party services</h2>
<p>We rely on Supabase (authentication &amp; account database), Stripe (payments), Groq (AI processing),
Render (hosting) and Cloudflare (DNS/network). Each processes data only as needed to provide the service.</p>

<h2>7. Changes &amp; contact</h2>
<p>We may update this policy; material changes will be reflected here with a new date. Questions?
Email <a href="mailto:sales@souc.ai">sales@souc.ai</a>.</p>
"""
    return _legal_shell("Privacy Policy", body), 200, {"Content-Type": "text/html; charset=utf-8"}


@app.route("/terms")
def terms_page():
    body = """
<h1>Terms &amp; Conditions</h1>
<div class="updated">Last updated: 26 August 2026</div>
<p>By using Alimne (the "Service"), operated by souc ai, you agree to these terms.</p>

<h2>1. The Service</h2>
<p>Alimne converts uploaded slides, documents, pasted text and YouTube videos into study guides,
summaries, flashcards and quizzes using AI. Output is generated automatically and may contain
inaccuracies \u2014 always verify important information against the source material.</p>

<h2>2. Plans &amp; billing</h2>
<ul>
<li><strong>Free plan:</strong> 3 processing tokens per month. No credit card required.</li>
<li><strong>Pro plan:</strong> US$2.99 per month, billed via Stripe, including <strong>30 tokens per month</strong>
and priority processing. Your token allowance renews each billing cycle.</li>
<li>You can cancel anytime; access continues until the end of the paid period. Charges are non-refundable
except where required by law.</li>
</ul>

<h2>3. Acceptable use</h2>
<ul>
<li>Only upload content you have the right to use.</li>
<li>Do not use the Service for unlawful purposes or to process content that infringes others' rights.</li>
<li>Do not attempt to disrupt, overload, or reverse-engineer the Service.</li>
</ul>

<h2>4. Your content</h2>
<p>You retain all rights to the content you submit. As described in our
<a href="/privacy">Privacy Policy</a>, your content is processed in memory only and deleted within
15 minutes. You are responsible for keeping your own copies.</p>

<h2>5. Disclaimers &amp; liability</h2>
<p>The Service is provided "as is", without warranties of any kind. To the maximum extent permitted by
law, souc ai is not liable for any indirect or consequential damages, or for reliance on AI-generated
output. Third-party services (Stripe, Google, YouTube, etc.) are governed by their own terms.</p>

<h2>6. Changes &amp; governing law</h2>
<p>We may update these terms; continued use means acceptance. These terms are governed by the laws of
the United Arab Emirates (Dubai). Contact: <a href="mailto:sales@souc.ai">sales@souc.ai</a>.</p>
"""
    return _legal_shell("Terms & Conditions", body), 200, {"Content-Type": "text/html; charset=utf-8"}


@app.route("/", defaults={"path": ""})
@app.route("/<path:path>")
def serve(path):
    if path and os.path.exists(os.path.join(DIST, path)):
        return send_from_directory(DIST, path)
    return send_from_directory(DIST, "index.html")


if __name__ == "__main__":
    print(f"  Model : {OLLAMA_MODEL}")
    print(f"  Ollama: {'running' if ollama_running() else 'NOT running — run: ollama serve'}")
    app.run(debug=False, host="127.0.0.1", port=5000, threaded=True)
